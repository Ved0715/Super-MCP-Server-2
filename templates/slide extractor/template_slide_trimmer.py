import shutil
import os
from datetime import datetime
from pptx import Presentation

class FixedTemplateTrimmer:
    def __init__(self, template_path, output_path=None):
        self.template_path = template_path
        self.output_path = output_path or self.generate_output_filename()
        
    def generate_output_filename(self):
        """Generate timestamped output filename"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        base_name = os.path.splitext(os.path.basename(self.template_path))[0]
        return f"{base_name}_FIXED_20_SLIDES_{timestamp}.pptx"
    
    def analyze_slide_text_content(self, slide):
        """Calculate total text content in a slide"""
        total_chars = 0
        total_words = 0
        text_content = []
        
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    total_chars += len(text)
                    total_words += len(text.split())
                    text_content.append(text)
        
        return {
            'char_count': total_chars,
            'word_count': total_words,
            'text_content': ' '.join(text_content),
            'text_density': total_chars + (total_words * 2)
        }
    
    def is_unwanted_slide(self, slide, text_content):
        """Identify slides that should be excluded"""
        text_lower = text_content.lower()
        
        unwanted_keywords = [
            'icon pack', 'icons', 'infographic', 'template', 'color palette',
            'font', 'slidesgo', 'credits', 'attribution', 'resources',
            'about this presentation', 'how to use', 'instructions'
        ]
        
        if any(keyword in text_lower for keyword in unwanted_keywords):
            return True
        
        if len(text_content.strip()) < 15:
            return True
        
        text_shapes = sum(1 for shape in slide.shapes 
                         if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text_frame.text.strip())
        total_shapes = len(slide.shapes)
        
        if total_shapes > 10 and text_shapes < 3:
            return True
        
        return False
    
    def classify_slide(self, slide_idx, slide, analysis):
        """Classify slide type with flexible criteria"""
        text_content = analysis['text_content']
        text_lower = text_content.lower()
        
        if self.is_unwanted_slide(slide, text_content):
            return 'unwanted'
        
        if (slide_idx == 0 or 
            any(keyword in text_lower for keyword in ['title', 'welcome', 'introduction', 'overview']) and
            analysis['word_count'] < 50):
            return 'title'
        
        elif any(keyword in text_lower for keyword in ['thank you', 'thanks', 'questions', 'conclusion', 'contact us', 'end']):
            return 'thank_you'
        
        elif analysis['char_count'] >= 30 and analysis['word_count'] >= 8:
            return 'content'
        
        else:
            return 'unwanted'
    
    def select_slides_intelligently(self, target_slides=20):
        """Select slides with intelligent filtering"""
        prs = Presentation(self.template_path)
        all_slides = []
        
        print(f"🔍 Analyzing {len(prs.slides)} slides...")
        
        for slide_idx, slide in enumerate(prs.slides):
            analysis = self.analyze_slide_text_content(slide)
            slide_type = self.classify_slide(slide_idx, slide, analysis)
            
            all_slides.append({
                'index': slide_idx,
                'type': slide_type,
                'analysis': analysis
            })
        
        # Separate slides by type
        title_slides = [s for s in all_slides if s['type'] == 'title']
        thank_you_slides = [s for s in all_slides if s['type'] == 'thank_you']
        content_slides = [s for s in all_slides if s['type'] == 'content']
        
        print(f"📊 Found: {len(title_slides)} title, {len(content_slides)} content, {len(thank_you_slides)} thank you slides")
        
        # Sort content slides by quality
        content_slides.sort(key=lambda x: x['analysis']['text_density'], reverse=True)
        
        # Build selection ensuring we don't exceed available content
        final_selection = []
        
        # Add title slide
        if title_slides:
            final_selection.append(title_slides[0])
        elif content_slides:
            final_selection.append(content_slides.pop(0))
        
        # Calculate content slots available
        thank_you_selected = thank_you_slides[0] if thank_you_slides else None
        available_content_slots = target_slides - len(final_selection) - (1 if thank_you_selected else 0)
        content_to_use = min(available_content_slots, len(content_slides))
        
        # Add content slides
        final_selection.extend(content_slides[:content_to_use])
        
        # Add thank you slide
        if thank_you_selected:
            final_selection.append(thank_you_selected)
        
        # Fill any remaining slots if needed
        remaining_slots = target_slides - len(final_selection)
        if remaining_slots > 0 and len(content_slides) > content_to_use:
            additional_content = content_slides[content_to_use:content_to_use + remaining_slots]
            final_selection.extend(additional_content)
        
        print(f"📋 Selected {len(final_selection)} slides for final presentation")
        return final_selection
    
    def safe_slide_removal(self, presentation, slides_to_keep):
        """FIXED: Safe slide removal with proper index handling"""
        try:
            # Get total slide count BEFORE any removal
            total_slides = len(presentation.slides)
            
            # Create set of indices to keep for fast lookup
            keep_indices = set(slide_info['index'] for slide_info in slides_to_keep)
            
            # Identify slides to remove
            slides_to_remove = [idx for idx in range(total_slides) if idx not in keep_indices]
            
            print(f"🗑️ Removing {len(slides_to_remove)} slides...")
            
            # CRITICAL FIX: Remove slides in REVERSE ORDER to maintain valid indices
            xml_slides = presentation.slides._sldIdLst
            
            for idx in reversed(slides_to_remove):
                try:
                    # SAFETY CHECK: Verify index is still valid
                    if idx < len(xml_slides):
                        rId = xml_slides[idx].rId
                        presentation.part.drop_rel(rId)
                        xml_slides.remove(xml_slides[idx])
                        print(f"   ✅ Removed slide {idx + 1}")
                    else:
                        print(f"   ⚠️ Skipped slide {idx + 1} - index out of range")
                        
                except Exception as e:
                    print(f"   ❌ Error removing slide {idx + 1}: {e}")
                    continue
            
            print(f"📊 Slide removal completed. Final count: {len(presentation.slides)}")
            
        except Exception as e:
            print(f"❌ Critical error in slide removal: {e}")
            raise
    
    def create_trimmed_presentation(self):
        """Create trimmed presentation with safe slide removal"""
        # Step 1: Select slides intelligently
        selected_slides = self.select_slides_intelligently()
        
        if not selected_slides:
            print("❌ No suitable slides found")
            return None
        
        # Step 2: Copy template to preserve design
        print(f"📁 Copying template...")
        shutil.copy2(self.template_path, self.output_path)
        
        # Step 3: Load and process presentation
        presentation = Presentation(self.output_path)
        original_count = len(presentation.slides)
        
        # Step 4: SAFE slide removal
        self.safe_slide_removal(presentation, selected_slides)
        
        # Step 5: Save presentation
        try:
            presentation.save(self.output_path)
            final_count = len(presentation.slides)
            
            print(f"\n🎉 SUCCESS! Template trimmed successfully!")
            print(f"📁 Original slides: {original_count}")
            print(f"📁 Final slides: {final_count}")
            print(f"💾 Saved as: {self.output_path}")
            
            return self.output_path
            
        except Exception as e:
            print(f"❌ Error saving presentation: {e}")
            return None

def main():
    print("=== FIXED Template Trimmer - No Index Errors ===")
    
    template_path = input("Enter template file path: ").strip()
    if not template_path:
        template_path = "template.pptx"
    
    if not os.path.exists(template_path):
        print(f"❌ Template file '{template_path}' not found!")
        return
    
    print(f"\n{'='*60}")
    print("FIXED TRIMMING - INDEX ERROR RESOLVED")
    print(f"{'='*60}")
    print(f"📋 Template: {template_path}")
    print(f"🛠️  Fixed: List index out of range errors")
    print(f"✅ Safe removal in reverse order")
    print(f"🎯 Smart structure with available content")
    print(f"{'='*60}")
    
    confirm = input("\nProceed with fixed trimming? (y/n): ").lower().strip()
    if confirm not in ['y', 'yes']:
        print("Trimming cancelled.")
        return
    
    try:
        trimmer = FixedTemplateTrimmer(template_path)
        result = trimmer.create_trimmed_presentation()
        
        if result:
            print(f"\n✅ Template successfully processed!")
            
    except Exception as e:
        print(f"❌ Unexpected error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()
