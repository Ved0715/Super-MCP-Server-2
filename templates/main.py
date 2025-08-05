from .extract_metadata import extract_all_metadata, is_thank_you_slide
from .prompt_builder import build_prompt
from .ai_content_fetcher import get_ai_content, create_enhanced_fallback_content
from .replace_content import replace_with_ai_content
from .content_format_analyzer import content_analyzer
from pptx import Presentation
import os
import time

# Global variable to store user content
USER_CONTENT = ""

def get_user_content_and_topic():
    """Get presentation content and topic from user input"""
    print("=== AI PowerPoint Content Formatter ===")
    print("Enter your complete content to be formatted for the presentation:")
    print("(Paste your content, then press Enter twice to finish)")
    print("-" * 50)
    
    content_lines = []
    empty_line_count = 0
    
    while True:
        try:
            line = input()
            if line.strip() == "":
                empty_line_count += 1
                if empty_line_count >= 2:
                    break
                content_lines.append(line)
            else:
                empty_line_count = 0
                content_lines.append(line)
        except EOFError:
            break
    
    # Remove trailing empty lines
    while content_lines and content_lines[-1].strip() == "":
        content_lines.pop()
    
    user_content = "\n".join(content_lines).strip()
    
    if not user_content:
        print("❌ No content provided. Please try again.")
        return get_user_content_and_topic()
    
    # Get topic title for context
    topic_title = input("\nWhat's the main topic/title for this presentation? ").strip()
    if not topic_title:
        topic_title = "Presentation"
    
    return topic_title, user_content


def get_template_path():
    """Get the single template path"""
    template_path = "templates/test_slide.pptx"
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Templatt e not found: {template_path}")
    return template_path


def generate_output_filename(topic):
    """Generate output filename based on topic"""
    # Clean topic for filename (remove special characters)
    clean_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '-', '_')).rstrip()
    clean_topic = clean_topic.replace(' ', '_')
    return f"presentations/{clean_topic}_presentation.pptx"


def batch_metadata_by_token_estimate(metadata, max_tokens_per_batch=4000):
    """Batch metadata by estimated token count"""
    batches = []
    current_batch = []
    current_tokens = 0
    
    for field in metadata:
        approx_tokens = len(str(field)) // 4 + field["char_count"] // 4 + 20  # Extra buffer
        if current_tokens + approx_tokens > max_tokens_per_batch and current_batch:
            batches.append(current_batch)
            current_batch = [field]
            current_tokens = approx_tokens
        else:
            current_batch.append(field)
            current_tokens += approx_tokens
    
    if current_batch:
        batches.append(current_batch)
    return batches


def generate_presentation_api(content, topic):
    """
    API function to generate PowerPoint presentation
    
    Args:
        content (str): The content to format for presentation
        topic (str): Main topic/title for presentation
    
    Returns:
        str: Output filename of generated presentation
    """
    try:
        # Get template and output paths
        template_path = get_template_path()
        output_path = generate_output_filename(topic)
        
        # Extract metadata from template
        print(f"🚀 Starting AI-powered content formatting...")
        all_metadata = extract_all_metadata(template_path)
        
        # Filter out thank you slides
        metadata = []
        skipped_slides = []
        
        # Load presentation to check slides
        prs = Presentation(template_path)
        
        for item in all_metadata:
            slide_idx = item.get('slide', 0)
            if slide_idx < len(prs.slides) and is_thank_you_slide(prs.slides[slide_idx]):
                skipped_slides.append(slide_idx + 1)  # 1-indexed for display
            else:
                metadata.append(item)
        
        print(f"✅ Extracted {len(metadata)} text fields to replace")
        if skipped_slides:
            print(f"🔒 Preserved slides: {', '.join(map(str, skipped_slides))} (thank you/closing slides)")
        
        # Process in batches to avoid overwhelming the AI
        batch_size = 20  # Reduced from 40 for better AI completion
        batches = [metadata[i:i + batch_size] for i in range(0, len(metadata), batch_size)]
        all_ai_results = []
        total_batches = len(batches)
        
        print(f"📦 Processing {total_batches} batches...")
        
        for batch_index, batch in enumerate(batches):
            print(f"🔄 Processing batch {batch_index + 1} of {total_batches} ({len(batch)} fields)...")
            
            success = False
            for attempt in range(3):  # Increased attempts
                try:
                    # Use enhanced mode for first attempt, simple mode for retries
                    use_enhanced = attempt == 0
                    prompt = build_prompt(topic, batch, use_enhanced=use_enhanced, user_content=content)
                    ai_results = get_ai_content(prompt, batch, user_topic=topic, use_enhanced=use_enhanced, user_content=content)
                    
                    # Validate that we got results for all fields in the batch
                    if len(ai_results) < len(batch):
                        print(f"   ❌ CRITICAL: AI failed to generate all items. Expected {len(batch)}, got {len(ai_results)}")
                        print(f"   🔄 This should not happen with the new retry logic!")
                        # The enhanced AI retry logic in get_ai_content should handle missing fields
                    
                    all_ai_results.extend(ai_results)
                    success = True
                    print(f"   ✅ Batch {batch_index + 1} completed successfully ({len(ai_results)} fields processed)")
                    break
                except Exception as e:
                    print(f"   ❌ Batch {batch_index + 1} failed (attempt {attempt+1}): {e}")
                    if attempt < 2:
                        print(f"   🔄 Retrying with simpler prompt...")
                        time.sleep(2)  # Brief pause between retries
            
            if not success:
                print(f"   ⚠️ Batch {batch_index+1} permanently skipped. See logs.")
        
        print(f"\n📊 Generated content for {len(all_ai_results)} fields")
        
        # Replace content in PowerPoint
        print("📝 Replacing content in PowerPoint template...")
        replace_with_ai_content(template_path, all_ai_results, output_path)
        
        return output_path
        
    except Exception as e:
        print(f"❌ Error generating presentation: {e}")
        raise


def generate_presentation_api_enhanced(content, topic, user_prompt, num_slides):
    """
    Enhanced API function to generate PowerPoint presentation with user prompt decomposition
    Args:
        content (str): The content to format for presentation
        topic (str): Main topic/title for presentation
        user_prompt (str): User's presentation prompt/instructions
        num_slides (int): Number of slides to generate
    Returns:
        str: Output filename of generated presentation
    """
    try:
        from tqdm import tqdm
        import time
        
        template_path = get_template_path()
        output_path = generate_output_filename(topic)
        
        print(f"🚀 Starting Enhanced AI-powered presentation generation...")
        print(f"📋 Topic: {topic}")
        print(f"📝 User Prompt: {user_prompt}")
        print(f"📊 Number of Slides: {num_slides}")
        print(f"📁 Template: {template_path}")
        print(f"📁 Output: {output_path}")
        
        # Use the new question-based generation with progress bars
        from .prompt_builder import generate_presentation_with_prompt
        all_ai_results = generate_presentation_with_prompt(content, user_prompt, topic, num_slides)
        
        if not all_ai_results:
            raise ValueError("No content was generated")
        
        print(f"\n📝 Replacing content in PowerPoint template...")
        with tqdm(total=1, desc="💾 Saving presentation", bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt}") as pbar:
            replace_with_ai_content(template_path, all_ai_results, output_path)
            pbar.update(1)
        
        return output_path
        
    except Exception as e:
        print(f"❌ Error generating enhanced presentation: {e}")
        raise




def preview_prompt_breakdown(content, topic, user_prompt, num_slides):
    """
    Preview function to show prompt breakdown and content extraction without generating presentation
    Args:
        content (str): The content to analyze
        topic (str): Main topic/title for presentation
        user_prompt (str): User's presentation prompt/instructions
        num_slides (int): Number of slides to generate
    Returns:
        dict: Preview data with questions and extracted content
    """
    try:
        from tqdm import tqdm
        import time
        
        print(f"🎯 PREVIEW MODE: Analyzing prompt breakdown and content extraction...")
        print(f"📋 Topic: {topic}")
        print(f"📝 User Prompt: {user_prompt}")
        print(f"📊 Number of Slides: {num_slides}")
        print(f"📄 Content Length: {len(content)} characters")
        print("="*80)
        
        # Step 1: Decompose user prompt into questions
        from .prompt_builder import decompose_user_prompt
        print(f"\n🔍 STEP 1: Decomposing user prompt into {num_slides} questions...")
        
        with tqdm(total=1, desc="🤖 Generating questions", bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt}") as pbar:
            slide_questions = decompose_user_prompt(user_prompt, num_slides, topic)
            pbar.update(1)
        
        print(f"✅ Generated {len(slide_questions)} questions:")
        for i, question in enumerate(slide_questions):
            print(f"   Slide {i+1}: {question}")
        
        # Step 2: Extract relevant content for each question
        from .prompt_builder import extract_relevant_content
        print(f"\n🔍 STEP 2: Extracting relevant content for each question...")
        
        preview_data = {
            "topic": topic,
            "user_prompt": user_prompt,
            "num_slides": num_slides,
            "content_length": len(content),
            "questions": [],
            "total_extracted_content": 0
        }
        
        # Progress bar for content extraction
        with tqdm(total=len(slide_questions), desc="📝 Extracting content", bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt} questions") as pbar:
            for i, question in enumerate(slide_questions):
                pbar.set_description(f"📝 Processing Slide {i+1}")
                
                # Extract relevant content for this question
                relevant_content = extract_relevant_content(content, question, topic)
                
                # Analyze content format with flow (no template needed for method 3)
                format_analysis = content_analyzer.analyze_content_format_with_flow(
                    question, relevant_content, [], i + 1, num_slides, topic, template_inventory=None
                )
                
                # Store preview data
                question_data = {
                    "slide_num": i + 1,
                    "question": question,
                    "extracted_content": relevant_content,
                    "content_length": len(relevant_content),
                    "format_analysis": format_analysis
                }
                preview_data["questions"].append(question_data)
                preview_data["total_extracted_content"] += len(relevant_content)
                
                # Update progress bar
                pbar.update(1)
                
                # Small delay to show progress
                time.sleep(0.1)
        
        # Step 3: Show summary
        print(f"\n📊 PREVIEW SUMMARY:")
        print(f"   📋 Topic: {topic}")
        print(f"   🎯 User Prompt: {user_prompt}")
        print(f"   📊 Number of Slides: {num_slides}")
        print(f"   📄 Original Content: {len(content)} characters")
        print(f"   📝 Total Extracted Content: {preview_data['total_extracted_content']} characters")
        print(f"   📈 Content Coverage: {(preview_data['total_extracted_content'] / len(content) * 100):.1f}%")
        
        # Step 4: Display format analysis and extracted content for each question
        print(f"\n🎯 CONTENT FORMAT ANALYSIS & EXTRACTED CONTENT:")
        print("="*80)
        
        for question_data in preview_data["questions"]:
            format_analysis = question_data.get('format_analysis', {})
            
            print(f"\n📝 SLIDE {question_data['slide_num']}:")
            print(f"   🎯 Question: {question_data['question']}")
            print(f"   📊 Content Format: {format_analysis.get('content_format', 'Unknown').replace('_', ' ').title()}")
            print(f"   🎨 Recommended Slide Type: {format_analysis.get('recommended_slide_type', 'Unknown').replace('_', ' ').title()}")
            print(f"   🎯 Confidence: {format_analysis.get('confidence_score', 0):.2f}")
            print(f"   📋 Template Match: Slide {format_analysis.get('template_match', {}).get('slide_number', 0) + 1} ({format_analysis.get('template_match', {}).get('template_type', 'Unknown').replace('_', ' ').title()})")
            print(f"   🎭 Slide Role: {format_analysis.get('slide_role', 'Unknown').replace('_', ' ').title()}")
            print(f"   📈 Flow Adjustment: {format_analysis.get('flow_adjustment', 0):.2f}")
            print(f"   📄 Extracted Content ({question_data['content_length']} characters):")
            print(f"   {'─' * 60}")
            
            # Show full content
            content_lines = question_data['extracted_content'].split('\n')
            for line in content_lines:
                if line.strip():
                    print(f"   {line}")
            print(f"   {'─' * 60}")
        
        print(f"\n🎉 PREVIEW COMPLETED! No presentation was generated.")
        print(f"💡 To generate the actual presentation, use generate_presentation_api_enhanced()")
        
        return preview_data
        
    except Exception as e:
        print(f"❌ Error in preview mode: {e}")
        raise


def preview_prompt_breakdown_api(content, topic, user_prompt, num_slides):
    """
    API function to preview prompt breakdown and content extraction without generating presentation
    Args:
        content (str): The content to analyze
        topic (str): Main topic/title for presentation
        user_prompt (str): User's presentation prompt/instructions
        num_slides (int): Number of slides to generate
    Returns:
        dict: Preview data with questions and extracted content
    """
    try:
        # Use the same preview logic but return data instead of printing
        from .prompt_builder import decompose_user_prompt, extract_relevant_content
        
        # Step 1: Decompose user prompt into questions
        slide_questions = decompose_user_prompt(user_prompt, num_slides, topic)
        
        # Step 2: Extract relevant content for each question
        preview_data = {
            "topic": topic,
            "user_prompt": user_prompt,
            "num_slides": num_slides,
            "content_length": len(content),
            "questions": [],
            "total_extracted_content": 0,
            "success": True
        }
        
        for i, question in enumerate(slide_questions):
            # Extract relevant content for this question
            relevant_content = extract_relevant_content(content, question, topic)
            
            # Store preview data
            question_data = {
                "slide_num": i + 1,
                "question": question,
                "extracted_content": relevant_content,
                "content_length": len(relevant_content)
            }
            preview_data["questions"].append(question_data)
            preview_data["total_extracted_content"] += len(relevant_content)
        
        # Calculate content coverage
        preview_data["content_coverage_percentage"] = (preview_data["total_extracted_content"] / len(content) * 100) if len(content) > 0 else 0
        
        return preview_data
        
    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "topic": topic,
            "user_prompt": user_prompt,
            "num_slides": num_slides,
            "content_length": len(content) if content else 0
        }


def generate_presentation_with_preview_api(content, topic, user_prompt, num_slides):
    """
    API function for Preview Mode with PowerPoint generation (Option 3 from menu)
    Combines preview analysis with direct PowerPoint generation
    
    Args:
        content (str): The content to format for presentation
        topic (str): Main topic/title for presentation
        user_prompt (str): User's presentation prompt/instructions
        num_slides (int): Number of slides to generate
    
    Returns:
        dict: Result with preview data and output file path
    """
    try:
        # Step 1: Get preview data (same as preview_prompt_breakdown)
        from .prompt_builder import decompose_user_prompt, extract_relevant_content
        
        # Decompose user prompt into questions
        slide_questions = decompose_user_prompt(user_prompt, num_slides, topic)
        
        # Extract relevant content and analyze format for each question
        preview_data = {
            "topic": topic,
            "user_prompt": user_prompt,
            "num_slides": num_slides,
            "content_length": len(content),
            "questions": [],
            "total_extracted_content": 0,
            "success": True
        }
        
        content_analysis_results = []
        
        for i, question in enumerate(slide_questions):
            # Extract relevant content for this question
            relevant_content = extract_relevant_content(content, question, topic)
            
            # Analyze content format with flow (no template needed for method 3)
            format_analysis = content_analyzer.analyze_content_format_with_flow(
                question, relevant_content, [], i + 1, num_slides, topic, template_inventory=None
            )
            
            # Store preview data
            question_data = {
                "slide_num": i + 1,
                "question": question,
                "extracted_content": relevant_content,
                "content_length": len(relevant_content),
                "format_analysis": format_analysis
            }
            preview_data["questions"].append(question_data)
            preview_data["total_extracted_content"] += len(relevant_content)
            
            # Prepare slide data for PowerPoint generation
            slide_data = {
                'format_type': format_analysis.get('content_format', 'paragraph'),
                'content': relevant_content,
                'question': question,
                'slide_num': i + 1,
                'confidence': format_analysis.get('confidence_score', 0.0),
                'recommended_slide_type': format_analysis.get('recommended_slide_type', 'content_slide'),
                'template_match': format_analysis.get('template_match', {}),
                'slide_role': format_analysis.get('slide_role', ''),
                'flow_adjustment': format_analysis.get('flow_adjustment', 0.0)
            }
            content_analysis_results.append(slide_data)
        
        # Calculate content coverage
        preview_data["content_coverage_percentage"] = (preview_data["total_extracted_content"] / len(content) * 100) if len(content) > 0 else 0
        
        # Step 2: Generate PowerPoint presentation
        output_path = generate_output_filename(topic)
        
        # Create PowerPoint with title slide, content slides, and thank you slide
        from .replace_content import create_ppt_with_title_and_thanks
        create_ppt_with_title_and_thanks(content_analysis_results, output_path, topic, user_prompt)
        
        # Return combined result
        return {
            "success": True,
            "output_file": output_path,
            "preview_data": preview_data,
            "slides_generated": num_slides,
            "format_analysis": [q["format_analysis"] for q in preview_data["questions"]]
        }
        
    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "topic": topic,
            "user_prompt": user_prompt,
            "num_slides": num_slides,
            "content_length": len(content) if content else 0,
            "output_file": None
        }


def main():
    """CLI version for interactive use"""
    print("=== AI PowerPoint Content Formatter ===")
    print("Choose generation mode:")
    print("1. Standard Mode (content formatting only)")
    print("2. Enhanced Mode (user prompt + question-based generation)")
    print("3. Preview Mode (show prompt breakdown and content extraction + create PPT)")
    
    mode_choice = input("\nEnter your choice (1, 2, or 3): ").strip()
    
    if mode_choice == "1":
        # Standard mode
        user_topic, user_content = get_user_content_and_topic()
        template_path = get_template_path()
        output_path = generate_output_filename(user_topic)
        
        print(f"\n📋 Processing Summary:")
        print(f"📁 Template: {template_path}")
        print(f"📁 Output: {output_path}")
        print(f"📊 Topic: {user_topic}")
        print(f"📝 Content Length: {len(user_content)} characters")
        
        confirm = input("\nProceed with processing? (y/n): ").lower().strip()
        if confirm != 'y':
            print("❌ Processing cancelled.")
            return
        
        try:
            output_file = generate_presentation_api(user_content, user_topic)
            print(f"\n🎉 FORMATTING COMPLETED! 🎉")
            print(f"📁 Output saved as: {output_file}")
            print(f"📊 Topic: {user_topic}")
        except Exception as e:
            print(f"\n❌ Error: {e}")
            print("Please check the logs above for details.")
    
    elif mode_choice == "2":
        # Enhanced mode
        print("\n=== Enhanced Mode: User Prompt + Question-Based Generation ===")
        
        user_topic = input("Enter presentation topic: ").strip()
        if not user_topic:
            print("❌ Topic is required.")
            return
        
        print("\nEnter your content (press Enter twice to finish):")
        user_content_lines = []
        while True:
            line = input()
            if line == "" and user_content_lines and user_content_lines[-1] == "":
                break
            user_content_lines.append(line)
        
        user_content = "\n".join(user_content_lines[:-1])  # Remove the last empty line
        if not user_content.strip():
            print("❌ Content is required.")
            return
        
        print("\nEnter your presentation prompt/instructions:")
        print("Example: 'Start with introduction, then explain methodology, present results, discuss implications, and conclude'")
        user_prompt = input("Prompt: ").strip()
        if not user_prompt:
            print("❌ User prompt is required for enhanced mode.")
            return
        
        try:
            num_slides = int(input("Enter number of slides: ").strip())
            if num_slides < 3:
                print("❌ Number of slides must be at least 3 (1 title + 1 content + 1 thank you).")
                return
        except ValueError:
            print("❌ Please enter a valid number of slides.")
            return
        
        template_path = get_template_path()
        output_path = generate_output_filename(user_topic)
        
        print(f"\n📋 Enhanced Processing Summary:")
        print(f"📁 Template: {template_path}")
        print(f"📁 Output: {output_path}")
        print(f"📊 Topic: {user_topic}")
        print(f"📝 Content Length: {len(user_content)} characters")
        print(f"🎯 User Prompt: {user_prompt}")
        print(f"📊 Number of Slides: {num_slides}")
        
        confirm = input("\nProceed with enhanced processing? (y/n): ").lower().strip()
        if confirm != 'y':
            print("❌ Processing cancelled.")
            return
        
        try:
            output_file = generate_presentation_api_enhanced(user_content, user_topic, user_prompt, num_slides)
            print(f"\n🎉 ENHANCED GENERATION COMPLETED! 🎉")
            print(f"📁 Output saved as: {output_file}")
            print(f"📊 Topic: {user_topic}")
            print(f"🎯 Generated {num_slides} slides based on your prompt!")
        except Exception as e:
            print(f"\n❌ Error: {e}")
            print("Please check the logs above for details.")
    
    elif mode_choice == "3":
        # Preview mode
        print("\n=== Preview Mode: Show Prompt Breakdown and Content Extraction ===")
        
        user_topic = input("Enter presentation topic: ").strip()
        if not user_topic:
            print("❌ Topic is required.")
            return
        
        print("\nEnter your content (press Enter twice to finish):")
        user_content_lines = []
        while True:
            line = input()
            if line == "" and user_content_lines and user_content_lines[-1] == "":
                break
            user_content_lines.append(line)
        
        user_content = "\n".join(user_content_lines[:-1])  # Remove the last empty line
        if not user_content.strip():
            print("❌ Content is required.")
            return
        
        print("\nEnter your presentation prompt/instructions:")
        print("Example: 'Start with introduction, then explain methodology, present results, discuss implications, and conclude'")
        user_prompt = input("Prompt: ").strip()
        if not user_prompt:
            print("❌ User prompt is required for preview mode.")
            return
        
        try:
            num_slides = int(input("Enter number of slides: ").strip())
            if num_slides < 3:
                print("❌ Number of slides must be at least 3 (1 title + 1 content + 1 thank you).")
                return
        except ValueError:
            print("❌ Please enter a valid number of slides.")
            return
        
        print(f"\n📋 Preview Summary:")
        print(f"📊 Topic: {user_topic}")
        print(f"📝 Content Length: {len(user_content)} characters")
        print(f"🎯 User Prompt: {user_prompt}")
        print(f"📊 Number of Slides: {num_slides}")
        
        confirm = input("\nProceed with preview? (y/n): ").lower().strip()
        if confirm != 'y':
            print("❌ Preview cancelled.")
            return
        
        try:
            preview_data = preview_prompt_breakdown(user_content, user_topic, user_prompt, num_slides)
            print(f"\n🎉 PREVIEW COMPLETED! 🎉")
            print(f"📊 Topic: {user_topic}")
            print(f"🎯 Analyzed {num_slides} questions from your prompt!")
            
            # Ask if user wants to create PowerPoint
            create_ppt = input(f"\n🎨 Would you like to create a PowerPoint presentation with format-specific layouts? (y/n): ").lower().strip()
            if create_ppt == 'y':
                print(f"\n🎨 Creating PowerPoint with format-specific layouts...")
                output_path = generate_output_filename(user_topic)
                
                # Create content analysis results from preview data
                content_analysis_results = []
                for i, question_data in enumerate(preview_data["questions"]):
                    slide_data = {
                        'format_type': question_data["format_analysis"].get('content_format', 'paragraph'),
                        'content': question_data["extracted_content"],
                        'question': question_data["question"],
                        'slide_num': i + 1,
                        'confidence': question_data["format_analysis"].get('confidence_score', 0.0),
                        'recommended_slide_type': question_data["format_analysis"].get('recommended_slide_type', 'content_slide'),
                        'template_match': question_data["format_analysis"].get('template_match', {}),
                        'slide_role': question_data["format_analysis"].get('slide_role', ''),
                        'flow_adjustment': question_data["format_analysis"].get('flow_adjustment', 0.0)
                    }
                    content_analysis_results.append(slide_data)
                
                # Create PowerPoint with title slide, content slides, and thank you slide
                from .replace_content import create_ppt_with_title_and_thanks
                from tqdm import tqdm
                
                with tqdm(total=1, desc="💾 Creating presentation", bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt}") as pbar:
                    create_ppt_with_title_and_thanks(content_analysis_results, output_path, user_topic, user_prompt)
                    pbar.update(1)
                
                print(f"\n🎉 POWERPOINT CREATION COMPLETED! 🎉")
                print(f"📁 Output saved as: {output_path}")
                print(f"📊 Topic: {user_topic}")
                print(f"🎯 Generated {num_slides} slides with format-specific layouts!")
                print(f"🎨 Each slide has custom visual layout matching its content format!")
            else:
                print(f"💡 To generate the actual presentation later, use this same option (3) and choose 'y' when prompted.")
                
        except Exception as e:
            print(f"\n❌ Error: {e}")
            print("Please check the logs above for details.")
    
    else:
        print("❌ Invalid choice. Please enter 1, 2, or 3.")


if __name__ == "__main__":
    main()
