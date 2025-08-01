from pptx import Presentation

def extract_text_shapes(shapes, slide_idx, source, source_index):
    metadata = []
    for shape_idx, shape in enumerate(shapes):
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                entry = {
                    "source": source,
                    "source_index": source_index,
                    "slide": slide_idx,
                    "shape": shape_idx,
                    "context": sanitize_context(txt),  # Clean context
                    "char_count": len(txt),
                    "field_type": classify_field_type(txt, shape)
                }
                metadata.append(entry)
    return metadata

def sanitize_context(text):
    """Remove template-specific phrases that might confuse the AI"""
    template_phrases = [
        "Slidesgo", "template", "presentation", "here's what you'll find",
        "below is the content", "visit", "read more", "click here"
    ]
    clean_text = text
    for phrase in template_phrases:
        clean_text = clean_text.replace(phrase, "").strip()
    return clean_text if clean_text else "content"

def classify_field_type(text, shape):
    """Classify the type of field based on text length and position"""
    text_len = len(text)
    if text_len < 30:
        return "title"
    elif text_len < 80:
        return "subtitle" 
    else:
        return "body"

def is_thank_you_slide(slide):
    """Detect if a slide is a thank you/closing slide"""
    thank_you_keywords = [
        'thank you', 'thanks', 'questions', 'contact', 'reach out',
        'get in touch', 'discussion', 'q&a', 'questions?', 'any questions',
        'contact us', 'contact me', 'email', '@', 'phone', 'linkedin'
    ]
    
    # Collect all text from the slide
    all_text = ""
    for shape in slide.shapes:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            all_text += " " + shape.text_frame.text.lower()
    
    # Check if any thank you keywords are present
    return any(keyword in all_text for keyword in thank_you_keywords)

def extract_all_metadata(pptx_path):
    prs = Presentation(pptx_path)
    metadata = []
    metadata += extract_text_shapes(prs.slide_master.shapes, -1, "master", 0)
    for layout_idx, layout in enumerate(prs.slide_layouts):
        metadata += extract_text_shapes(layout.shapes, -1, "layout", layout_idx)
    for slide_idx, slide in enumerate(prs.slides):
        if is_thank_you_slide(slide):
            print(f"🔒 Skipping slide {slide_idx + 1} (detected as thank you/closing slide)")
            continue
        metadata += extract_text_shapes(slide.shapes, slide_idx, "slide", slide_idx)
    return metadata

def extract_all_metadata(template_path):
    """
    Extract metadata for all slides in the PowerPoint template.
    Returns a list of dicts, one per slide, with field info.
    """
    prs = Presentation(template_path)
    all_metadata = []
    for i, slide in enumerate(prs.slides):
        slide_metadata = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            char_count = len(text)
            context = shape.name if hasattr(shape, 'name') else ''
            positioning = {
                'top': getattr(shape, 'top', 0),
                'left': getattr(shape, 'left', 0)
            }
            slide_metadata.append({
                'slide': i,
                'shape': shape.shape_id,
                'text': text,
                'char_count': char_count,
                'context': context,
                'positioning': positioning
            })
        all_metadata.append(slide_metadata)
    return all_metadata

def group_by_slide(all_metadata):
    """
    Group metadata by slide number.
    Returns a list of (slide_num, slide_metadata) tuples.
    """
    return [(i, slide_metadata) for i, slide_metadata in enumerate(all_metadata)]

def classify_slide_layout(slide_metadata):
    """
    Classify slide layout based on field characteristics.
    """
    field_count = len(slide_metadata)
    char_counts = [item.get("char_count", 0) for item in slide_metadata]
    has_title = any(count < 50 for count in char_counts)
    has_body = any(count > 100 for count in char_counts)
    if field_count == 1:
        return "title_slides"
    elif field_count == 2 and has_title and has_body:
        return "content_slides"
    elif field_count >= 3 and any(count < 80 for count in char_counts):
        return "list_slides"
    elif field_count >= 4 and has_title:
        return "comparison_slides"
    elif field_count >= 5:
        return "table_slides"
    else:
        return "content_slides"

def analyze_template_inventory(template_path):
    """
    Analyze all available slides in the PowerPoint template and categorize them.
    Returns a dict of slide categories with metadata.
    """
    all_metadata = extract_all_metadata(template_path)
    slide_categories = {
        "title_slides": [],
        "content_slides": [],
        "list_slides": [],
        "comparison_slides": [],
        "table_slides": [],
        "timeline_slides": [],
        "overview_slides": [],
        "conclusion_slides": []
    }
    for slide_num, slide_metadata in group_by_slide(all_metadata):
        slide_type = classify_slide_layout(slide_metadata)
        slide_categories[slide_type].append({
            "slide_number": slide_num,
            "template_type": slide_type,
            "field_count": len(slide_metadata),
            "metadata": slide_metadata
        })
    return slide_categories

if __name__ == "__main__":
    import sys, json
    if len(sys.argv) < 2:
        print("Usage: python extract_metadata.py template.pptx")
    else:
        print(json.dumps(extract_all_metadata(sys.argv[1]), indent=2))
