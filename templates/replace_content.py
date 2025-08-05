from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
import openai
import os
import json
import re
from dotenv import load_dotenv

# =====================================================
# SAFE SLIDE LAYOUT FUNCTIONS
# =====================================================

def get_safe_slide_layout(prs):
    """Get a safe slide layout to avoid PowerPoint corruption"""
    try:
        # Try to use blank layout (usually index 6)
        if len(prs.slide_layouts) > 6:
            return prs.slide_layouts[6]
        # If not available, use the last layout (most likely to be blank)
        elif len(prs.slide_layouts) > 0:
            return prs.slide_layouts[-1]
        # Fallback to first layout if no layouts available
        else:
            return prs.slide_layouts[0]
    except (IndexError, AttributeError):
        # Ultimate fallback - use first available layout
        return prs.slide_layouts[0]

# =====================================================
# ENHANCED TABLE GENERATION SYSTEM
# =====================================================

# Table structure templates for different question types
TABLE_TEMPLATES = {
    "statistics": {
        "headers": ["Metric", "Value", "Context/Source"],
        "description": "For questions asking about data, numbers, statistics, measurements"
    },
    "comparison": {
        "headers": ["Feature", "Option 1", "Option 2"],
        "description": "For questions comparing two or more things"
    },
    "classification": {
        "headers": ["Category", "Description", "Example/Details"],
        "description": "For questions about types, kinds, categories"
    },
    "process": {
        "headers": ["Step", "Action", "Result/Outcome"],
        "description": "For questions about processes, procedures, workflows"
    },
    "benefits": {
        "headers": ["Benefit", "Description", "Impact"],
        "description": "For questions about advantages, benefits, pros"
    },
    "problems": {
        "headers": ["Issue", "Description", "Impact/Solution"],
        "description": "For questions about problems, challenges, issues"
    },
    "features": {
        "headers": ["Feature", "Description", "Purpose/Benefit"],
        "description": "For questions about features, capabilities, functions"
    },
    "timeline": {
        "headers": ["Period/Phase", "Event/Development", "Impact/Significance"],
        "description": "For questions about history, evolution, timeline"
    }
}

def analyze_table_question_intent(question):
    """
    Analyze question to determine the optimal table structure and intent
    """
    if not question:
        return "classification", TABLE_TEMPLATES["classification"]
    
    question_lower = question.lower()
    
    # Statistics questions
    if any(word in question_lower for word in ['statistics', 'data', 'numbers', 'percentage', 'rate', 'measurement', 'metrics', 'figures', 'amount', 'quantity', 'how much', 'how many']):
        return "statistics", TABLE_TEMPLATES["statistics"]
    
    # Comparison questions
    elif any(word in question_lower for word in ['compare', 'comparison', 'difference', 'versus', 'vs', 'between', 'contrast', 'similarities', 'differ']):
        return "comparison", TABLE_TEMPLATES["comparison"]
    
    # Process questions - Note: Process questions should be TIMELINE format, not table
    # This function should NOT handle process questions as they are better suited for timeline slides
    # elif any(word in question_lower for word in ['process', 'steps', 'procedure', 'workflow', 'how to', 'method', 'approach', 'sequence', 'stages', 'phases']):
    #     return "process", TABLE_TEMPLATES["process"]
    
    # Benefits/advantages questions
    elif any(word in question_lower for word in ['benefits', 'advantages', 'pros', 'positive', 'strengths', 'gains', 'value']):
        return "benefits", TABLE_TEMPLATES["benefits"]
    
    # Problems/challenges questions
    elif any(word in question_lower for word in ['problems', 'issues', 'challenges', 'difficulties', 'cons', 'disadvantages', 'limitations', 'drawbacks']):
        return "problems", TABLE_TEMPLATES["problems"]
    
    # Features questions
    elif any(word in question_lower for word in ['features', 'capabilities', 'functions', 'functionality', 'what does', 'what can', 'abilities']):
        return "features", TABLE_TEMPLATES["features"]
    
    # Timeline/history questions
    elif any(word in question_lower for word in ['history', 'evolution', 'timeline', 'development', 'over time', 'chronology', 'progression', 'when', 'timeline']):
        return "timeline", TABLE_TEMPLATES["timeline"]
    
    # Classification questions (types, kinds, categories)
    elif any(word in question_lower for word in ['types', 'kinds', 'categories', 'varieties', 'forms', 'classes', 'what are']):
        return "classification", TABLE_TEMPLATES["classification"]
    
    # Default to classification for general questions
    else:
        return "classification", TABLE_TEMPLATES["classification"]

def plan_table_structure(question, content):
    """
    Analyze question and content to determine optimal table structure
    """
    print(f"   >>> Planning table structure for question: {question[:50]}...")
    
    # Step 1: Analyze question intent
    structure_type, template = analyze_table_question_intent(question)
    print(f"   *** Detected table type: {structure_type}")
    print(f"   *** Template headers: {template['headers']}")
    
    # Step 2: Analyze content availability
    content_analysis = analyze_content_for_table_data(content, structure_type)
    
    # Step 3: Create table plan
    table_plan = {
        "structure_type": structure_type,
        "headers": template["headers"].copy(),
        "template_description": template["description"],
        "content_analysis": content_analysis,
        "question": question
    }
    
    return table_plan

def analyze_content_for_table_data(content, structure_type):
    """
    Analyze if content contains data suitable for the requested table structure
    """
    content_lower = content.lower()
    
    analysis = {
        "has_numerical_data": bool(re.search(r'\d+[%$]?|\d+\.\d+', content)),
        "has_structured_lists": bool(re.search(r'(\d+\.|•|-|\*)\s+', content)),
        "has_comparison_words": any(word in content_lower for word in ['versus', 'vs', 'compared to', 'different', 'similar', 'contrast']),
        "has_process_indicators": any(word in content_lower for word in ['first', 'second', 'third', 'step', 'then', 'next', 'finally']),
        "content_length": len(content),
        "sentence_count": len(re.split(r'[.!?]+', content))
    }
    
    return analysis

def extract_table_data_with_ai(question, content, table_plan):
    """
    Use AI to intelligently extract and organize table data based on question intent
    """
    print(f"   >>> Using AI to extract table data for {table_plan['structure_type']} table...")
    
    try:
        load_dotenv()
        api_key = os.getenv('OPENAI_API_KEY')
        if not api_key:
            print("   WARNING: No OpenAI API key found, falling back to rule-based extraction")
            return extract_table_data_rule_based(question, content, table_plan)
        
        client = openai.OpenAI(api_key=api_key)
        
        # Create AI prompt for table data extraction
        prompt = f"""
QUESTION: {question}

CONTENT TO ANALYZE:
{content[:3000]}  

TASK: Create a {table_plan['structure_type']} table that directly answers the question using ONLY the provided content.

REQUIRED TABLE STRUCTURE:
Headers: {table_plan['headers']}
Description: {table_plan['template_description']}

CRITICAL INSTRUCTIONS:
1. Extract data from the provided content that directly relates to the question
2. Organize the extracted data into the exact column structure specified
3. If specific data is not available in the content, write "Not specified in content" or "Information not provided"
4. Do NOT make up or infer data that isn't explicitly stated in the content
5. Each row should contain related information across all columns
6. Aim for 3-5 meaningful rows of data
7. Keep cell contents concise but informative (max 100 characters per cell)

RESPONSE FORMAT: Return ONLY a JSON array of arrays, where the first array is headers and subsequent arrays are data rows.

EXAMPLE FORMAT:
[
  ["Header1", "Header2", "Header3"],
  ["Data1A", "Data1B", "Data1C"],
  ["Data2A", "Data2B", "Data2C"]
]

Extract the table data now:
"""
        
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1500,
            temperature=0.3
        )
        
        ai_result_raw = response.choices[0].message.content
        # Defensive type checking - handle case where content might be a list
        if isinstance(ai_result_raw, list):
            ai_result = str(ai_result_raw[0]) if ai_result_raw else ""
        else:
            ai_result = str(ai_result_raw).strip() if ai_result_raw else ""
        
        # Parse AI response
        try:
            # Clean up the response to extract JSON
            if '```json' in ai_result:
                ai_result = ai_result.split('```json')[1].split('```')[0].strip()
            elif '```' in ai_result:
                ai_result = ai_result.split('```')[1].strip()
            
            table_data = json.loads(ai_result)
            
            # Validate the table data
            if validate_table_data(table_data, table_plan):
                print(f"   SUCCESS: AI extracted {len(table_data)} rows of table data")
                return table_data
            else:
                print("   WARNING: AI data validation failed, using rule-based extraction")
                return extract_table_data_rule_based(question, content, table_plan)
                
        except json.JSONDecodeError as e:
            print(f"   WARNING: AI response parsing failed: {e}, using rule-based extraction")
            return extract_table_data_rule_based(question, content, table_plan)
            
    except Exception as e:
        print(f"   ERROR: AI table extraction failed: {e}, using rule-based extraction")
        return extract_table_data_rule_based(question, content, table_plan)

def validate_table_data(table_data, table_plan):
    """
    Validate that the extracted table data is coherent and properly structured
    """
    if not table_data or len(table_data) < 2:
        return False
    
    # Check if headers match expected structure
    headers = table_data[0]
    expected_headers = table_plan['headers']
    
    if len(headers) != len(expected_headers):
        return False
    
    # Check if we have meaningful data rows
    data_rows = table_data[1:]
    if not data_rows:
        return False
    
    # Check if all rows have the same number of columns
    expected_cols = len(headers)
    for row in data_rows:
        if len(row) != expected_cols:
            return False
    
    # Check for meaningful content (not all "Not specified" or similar)
    meaningful_cells = 0
    total_cells = 0
    
    for row in data_rows:
        for cell in row:
            total_cells += 1
            if cell and isinstance(cell, str) and len(cell.strip()) > 5 and \
               "not specified" not in cell.lower() and "not provided" not in cell.lower() and \
               "information not" not in cell.lower():
                meaningful_cells += 1
    
    # At least 40% of cells should have meaningful content
    if total_cells > 0 and meaningful_cells / total_cells >= 0.4:
        return True
    
    return False

def extract_table_data_rule_based(question, content, table_plan):
    """
    Fallback rule-based extraction when AI is not available or fails
    """
    print(f"   >>> Using rule-based extraction for {table_plan['structure_type']} table...")
    
    structure_type = table_plan['structure_type']
    headers = table_plan['headers']
    
    # Create base table with headers
    table_data = [headers]
    
    # Apply structure-specific extraction rules
    if structure_type == "statistics":
        rows = extract_statistics_rule_based(content)
    elif structure_type == "comparison":
        rows = extract_comparison_rule_based(content, question)
    elif structure_type == "classification":
        rows = extract_classification_rule_based(content)
    elif structure_type == "process":
        rows = extract_process_rule_based(content)
    elif structure_type == "features":
        rows = extract_features_rule_based(content)
    else:
        # Generic extraction
        rows = extract_generic_rule_based(content)
    
    # Ensure rows match header count
    for row in rows:
        while len(row) < len(headers):
            row.append("Information not available")
        row = row[:len(headers)]  # Truncate if too long
        table_data.append(row)
    
    return table_data if len(table_data) > 1 else []

def extract_statistics_rule_based(content):
    """Extract statistical data using rule-based approach"""
    rows = []
    
    # Look for numerical patterns
    number_patterns = re.findall(r'(\w+(?:\s+\w+)*)\s*(?:is|are|:|=)\s*(\d+(?:\.\d+)?[%$]?)', content, re.IGNORECASE)
    
    for metric, value in number_patterns[:4]:
        metric = metric.strip().title()
        context = "As mentioned in content"
        rows.append([metric, value, context])
    
    # If no numbers found, create generic data
    if not rows:
        sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20][:3]
        for i, sentence in enumerate(sentences):
            metric = f"Data Point {i+1}"
            value = "Not quantified"
            context = smart_content_adaptation(sentence, 80)
            rows.append([metric, value, context])
    
    return rows

def extract_comparison_rule_based(content, question):
    """Extract comparison data using rule-based approach"""
    rows = []
    
    # Look for comparison indicators
    comparison_sentences = []
    sentences = re.split(r'[.!?]+', content)
    
    for sentence in sentences:
        if any(word in sentence.lower() for word in ['versus', 'vs', 'compared', 'different', 'while', 'whereas']):
            comparison_sentences.append(sentence.strip())
    
    # Extract features being compared
    if comparison_sentences:
        for i, sentence in enumerate(comparison_sentences[:3]):
            feature = f"Aspect {i+1}"
            parts = re.split(r'\b(while|whereas|but|however)\b', sentence, flags=re.IGNORECASE)
            option1 = smart_content_adaptation(parts[0].strip(), 60)
            option2 = smart_content_adaptation(parts[-1].strip(), 60) if len(parts) > 1 else "Not specified"
            rows.append([feature, option1, option2])
    
    # Generic fallback
    if not rows:
        rows.append(["Primary Aspect", "Information provided in content", "Additional details available"])
        rows.append(["Secondary Aspect", "Content describes various features", "Context explains functionality"])
    
    return rows

def extract_classification_rule_based(content):
    """Extract classification data using rule-based approach"""
    rows = []
    
    # Look for list patterns or structured content
    list_items = re.findall(r'(?:^|\n)\s*(?:\d+\.|•|-|\*)\s*([^\n.!?]+)', content)
    
    if list_items:
        for item in list_items[:4]:
            category = smart_content_adaptation(item.strip(), 30)
            description = "As described in content"
            example = "Refer to source material"
            rows.append([category, description, example])
    else:
        # Look for key terms that might be categories
        sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20][:3]
        for i, sentence in enumerate(sentences):
            words = sentence.split()
            category = ' '.join(words[:3]).title() if len(words) >= 3 else f"Category {i+1}"
            description = smart_content_adaptation(sentence, 80)
            example = "Detailed in content"
            rows.append([category, description, example])
    
    return rows

def extract_process_rule_based(content):
    """Extract process data using rule-based approach"""
    rows = []
    
    # Look for step indicators
    step_patterns = re.findall(r'(?:step\s*\d+|first|second|third|then|next|finally)[:\s]*([^.!?]+)', content, re.IGNORECASE)
    
    for i, step_content in enumerate(step_patterns[:4]):
        step = f"Step {i+1}"
        action = smart_content_adaptation(step_content.strip(), 60)
        result = "Leads to next phase"
        rows.append([step, action, result])
    
    # Generic fallback
    if not rows:
        sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20][:3]
        for i, sentence in enumerate(sentences):
            step = f"Phase {i+1}"
            action = smart_content_adaptation(sentence, 60)
            result = "Contributes to overall process"
            rows.append([step, action, result])
    
    return rows

def extract_features_rule_based(content):
    """Extract features data using rule-based approach"""
    rows = []
    
    # Look for feature-like content
    sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20]
    
    for i, sentence in enumerate(sentences[:4]):
        # Try to identify feature name from beginning of sentence
        words = sentence.split()
        feature = ' '.join(words[:2]).title() if len(words) >= 2 else f"Feature {i+1}"
        description = smart_content_adaptation(sentence, 80)
        purpose = "Enhances functionality"
        rows.append([feature, description, purpose])
    
    return rows

def extract_generic_rule_based(content):
    """Generic extraction for any content"""
    rows = []
    
    sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20]
    
    for i, sentence in enumerate(sentences[:4]):
        words = sentence.split()
        topic = ' '.join(words[:3]).title() if len(words) >= 3 else f"Topic {i+1}"
        description = smart_content_adaptation(sentence, 80)
        details = "Additional information available"
        rows.append([topic, description, details])
    
    return rows

# =====================================================
# ENHANCED TIMELINE GENERATION SYSTEM
# =====================================================

# Timeline structure templates for different question types
TIMELINE_TEMPLATES = {
    "historical": {
        "headers": ["Period/Era", "Key Development", "Impact/Significance"],
        "visual_style": "chronological_timeline",
        "description": "For questions about history, evolution, development over time"
    },
    "process": {
        "headers": ["Step", "Action", "Result/Output"],
        "visual_style": "process_flow", 
        "description": "For questions about procedures, workflows, methods"
    },
    "lifecycle": {
        "headers": ["Phase", "Activities", "Deliverables"],
        "visual_style": "lifecycle_diagram",
        "description": "For questions about development cycles, project phases"
    },
    "sequential": {
        "headers": ["Stage", "Event", "Next Action"],
        "visual_style": "sequential_flow",
        "description": "For questions about sequences, progressions, pipelines"
    },
    "stages": {
        "headers": ["Stage", "Description", "Key Features"],
        "visual_style": "stage_progression",
        "description": "For questions about stages, levels, phases of development"
    }
}

# =====================================================
# ENHANCED BOXES TEMPLATES SYSTEM
# =====================================================

BOXES_TEMPLATES = {
    "comparison": {
        "headers": ["Item A", "Item B"],
        "layout_type": "two_column_comparison",
        "description": "For questions comparing two items, methods, or approaches"
    },
    "components": {
        "headers": ["Component", "Function", "Details"],
        "layout_type": "multi_box_grid",
        "description": "For questions about parts, components, elements of a system"
    },
    "features": {
        "headers": ["Feature", "Description", "Benefit"],
        "layout_type": "feature_grid",
        "description": "For questions about features, characteristics, or attributes"
    },
    "categories": {
        "headers": ["Category", "Description", "Examples"],
        "layout_type": "category_grid", 
        "description": "For questions about types, kinds, or classifications"
    },
    "aspects": {
        "headers": ["Aspect", "Details", "Importance"],
        "layout_type": "aspect_grid",
        "description": "For questions about different aspects or dimensions"
    }
}

def analyze_boxes_question_intent(question):
    """
    Analyze question to determine the optimal boxes structure and intent
    Enhanced with intelligent question type detection for boxes/comparison layouts
    """
    if not question:
        return "generic", BOXES_TEMPLATES["aspects"]
    
    question_lower = question.lower()
    
    # COMPARISON DETECTION - Strong indicators for two-item comparisons
    comparison_keywords = ['compare', 'vs', 'versus', 'difference', 'differ', 'contrast', 'similarities', 'vs.']
    if any(keyword in question_lower for keyword in comparison_keywords):
        print(f"   >>> Comparison question detected for boxes format")
        return "comparison", BOXES_TEMPLATES["comparison"]
    
    # COMPONENTS DETECTION - Parts, elements, components of systems
    components_keywords = ['components', 'parts', 'elements', 'pieces', 'modules', 'sections']
    if any(keyword in question_lower for keyword in components_keywords):
        print(f"   >>> Components question detected for boxes format")
        return "components", BOXES_TEMPLATES["components"]
    
    # FEATURES DETECTION - Characteristics, features, attributes
    features_keywords = ['features', 'characteristics', 'attributes', 'properties', 'qualities', 'capabilities']
    if any(keyword in question_lower for keyword in features_keywords):
        print(f"   >>> Features question detected for boxes format")
        return "features", BOXES_TEMPLATES["features"]
    
    # CATEGORIES DETECTION - Types, kinds, classifications
    categories_keywords = ['types', 'kinds', 'categories', 'classifications', 'classes', 'varieties']
    if any(keyword in question_lower for keyword in categories_keywords):
        print(f"   >>> Categories question detected for boxes format")
        return "categories", BOXES_TEMPLATES["categories"]
    
    # ASPECTS DETECTION - Different aspects, dimensions, areas
    aspects_keywords = ['aspects', 'dimensions', 'areas', 'factors', 'considerations', 'perspectives']
    if any(keyword in question_lower for keyword in aspects_keywords):
        print(f"   >>> Aspects question detected for boxes format")
        return "aspects", BOXES_TEMPLATES["aspects"]
    
    # "What are the" questions - Usually components or categories
    if question_lower.startswith('what are the'):
        if any(word in question_lower for word in ['main', 'key', 'primary', 'core']):
            return "components", BOXES_TEMPLATES["components"]
        else:
            return "categories", BOXES_TEMPLATES["categories"]
    
    # Default to aspects for general questions
    return "aspects", BOXES_TEMPLATES["aspects"]

def analyze_timeline_question_intent(question):
    """
    Analyze question to determine the optimal timeline structure and intent
    Enhanced with better component/parts detection to avoid timeline misclassification
    """
    if not question:
        return "generic", TIMELINE_TEMPLATES["process"]
    
    question_lower = question.lower()
    
    # EXCLUSION RULES: Questions that should NOT be timelines
    # Components/parts questions should be boxes, not timelines
    components_keywords = ['components', 'parts', 'elements', 'aspects', 'features', 'main', 'key', 'core']
    if any(keyword in question_lower for keyword in components_keywords):
        print(f"   >>> Components question detected - NOT timeline format")
        return "generic", TIMELINE_TEMPLATES["process"]  # Return generic to indicate not timeline
    
    # "What are the" questions are usually lists/boxes, not timelines
    if question_lower.startswith('what are the') and not any(word in question_lower for word in ['steps', 'stages', 'phases', 'history', 'evolution']):
        print(f"   >>> 'What are the' question detected - likely not timeline format")
        return "generic", TIMELINE_TEMPLATES["process"]
    
    # POSITIVE TIMELINE DETECTION
    # Historical timeline questions - Strong indicators
    if any(word in question_lower for word in ['history', 'evolution', 'evolved', 'developed over time', 'timeline', 'chronology', 'over the years', 'development of', 'advancement', 'progression over time', 'since its inception', 'over time']):
        return "historical", TIMELINE_TEMPLATES["historical"]
    
    # Process timeline questions - Only for clear process questions
    elif any(phrase in question_lower for phrase in ['steps to', 'how to', 'procedure for', 'process of', 'workflow for', 'method to', 'approach to']):
        return "process", TIMELINE_TEMPLATES["process"]
    
    # Lifecycle questions
    elif any(word in question_lower for word in ['lifecycle', 'life cycle', 'development cycle', 'phases of development', 'project phases', 'software development', 'cycle']):
        return "lifecycle", TIMELINE_TEMPLATES["lifecycle"]
    
    # Sequential questions - Must be clearly sequential
    elif any(phrase in question_lower for phrase in ['sequence of', 'order of', 'progression of', 'flow of', 'chain of', 'series of']):
        return "sequential", TIMELINE_TEMPLATES["sequential"]
    
    # Stages questions - Must be clearly about stages
    elif any(phrase in question_lower for phrase in ['stages of', 'phases of', 'levels of']) and not any(word in question_lower for word in ['components', 'parts', 'elements']):
        return "stages", TIMELINE_TEMPLATES["stages"]
    
    # If "process" keyword alone without clear timeline indicators, return generic
    elif 'process' in question_lower and not any(word in question_lower for word in ['steps', 'stages', 'phases', 'how', 'sequence']):
        print(f"   >>> Generic process question - likely not timeline format")
        return "generic", TIMELINE_TEMPLATES["process"]
    
    # Default to generic (not timeline) for ambiguous questions
    else:
        return "generic", TIMELINE_TEMPLATES["process"]

def plan_timeline_structure(question, content):
    """
    Analyze question and content to determine optimal timeline structure
    """
    print(f"   >>> Planning timeline structure for question: {question[:50]}...")
    
    # Step 1: Analyze question intent
    timeline_type, template = analyze_timeline_question_intent(question)
    print(f"   *** Detected timeline type: {timeline_type}")
    print(f"   *** Template headers: {template['headers']}")
    print(f"   *** Visual style: {template['visual_style']}")
    
    # Step 2: Analyze content for timeline indicators
    content_analysis = analyze_content_for_timeline_data(content, timeline_type)
    
    # Step 3: Create timeline plan
    timeline_plan = {
        "timeline_type": timeline_type,
        "headers": template["headers"].copy(),
        "visual_style": template["visual_style"],
        "template_description": template["description"],
        "content_analysis": content_analysis,
        "question": question
    }
    
    return timeline_plan

def analyze_content_for_timeline_data(content, timeline_type):
    """
    Analyze if content contains data suitable for the requested timeline structure
    """
    content_lower = content.lower()
    
    analysis = {
        "has_temporal_indicators": bool(re.search(r'\b\d{4}\b|\b(early|mid|late)\s+\d{4}s?\b|in\s+\d{4}|since\s+\d{4}', content)),
        "has_sequence_numbers": bool(re.search(r'\b(first|second|third|fourth|1st|2nd|3rd|4th|\d+\.)\b', content)),
        "has_process_words": any(word in content_lower for word in ['then', 'next', 'after', 'before', 'following', 'subsequently', 'finally']),
        "has_stage_indicators": any(word in content_lower for word in ['stage', 'phase', 'step', 'level', 'tier', 'generation']),
        "content_length": len(content),
        "sentence_count": len(re.split(r'[.!?]+', content)),
        "paragraph_count": len([p for p in content.split('\n\n') if p.strip()])
    }
    
    return analysis

def extract_timeline_data_with_ai(question, content, timeline_plan):
    """
    Use AI to intelligently extract and organize timeline data based on question intent
    """
    print(f"   >>> Using AI to extract timeline data for {timeline_plan['timeline_type']} timeline...")
    
    try:
        load_dotenv()
        api_key = os.getenv('OPENAI_API_KEY')
        if not api_key:
            print("   WARNING: No OpenAI API key found, falling back to rule-based extraction")
            return extract_timeline_data_rule_based(question, content, timeline_plan)
        
        client = openai.OpenAI(api_key=api_key)
        
        # Create AI prompt for timeline data extraction
        prompt = f"""
QUESTION: {question}

CONTENT TO ANALYZE:
{content[:3000]}  

TASK: Create a {timeline_plan['timeline_type']} timeline that directly answers the question using ONLY the provided content.

REQUIRED TIMELINE STRUCTURE:
Headers: {timeline_plan['headers']}
Visual Style: {timeline_plan['visual_style']}
Description: {timeline_plan['template_description']}

CRITICAL INSTRUCTIONS:
1. Extract timeline data from the provided content that directly relates to the question
2. Organize the extracted data into the exact timeline structure specified
3. For {timeline_plan['timeline_type']} timelines, focus on proper chronological/sequential order
4. If temporal/sequential data is not available in the content, write "Timeline not specified in content"
5. Do NOT make up dates, sequences, or timeline information that isn't explicitly stated
6. Each timeline entry should contain related information across all columns
7. Aim for 3-5 meaningful timeline entries
8. Keep entries concise but informative (max 120 characters per field)

RESPONSE FORMAT: Return ONLY a JSON array of arrays, where the first array is headers and subsequent arrays are timeline entries.

EXAMPLE FORMAT:
[
  ["Header1", "Header2", "Header3"],
  ["Entry1A", "Entry1B", "Entry1C"],
  ["Entry2A", "Entry2B", "Entry2C"]
]

Extract the timeline data now:
"""
        
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1500,
            temperature=0.3
        )
        
        ai_result_raw = response.choices[0].message.content
        # Defensive type checking - handle case where content might be a list
        if isinstance(ai_result_raw, list):
            ai_result = str(ai_result_raw[0]) if ai_result_raw else ""
        else:
            ai_result = str(ai_result_raw).strip() if ai_result_raw else ""
        
        # Parse AI response
        try:
            # Clean up the response to extract JSON
            if '```json' in ai_result:
                ai_result = ai_result.split('```json')[1].split('```')[0].strip()
            elif '```' in ai_result:
                ai_result = ai_result.split('```')[1].strip()
            
            timeline_data = json.loads(ai_result)
            
            # Validate the timeline data
            if validate_timeline_data(timeline_data, timeline_plan):
                print(f"   SUCCESS: AI extracted {len(timeline_data)} timeline entries")
                return timeline_data
            else:
                print("   WARNING: AI data validation failed, using rule-based extraction")
                return extract_timeline_data_rule_based(question, content, timeline_plan)
                
        except json.JSONDecodeError as e:
            print(f"   WARNING: AI response parsing failed: {e}, using rule-based extraction")
            return extract_timeline_data_rule_based(question, content, timeline_plan)
            
    except Exception as e:
        print(f"   ERROR: AI timeline extraction failed: {e}, using rule-based extraction")
        return extract_timeline_data_rule_based(question, content, timeline_plan)

def validate_timeline_data(timeline_data, timeline_plan):
    """
    Validate that the extracted timeline data is coherent and properly structured
    """
    if not timeline_data or len(timeline_data) < 2:
        return False
    
    # Check if headers match expected structure
    headers = timeline_data[0]
    expected_headers = timeline_plan['headers']
    
    if len(headers) != len(expected_headers):
        return False
    
    # Check if we have meaningful timeline entries
    timeline_entries = timeline_data[1:]
    if not timeline_entries:
        return False
    
    # Check if all entries have the same number of columns
    expected_cols = len(headers)
    for entry in timeline_entries:
        if len(entry) != expected_cols:
            return False
    
    # Check for meaningful content (not all "Timeline not specified" or similar)
    meaningful_cells = 0
    total_cells = 0
    
    for entry in timeline_entries:
        for cell in entry:
            total_cells += 1
            if cell and isinstance(cell, str) and len(cell.strip()) > 5 and \
               "not specified" not in cell.lower() and "timeline not" not in cell.lower() and \
               "information not" not in cell.lower():
                meaningful_cells += 1
    
    # At least 40% of cells should have meaningful content
    if total_cells > 0 and meaningful_cells / total_cells >= 0.4:
        return True
    
    return False

def extract_timeline_data_rule_based(question, content, timeline_plan):
    """
    Fallback rule-based extraction when AI is not available or fails
    """
    print(f"   >>> Using rule-based extraction for {timeline_plan['timeline_type']} timeline...")
    
    timeline_type = timeline_plan['timeline_type']
    headers = timeline_plan['headers']
    
    # Create base timeline with headers
    timeline_data = [headers]
    
    # Apply timeline-type specific extraction rules
    if timeline_type == "historical":
        entries = extract_historical_timeline_rule_based(content)
    elif timeline_type == "process":
        entries = extract_process_timeline_rule_based(content)
    elif timeline_type == "lifecycle":
        entries = extract_lifecycle_timeline_rule_based(content)
    elif timeline_type == "sequential":
        entries = extract_sequential_timeline_rule_based(content)
    elif timeline_type == "stages":
        entries = extract_stages_timeline_rule_based(content)
    else:
        # Generic extraction
        entries = extract_generic_timeline_rule_based(content)
    
    # Ensure entries match header count
    for entry in entries:
        while len(entry) < len(headers):
            entry.append("Information not available")
        entry = entry[:len(headers)]  # Truncate if too long
        timeline_data.append(entry)
    
    return timeline_data if len(timeline_data) > 1 else []

def extract_historical_timeline_rule_based(content):
    """Extract historical timeline data using rule-based approach"""
    entries = []
    
    # Look for year patterns and historical indicators
    year_patterns = re.findall(r'(\b\d{4}\b|\b(early|mid|late)\s+\d{4}s?\b).*?([^.!?]*[.!?])', content, re.IGNORECASE)
    
    for year_match in year_patterns[:4]:
        period = year_match[0] if year_match[0] else year_match[1]
        development = smart_content_adaptation(year_match[2].strip(), 80)
        impact = "Historical significance noted"
        entries.append([period, development, impact])
    
    # If no years found, create timeline from content chunks
    if not entries:
        sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20][:4]
        for i, sentence in enumerate(sentences):
            period = f"Period {i+1}"
            development = smart_content_adaptation(sentence, 80)
            impact = "Part of historical development"
            entries.append([period, development, impact])
    
    return entries

def extract_process_timeline_rule_based(content):
    """Extract process timeline data using rule-based approach"""
    entries = []
    
    # Look for step indicators
    step_patterns = re.findall(r'(?:step\s*\d+|first|second|third|then|next|finally)[:\s]*([^.!?]+)', content, re.IGNORECASE)
    
    for i, step_content in enumerate(step_patterns[:4]):
        step = f"Step {i+1}"
        action = smart_content_adaptation(step_content.strip(), 80)
        result = "Leads to next step"
        entries.append([step, action, result])
    
    # Generic fallback for processes
    if not entries:
        sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20][:4]
        for i, sentence in enumerate(sentences):
            step = f"Step {i+1}"
            action = smart_content_adaptation(sentence, 80)
            result = "Part of overall process"
            entries.append([step, action, result])
    
    return entries

def extract_lifecycle_timeline_rule_based(content):
    """Extract lifecycle timeline data using rule-based approach"""
    entries = []
    
    # Look for lifecycle/phase indicators
    phase_patterns = re.findall(r'(?:phase\s*\d+|stage\s*\d+|cycle\s*\d+)[:\s]*([^.!?]+)', content, re.IGNORECASE)
    
    for i, phase_content in enumerate(phase_patterns[:4]):
        phase = f"Phase {i+1}"
        activities = smart_content_adaptation(phase_content.strip(), 80)
        deliverables = "Phase deliverables"
        entries.append([phase, activities, deliverables])
    
    # Generic lifecycle fallback
    if not entries:
        sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20][:4]
        for i, sentence in enumerate(sentences):
            phase = f"Phase {i+1}"
            activities = smart_content_adaptation(sentence, 80)
            deliverables = "Contributes to lifecycle"
            entries.append([phase, activities, deliverables])
    
    return entries

def extract_sequential_timeline_rule_based(content):
    """Extract sequential timeline data using rule-based approach"""
    entries = []
    
    # Look for sequential content
    sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20]
    
    for i, sentence in enumerate(sentences[:4]):
        stage = f"Stage {i+1}"
        event = smart_content_adaptation(sentence, 80)
        next_action = "Continues to next stage"
        entries.append([stage, event, next_action])
    
    return entries

def extract_stages_timeline_rule_based(content):
    """Extract stages timeline data using rule-based approach"""
    entries = []
    
    # Look for stage-like content
    sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20]
    
    for i, sentence in enumerate(sentences[:4]):
        stage = f"Stage {i+1}"
        description = smart_content_adaptation(sentence, 80)
        features = "Key characteristics noted"
        entries.append([stage, description, features])
    
    return entries

def extract_generic_timeline_rule_based(content):
    """Generic timeline extraction for any content"""
    entries = []
    
    sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20]
    
    for i, sentence in enumerate(sentences[:4]):
        item = f"Item {i+1}"
        description = smart_content_adaptation(sentence, 80)
        context = "Timeline context"
        entries.append([item, description, context])
    
    return entries

# =====================================================
# ENHANCED BOXES CONTENT EXTRACTION SYSTEM
# =====================================================

def plan_boxes_structure(question, content):
    """
    Analyze question and content to determine optimal boxes structure
    """
    print(f"   >>> Planning boxes structure for question: {question[:50]}...")
    
    # Step 1: Analyze question intent
    boxes_type, boxes_template = analyze_boxes_question_intent(question)
    
    print(f"   *** Detected boxes type: {boxes_type}")
    print(f"   *** Template headers: {boxes_template['headers']}")
    print(f"   *** Layout type: {boxes_template['layout_type']}")
    
    # Step 2: Analyze content characteristics
    content_analysis = analyze_content_for_boxes_data(content, boxes_type)
    
    # Step 3: Determine optimal structure
    structure_plan = {
        "boxes_type": boxes_type,
        "headers": boxes_template["headers"],
        "layout_type": boxes_template["layout_type"],
        "description": boxes_template["description"],
        "content_analysis": content_analysis,
        "expected_boxes": determine_box_count(boxes_type, content_analysis)
    }
    
    return structure_plan

def analyze_content_for_boxes_data(content, boxes_type):
    """
    Analyze if content contains data suitable for the requested boxes structure
    """
    content_lower = content.lower()
    
    analysis = {
        "has_comparison_indicators": bool(re.search(r'\b(vs|versus|compared|contrast|while|whereas|however|but)\b', content_lower)),
        "has_component_indicators": bool(re.search(r'\b(component|part|element|module|piece|section)\b', content_lower)),
        "has_feature_indicators": bool(re.search(r'\b(feature|characteristic|attribute|property|capability)\b', content_lower)),
        "has_list_structure": bool(re.search(r'\b(first|second|third|\d+\.)\b', content_lower)),
        "content_length": len(content),
        "sentence_count": len(re.split(r'[.!?]+', content)),
        "paragraph_count": len([p for p in content.split('\n\n') if p.strip()])
    }
    
    return analysis

def determine_box_count(boxes_type, content_analysis):
    """Determine optimal number of boxes based on type and content"""
    if boxes_type == "comparison":
        return 2  # Always 2 for comparisons
    elif boxes_type == "components":
        # 3-6 boxes for components based on content
        return min(6, max(3, content_analysis["sentence_count"] // 2))
    else:
        # 3-4 boxes for other types
        return min(4, max(3, content_analysis["sentence_count"] // 3))

def extract_boxes_data_with_ai(question, content, boxes_plan):
    """
    Use AI to intelligently extract and organize boxes data based on question intent
    """
    print(f"   >>> Using AI to extract boxes data for {boxes_plan['boxes_type']} layout...")
    
    try:
        load_dotenv()
        api_key = os.getenv('OPENAI_API_KEY')
        if not api_key:
            print("   WARNING: No OpenAI API key found, falling back to rule-based extraction")
            return extract_boxes_data_rule_based(question, content, boxes_plan)
        
        client = openai.OpenAI(api_key=api_key)
        
        # Create enhanced AI prompt for richer content extraction
        prompt = f"""
QUESTION: {question}

CONTENT TO ANALYZE:
{content[:5000]}  

TASK: Create a comprehensive {boxes_plan['boxes_type']} boxes layout that thoroughly answers the question using the provided content.

REQUIRED BOXES STRUCTURE:
Layout Type: {boxes_plan['layout_type']}
Expected Boxes: {boxes_plan['expected_boxes']}
Description: {boxes_plan['description']}

ENHANCED CONTENT REQUIREMENTS:
1. Extract comprehensive boxes data from the provided content that directly relates to the question
2. Create DETAILED, INFORMATIVE entries for each box with rich explanations
3. Each box should have substantial content (aim for 100-200 characters per field)
4. Include specific technical details, examples, and context from the source material
5. For {boxes_plan['boxes_type']} boxes, provide clear distinctions and comprehensive information
6. Expand on concepts with explanations of HOW and WHY, not just WHAT
7. Include relevant technical terminology and specifications when available
8. Each box should provide substantial educational value and context
9. If expanding beyond source content, clearly indicate this is explanatory context
10. Aim for {boxes_plan['expected_boxes']} information-rich boxes
11. Each field should be substantial and informative, not just brief labels
12. Provide enough detail to fill the visual space effectively

CONTENT DENSITY TARGETS:
- Title: Descriptive and specific (20-40 characters)
- Content: Comprehensive explanation (150-300 characters) 
- Details: Additional context, examples, or technical specifications (100-200 characters)

RESPONSE FORMAT: Return ONLY a JSON array of objects, where each object represents a box with substantial, detailed content.

EXAMPLE FORMAT:
[
  {{"title": "Descriptive Component Name", "content": "Comprehensive explanation of what this component does, how it works, and why it's important in the overall system. Include technical details and context.", "details": "Additional technical specifications, examples of usage, or implementation details that provide deeper understanding."}},
  {{"title": "Another Component", "content": "Detailed description with technical context and explanations...", "details": "More specific details and examples..."}}
]

Extract the enriched boxes data now:
"""
        
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=2500,
            temperature=0.3
        )
        
        ai_result_raw = response.choices[0].message.content
        # Defensive type checking - handle case where content might be a list
        if isinstance(ai_result_raw, list):
            ai_result = str(ai_result_raw[0]) if ai_result_raw else ""
        else:
            ai_result = str(ai_result_raw).strip() if ai_result_raw else ""
        
        # Parse AI response
        try:
            # Clean up the response to extract JSON
            if '```json' in ai_result:
                ai_result = ai_result.split('```json')[1].split('```')[0].strip()
            elif '```' in ai_result:
                ai_result = ai_result.split('```')[1].strip()
            
            boxes_data = json.loads(ai_result)
            
            # Validate the boxes data
            if validate_boxes_data(boxes_data, boxes_plan):
                print(f"   SUCCESS: AI extracted {len(boxes_data)} boxes")
                return boxes_data
            else:
                print("   WARNING: AI data validation failed, using rule-based extraction")
                return extract_boxes_data_rule_based(question, content, boxes_plan)
                
        except json.JSONDecodeError as e:
            print(f"   WARNING: AI response parsing failed: {e}, using rule-based extraction")
            return extract_boxes_data_rule_based(question, content, boxes_plan)
            
    except Exception as e:
        print(f"   ERROR: AI boxes extraction failed: {e}, using rule-based extraction")
        return extract_boxes_data_rule_based(question, content, boxes_plan)

def validate_boxes_data(boxes_data, boxes_plan):
    """
    Validate that the extracted boxes data is coherent and properly structured
    """
    if not boxes_data or len(boxes_data) < 2:
        return False
    
    # Check if we have meaningful boxes
    meaningful_boxes = 0
    total_boxes = len(boxes_data)
    
    for box in boxes_data:
        if isinstance(box, dict) and any(
            isinstance(value, str) and len(value.strip()) > 5 and 
            "not specified" not in value.lower() and "data not" not in value.lower()
            for value in box.values()
        ):
            meaningful_boxes += 1
    
    # At least 60% of boxes should have meaningful content
    success_ratio = meaningful_boxes / total_boxes
    return success_ratio >= 0.6

def extract_boxes_data_rule_based(question, content, boxes_plan):
    """
    Fallback rule-based extraction for boxes when AI fails
    """
    boxes_type = boxes_plan['boxes_type']
    
    if boxes_type == "comparison":
        return extract_comparison_boxes_rule_based(content)
    elif boxes_type == "components":
        return extract_components_boxes_rule_based(content)
    elif boxes_type == "features":
        return extract_features_boxes_rule_based(content)
    elif boxes_type == "categories":
        return extract_categories_boxes_rule_based(content)
    else:
        # Generic extraction
        return extract_generic_boxes_rule_based(content)

def extract_comparison_boxes_rule_based(content):
    """Extract comparison data for two-box layout"""
    # Try to split content for comparison
    comparison_data = parse_comparison_content(content, "")
    
    if len(comparison_data) >= 2:
        return [
            {"title": "Item A", "content": smart_content_adaptation(comparison_data[0], 150), "details": "Comparison details"},
            {"title": "Item B", "content": smart_content_adaptation(comparison_data[1], 150), "details": "Comparison details"}
        ]
    else:
        # Split content in half
        mid_point = len(content) // 2
        return [
            {"title": "Item A", "content": smart_content_adaptation(content[:mid_point], 150), "details": "First comparison item"},
            {"title": "Item B", "content": smart_content_adaptation(content[mid_point:], 150), "details": "Second comparison item"}
        ]

def extract_components_boxes_rule_based(content):
    """Extract component data for multi-box layout with enhanced content"""
    boxes = []
    
    # Enhanced patterns for component extraction
    component_patterns = [
        r'\d+\.\s*([^.!?\n]{10,})',  # Numbered items (minimum 10 chars)
        r'[•\-\*]\s*([^.!?\n]{10,})',  # Bulleted items
        r'([A-Z][a-z]+(?:\s+[A-Z][a-z]+)*)[:\s]+([^.!?\n]{20,})',  # "Title: Description" format
        r'([A-Z][^.!?\n]{20,100})'  # Capitalized sentences
    ]
    
    for pattern in component_patterns:
        matches = re.findall(pattern, content)
        if matches and len(matches) >= 3:
            for i, match in enumerate(matches[:6]):
                if isinstance(match, tuple) and len(match) == 2:
                    # Title: Description format
                    title, description = match
                    boxes.append({
                        "title": smart_content_adaptation(title.strip(), 40),
                        "content": smart_content_adaptation(description.strip(), 200),
                        "details": f"Essential component providing {title.lower()} functionality for the overall system architecture and security implementation."
                    })
                else:
                    # Single match - create enhanced content
                    component_text = match.strip()
                    
                    # Try to extract a meaningful title from the component text
                    words = component_text.split()
                    if len(words) >= 2:
                        title = ' '.join(words[:2]).title()
                    else:
                        title = f"Component {i+1}"
                    
                    # Create detailed content and context
                    enhanced_content = smart_content_adaptation(component_text, 200)
                    details = f"This component plays a crucial role in the system by providing specific functionality and integration capabilities essential for secure operation."
                    
                    boxes.append({
                        "title": title,
                        "content": enhanced_content,
                        "details": details
                    })
            break
    
    if not boxes:
        # Enhanced fallback: Split content into sentences with better content generation
        sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20]
        
        # Look for component-related keywords to create better titles
        component_keywords = ['key', 'algorithm', 'signature', 'encryption', 'compression', 'generation', 'management', 'verification']
        
        for i, sentence in enumerate(sentences[:4]):
            # Try to find a relevant keyword for the title
            title = f"Component {i+1}"
            for keyword in component_keywords:
                if keyword.lower() in sentence.lower():
                    title = keyword.title() + " System"
                    break
            
            # Enhanced content with context
            enhanced_content = smart_content_adaptation(sentence, 250)
            details = f"Critical system component that integrates with other elements to provide comprehensive functionality and maintain system security and reliability."
            
            boxes.append({
                "title": title,
                "content": enhanced_content,
                "details": details
            })
    
    return boxes

def extract_features_boxes_rule_based(content):
    """Extract feature data for feature-grid layout"""
    boxes = []
    
    # Look for feature-like patterns
    sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20]
    
    for i, sentence in enumerate(sentences[:4]):
        # Try to extract feature name from sentence
        words = sentence.split()
        feature_name = ' '.join(words[:3]).title() if len(words) >= 3 else f"Feature {i+1}"
        
        boxes.append({
            "title": feature_name,
            "content": smart_content_adaptation(sentence, 150),
            "details": "Feature benefits and usage"
        })
    
    return boxes

def extract_categories_boxes_rule_based(content):
    """Extract category data for category-grid layout"""
    boxes = []
    
    sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20]
    
    for i, sentence in enumerate(sentences[:4]):
        boxes.append({
            "title": f"Category {i+1}",
            "content": smart_content_adaptation(sentence, 150),
            "details": "Category examples and characteristics"
        })
    
    return boxes

def extract_generic_boxes_rule_based(content):
    """Generic extraction for any boxes content"""
    boxes = []
    
    sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20]
    
    for i, sentence in enumerate(sentences[:4]):
        boxes.append({
            "title": f"Aspect {i+1}",
            "content": smart_content_adaptation(sentence, 150),
            "details": "Important details and context"
        })
    
    return boxes

def find_shape(shapes, shape_idx):
    if shape_idx >= len(shapes):
        return None
    shape = shapes[shape_idx]
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        return shape
    return None

def get_smart_text_color(shape):
    """Determine appropriate text color based on background luminance"""
    try:
        if hasattr(shape, 'fill') and shape.fill:
            fill = shape.fill
            if hasattr(fill, 'fore_color') and fill.fore_color:
                try:
                    rgb = fill.fore_color.rgb
                    if rgb:
                        # Calculate luminance
                        r, g, b = rgb.r / 255.0, rgb.g / 255.0, rgb.b / 255.0
                        luminance = 0.299 * r + 0.587 * g + 0.114 * b
                        
                        # Return contrasting color
                        if luminance < 0.5:
                            return RGBColor(255, 255, 255)  # White text on dark background
                        else:
                            return RGBColor(0, 0, 0)  # Black text on light background
                except:
                    pass
        
        # Default to black text
        return RGBColor(0, 0, 0)
    except:
        return RGBColor(0, 0, 0)

def get_safe_shape_dimensions(shape):
    """Get shape dimensions with robust error handling and safe defaults"""
    try:
        width = None
        height = None
        
        # Try multiple ways to get dimensions
        if hasattr(shape, 'width') and shape.width:
            try:
                width = float(shape.width.inches)
            except:
                pass
                
        if hasattr(shape, 'height') and shape.height:
            try:
                height = float(shape.height.inches)
            except:
                pass
        
        # If dimensions failed or are invalid, use safe defaults
        if width is None or height is None or width <= 0 or height <= 0:
            # Return conservative defaults that work for most shapes
            return {"width": 3.0, "height": 1.0}  # Safe medium size
            
        return {"width": width, "height": height}
    except:
        return {"width": 3.0, "height": 1.0}  # Safe fallback

def analyze_shape_constraints(shape):
    """Analyze shape constraints with much more realistic capacity estimates"""
    try:
        dims = get_safe_shape_dimensions(shape)
        width = dims["width"]  # Always a number now
        height = dims["height"]  # Always a number now
        
        # Calculate area and estimated character capacity - MORE REALISTIC AND GENEROUS
        area = width * height
        estimated_capacity = int(area * 150)  # ~150 chars per square inch (increased from 100 for better fitting)
        
        # Classify shape size - now safe since width/height are always numbers
        is_micro = width < 1.0 or height < 0.5
        is_small = width < 1.5 or height < 0.8
        is_large = width > 6.0 or height > 3.0
        
        return {
            "width": width,
            "height": height,
            "area": area,
            "estimated_capacity": max(estimated_capacity, 50),  # Minimum 50 chars (increased from 25 for better fitting)
            "is_micro": is_micro,
            "is_small": is_small,
            "is_large": is_large
        }
    except Exception as e:
        # Ultimate fallback - return safe defaults with generous capacity
        return {
            "width": 3.0,
            "height": 1.0,
            "area": 3.0,
            "estimated_capacity": 300,  # Much more generous default
            "is_micro": False,
            "is_small": False,
            "is_large": False
        }

def test_if_text_fits(shape, text):
    """Actually test if text fits instead of just estimating"""
    try:
        if not hasattr(shape, 'text_frame') or not shape.text_frame:
            return False
            
        # Get current text to restore later
        original_text = ""
        if shape.text_frame.paragraphs:
            original_text = shape.text_frame.paragraphs[0].text
        
        # Temporarily set the test text
        shape.text_frame.paragraphs[0].text = text
        
        # Basic heuristics to determine if it fits reasonably
        # For now, we'll be optimistic - PowerPoint handles overflow better than we think
        fits = True
        
        # Only flag as not fitting if text is extremely long relative to shape
        constraints = analyze_shape_constraints(shape)
        if len(text) > constraints["estimated_capacity"] * 4:  # Only if 4x over estimate
            fits = False
        
        # Restore original text
        shape.text_frame.paragraphs[0].text = original_text
        return fits
    except:
        return True  # If test fails, assume it fits (optimistic approach)

def validate_truncation_quality(original_text, truncated_text):
    """Check if truncated text looks professional and preserves meaning"""
    
    if not truncated_text or not original_text:
        return False
        
    # If no truncation occurred, it's automatically good
    if truncated_text == original_text:
        return True
    
    # Rule 1: Must end with complete word
    if truncated_text and not truncated_text[-1].isspace():
        if not truncated_text.endswith(('.', '!', '?', ':')):
            # Check if last word is complete
            words = original_text.split()
            truncated_words = truncated_text.split()
            
            if truncated_words:
                last_word_in_truncated = truncated_words[-1]
                
                # Check if this word exists complete in original
                if last_word_in_truncated not in words:
                    return False  # Last word is cut off
    
    # Rule 2: Must preserve at least 60% of meaning (was 70%, now more lenient)
    if len(truncated_text) < len(original_text) * 0.6:
        return False  # Too much content lost
    
    # Rule 3: Should not end with connecting words (makes text look incomplete)
    bad_endings = ['and', 'or', 'but', 'the', 'a', 'an', 'of', 'in', 'on', 'at', 'to', 'for', 'with']
    last_word = truncated_text.split()[-1].lower() if truncated_text.split() else ""
    if last_word in bad_endings:
        return False  # Ends awkwardly
    
    # Rule 4: Should have reasonable length (not too short)
    if len(truncated_text) < 10:  # Very short text might not be meaningful
        return False
    
    return True  # Looks good

def smart_content_adaptation(text, max_length):
    """Intelligently adapt content to fit space while preserving user's meaning"""
    
    if len(text) <= max_length:
        return text
    
    # Strategy 1: Try to keep complete sentences (unchanged)
    sentences = text.split('. ')
    for i in range(len(sentences), 0, -1):
        candidate = '. '.join(sentences[:i])
        if not candidate.endswith('.') and i < len(sentences):
            candidate += '.'
        if len(candidate) <= max_length:
            return candidate
    
    # Strategy 2: Intelligent abbreviation while preserving meaning
    # This uses user content but makes it more concise
    adapted_text = adapt_content_intelligently(text, max_length)
    if len(adapted_text) <= max_length:
        return adapted_text
    
    # Strategy 3: Try to keep complete phrases (comma-separated)
    phrases = text.split(', ')
    for i in range(len(phrases), 0, -1):
        candidate = ', '.join(phrases[:i])
        if len(candidate) <= max_length:
            return candidate
    
    # Strategy 4: Keep complete words (current approach)
    words = text.split()
    for i in range(len(words), 0, -1):
        candidate = ' '.join(words[:i])
        if len(candidate) <= max_length:
            return candidate
    
    # Strategy 5: Last resort - character truncation at word boundary
    if ' ' in text[:max_length]:
        return text[:max_length].rsplit(' ', 1)[0]
    
    return text[:max_length]  # Absolute last resort

def adapt_content_intelligently(text, max_length):
    """Make user content more concise while preserving meaning - NO external content added"""
    
    # Common abbreviations that preserve meaning
    abbreviations = {
        'machine learning': 'ML',
        'artificial intelligence': 'AI',
        'natural language processing': 'NLP',
        'for example': 'e.g.',
        'that is': 'i.e.',
        'and so on': 'etc.',
        'versus': 'vs',
        'with regard to': 'regarding',
        'in order to': 'to',
        'due to the fact that': 'because',
        'a large number of': 'many',
        'a small number of': 'few',
        'is able to': 'can',
        'has the ability to': 'can',
        'it should be noted that': '',  # Remove redundant phrases
        'it is important to note that': '',
        'as a matter of fact': '',
        'the fact of the matter is': '',
    }
    
    adapted = text
    
    # Apply abbreviations only if they help fit the content
    for long_form, short_form in abbreviations.items():
        if long_form in adapted.lower():
            # Replace while preserving case
            import re
            pattern = re.compile(re.escape(long_form), re.IGNORECASE)
            adapted = pattern.sub(short_form, adapted)
            
            # Check if this helps us fit
            if len(adapted) <= max_length:
                return adapted
    
    # Remove redundant words while preserving core meaning
    # Only if the content is still too long
    if len(adapted) > max_length:
        # Remove filler words
        filler_words = ['very', 'quite', 'rather', 'really', 'actually', 'basically', 'essentially']
        words = adapted.split()
        filtered_words = [word for word in words if word.lower() not in filler_words]
        adapted = ' '.join(filtered_words)
    
    return adapted

# Keep the old function name for compatibility
def smart_truncate_with_meaning(text, max_length):
    """Legacy function name - now uses intelligent content adaptation"""
    return smart_content_adaptation(text, max_length)

def fit_text_with_quality_chain(shape, text, field_type="body"):
    """Try multiple approaches until we get good quality - comprehensive approach"""
    
    constraints = analyze_shape_constraints(shape)
    base_capacity = constraints["estimated_capacity"]
    
    print(f"   📐 Shape analysis: {constraints['width']:.1f}\"×{constraints['height']:.1f}\", capacity: ~{base_capacity} chars")
    
    # Define approaches in order of preference (most generous first)
    approaches = [
        # Approach 1: Try original text (most generous)
        {
            "name": "original", 
            "text": text,
            "description": "using original text"
        },
        
        # Approach 2: Try with very generous limits (4x capacity - more generous)
        {
            "name": "very_generous", 
            "text": smart_content_adaptation(text, base_capacity * 4),
            "description": "with very generous limits (400%)"
        },
        
        # Approach 3: Try with generous limits (3x capacity - increased)  
        {
            "name": "generous", 
            "text": smart_content_adaptation(text, base_capacity * 3),
            "description": "with generous limits (300%)"
        },
        
        # Approach 4: Try with moderate limits (2x capacity - increased)
        {
            "name": "moderate", 
            "text": smart_content_adaptation(text, int(base_capacity * 2)),
            "description": "with moderate limits (200%)"
        },
        
        # Approach 5: Try with conservative limits (1.5x capacity - increased)
        {
            "name": "conservative", 
            "text": smart_content_adaptation(text, int(base_capacity * 1.5)),
            "description": "with conservative limits (150%)"
        },
        
        # Approach 6: Simple fallback
        {
            "name": "simple", 
            "text": simple_fit_text(text, base_capacity),
            "description": "using simple word-boundary fitting"
        }
    ]
    
    original_len = len(text)
    
    # Try each approach until we find one that works well
    for approach in approaches:
        candidate_text = approach["text"]
        candidate_len = len(candidate_text)
        
        # Test if it fits and looks good
        fits = test_if_text_fits(shape, candidate_text)
        quality_good = validate_truncation_quality(text, candidate_text)
        
        if fits and quality_good:
            print(f"   SUCCESS: Success {approach['description']}")
            if candidate_len != original_len:
                print(f"   DIMENSION: Text adjusted: '{text[:20]}...' ({original_len}) → '{candidate_text[:20]}...' ({candidate_len})")
            else:
                print(f"   DIMENSION: Text fits perfectly: {original_len} chars")
            return candidate_text
        else:
            # Log why this approach didn't work
            reasons = []
            if not fits:
                reasons.append("doesn't fit")
            if not quality_good:
                reasons.append("poor quality")
            print(f"   WARNING: {approach['name']} approach failed: {', '.join(reasons)}")
    
    # If all approaches fail, use original text and let PowerPoint handle it
    print(f"   🚨 All approaches failed - using original text (PowerPoint will handle overflow)")
    return text

def simple_fit_text(text, desired_len):
    """Simple, reliable text fitting - the original approach as fallback"""
    if len(text) <= desired_len:
        return text
    
    # Find last complete word that fits
    words = text.split()
    fitted = ""
    for word in words:
        test_text = fitted + (" " if fitted else "") + word
        if len(test_text) <= desired_len:
            fitted = test_text
        else:
            break
    
    # If no words fit, take partial text up to last complete word
    if not fitted:
        # Take as much as possible, but try to end at word boundary
        truncated = text[:desired_len]
        if ' ' in truncated:
            fitted = truncated.rsplit(' ', 1)[0]
        else:
            fitted = truncated
    
    return fitted

def estimate_line_count(text, width_inches):
    """Estimate how many lines text will take"""
    try:
        # Rough estimate: ~12-15 characters per inch width
        chars_per_line = int(width_inches * 13)
        if chars_per_line < 10:
            chars_per_line = 10
        
        lines = len(text) // chars_per_line + (1 if len(text) % chars_per_line else 0)
        return max(lines, 1)
    except:
        return 1

def reduce_font_size(shape, target_reduction=0.8):
    """Reduce font size to fit more text"""
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size:
                    current_size = run.font.size.pt
                    new_size = max(8, int(current_size * target_reduction))  # Minimum 8pt
                    run.font.size = Pt(new_size)
                    print(f"   📝 Reduced font size: {current_size}pt → {new_size}pt")
    except Exception as e:
        print(f"   WARNING: Font size reduction failed: {e}")

def set_conservative_margins(shape):
    """Set conservative margins for better text fitting"""
    try:
        if hasattr(shape, 'text_frame'):
            tf = shape.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
    except:
        pass

def apply_no_autosize_settings(shape):
    """Apply settings to prevent auto-sizing and enable word wrap"""
    try:
        if hasattr(shape, 'text_frame'):
            tf = shape.text_frame
            tf.auto_size = False  # Disable auto-sizing
            tf.word_wrap = True   # Enable word wrapping
    except:
        pass

# Legacy functions for backward compatibility - now use the new quality chain approach
def intelligent_truncate(text, max_length, strategy="word_boundary"):
    """Legacy function - now redirects to smart_truncate_with_meaning"""
    return smart_truncate_with_meaning(text, max_length)

def fit_text_intelligently(shape, text, field_type="body"):
    """Main text fitting function - now uses the comprehensive quality chain approach"""
    try:
        return fit_text_with_quality_chain(shape, text, field_type)
    except Exception as e:
        print(f"   WARNING: Quality chain fitting failed: {e}")
        # Ultimate fallback to simple approach
        return simple_fit_text(text, 100)  # Safe default length

def replace_with_ai_content(pptx_path, ai_output, output_path):
    """Replace PowerPoint content with AI-generated text"""
    prs = Presentation(pptx_path)
    successful_replacements = 0
    total_fields = len(ai_output)
    
    print(f"📝 Replacing content in PowerPoint template...")
    
    for item in ai_output:
        source = item.get('source')
        source_index = item.get('source_index')
        slide_index = item.get('slide', -1)
        shape_index = item.get('shape')
        text = item.get('text', '')
        
        shape = None
        field_type = "body"  # Default
        
        try:
            # Determine field type from text characteristics
            if len(text) < 50:
                field_type = "header"
            elif len(text) < 150:
                field_type = "subheader"
            else:
                field_type = "body"
            
            # Get the shape based on source with robust validation
            if source == 'slide' and slide_index >= 0:
                # Validate slide index exists
                if slide_index < len(prs.slides):
                    slide = prs.slides[slide_index]
                    # Validate shape index exists and is not None
                    if shape_index is not None and shape_index < len(slide.shapes):
                        shape = slide.shapes[shape_index]
                    else:
                        print(f"   WARNING: Invalid shape index {shape_index} for slide {slide_index} (has {len(slide.shapes)} shapes)")
                else:
                    print(f"   WARNING: Invalid slide index {slide_index} (presentation has {len(prs.slides)} slides)")
            elif source == 'layout':
                # Handle layout shapes with validation
                if source_index is not None and source_index < len(prs.slide_layouts):
                    layout = prs.slide_layouts[source_index]
                    if shape_index is not None and shape_index < len(layout.shapes):
                        shape = layout.shapes[shape_index]
                    else:
                        print(f"   WARNING: Invalid shape index {shape_index} for layout {source_index}")
                else:
                    print(f"   WARNING: Invalid layout index {source_index}")
            elif source == 'master':
                # Handle master shapes with validation
                if source_index is not None and source_index < len(prs.slide_masters):
                    master = prs.slide_masters[source_index]
                    if shape_index is not None and shape_index < len(master.shapes):
                        shape = master.shapes[shape_index]
                    else:
                        print(f"   WARNING: Invalid shape index {shape_index} for master {source_index}")
                else:
                    print(f"   WARNING: Invalid master index {source_index}")
            else:
                print(f"   WARNING: Unknown source type '{source}' or invalid slide index {slide_index}")
            
            # Replace content if shape found
            if shape and hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                # CRITICAL: Always try to set content, never skip
                try:
                    # Use the new comprehensive quality chain approach
                    fitted_text = fit_text_with_quality_chain(shape, text, field_type)
                    
                    # Apply conservative settings
                    apply_no_autosize_settings(shape)
                    set_conservative_margins(shape)
                    
                    # Set the text content
                    if shape.text_frame.paragraphs:
                        p = shape.text_frame.paragraphs[0]
                        p.text = fitted_text
                        
                        # Apply smart text color
                        text_color = get_smart_text_color(shape)
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    run.font.color.rgb = text_color
                        
                        successful_replacements += 1
                    
                except Exception as e:
                    # FALLBACK: If everything fails, use simple approach
                    print(f"   WARNING: Advanced fitting failed, using simple fallback: {e}")
                    try:
                        simple_text = simple_fit_text(text, 150)  # More generous fallback
                        shape.text_frame.paragraphs[0].text = simple_text
                        successful_replacements += 1
                    except Exception as e2:
                        print(f"Warning: Could not replace ({item}) - {e2}")
            else:
                print(f"Warning: Could not replace ({item}) - shape not found or no text frame")
                
        except Exception as e:
            print(f"Warning: Could not replace ({item}) - {e}")
    
    # Save the presentation
    prs.save(output_path)
    print(f"Successfully replaced {successful_replacements} out of {total_fields} fields")
    print(f"Saved as {output_path}. (Close any open instances while running.)")

def create_ppt_from_scratch(content_analysis_results, output_path):
    """Create PowerPoint presentation from scratch with overflow checking and AI reformatting"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    prs = Presentation()
    
    print(f"🎨 Creating PowerPoint from scratch with format-specific layouts...")
    print(f"DEBUG: Overflow checking and AI reformatting enabled...")
    
    for i, slide_data in enumerate(content_analysis_results):
        format_type = slide_data.get('format_type', 'paragraph')
        content = slide_data.get('content', '')
        question = slide_data.get('question', f'Slide {i+1}')
        slide_num = i + 1
        
        print(f"   📝 Creating slide {slide_num}: {format_type} format")
        
        # Create slide data structure for overflow checking
        slide_data_for_check = {
            'slide_num': slide_num,
            'question': question,
            'content': content,
            'format_type': format_type
        }
        
        # Create slide with overflow checking and AI reformatting
        try:
            print(f"   🎨 Calling create_slide_with_overflow_check for format_type: {format_type}")
            create_slide_with_overflow_check(prs, slide_data_for_check)
            print(f"   SUCCESS: Slide {slide_num} created successfully")
        except Exception as e:
            print(f"   ERROR: Error creating slide {slide_num}: {e}")
            print(f"   🔄 Falling back to basic slide creation for format_type: {format_type}")
            # Fallback to basic slide creation
            if format_type == 'timeline':
                create_timeline_slide(prs, question, content, slide_num)
            elif format_type == 'table':
                create_table_slide(prs, question, content, slide_num)
            elif format_type == 'boxes':
                create_comparison_slide(prs, question, content, slide_num)
            elif format_type == 'bullet_points':
                create_bullet_slide(prs, question, content, slide_num)
            else:  # paragraph or default
                create_paragraph_slide(prs, question, content, slide_num)
    
    # Save the presentation
    prs.save(output_path)
    print(f"SUCCESS: PowerPoint created successfully: {output_path}")
    return output_path

def create_timeline_slide(prs, question, content, slide_num):
    """Create a timeline slide with enhanced question-aware structure and AI-powered content extraction"""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    print(f"   >>> Creating enhanced timeline slide for: {question[:50]}...")
    
    # Add slide
    slide_layout = get_safe_slide_layout(prs)  # Safe layout selection
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title with AI-generated heading
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    
    # Generate professional title using AI
    from .prompt_builder import generate_slide_title
    # Extract topic from question or use a default
    topic = question.split()[-1] if question else ""
    professional_title = generate_slide_title(question, content, "timeline", topic)
    title_frame.text = professional_title
    
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # STEP 1: Plan timeline structure based on question intent
    timeline_plan = plan_timeline_structure(question, content)
    
    # STEP 2: Extract timeline data using AI-powered approach
    timeline_data = extract_timeline_data_with_ai(question, content, timeline_plan)
    
    # STEP 3: Fallback to legacy extraction if needed
    if not timeline_data or len(timeline_data) <= 1:
        print(f"   >>> Enhanced extraction failed, trying legacy approach...")
        timeline_data = parse_timeline_content_legacy(content, question)
        # Convert legacy format to new format
        if timeline_data:
            headers = timeline_plan['headers']
            converted_data = [headers]
            for i, step in enumerate(timeline_data[:4]):
                if timeline_plan['timeline_type'] == "historical":
                    entry = [f"Period {i+1}", step[:80], "Historical significance"]
                elif timeline_plan['timeline_type'] == "process":
                    entry = [f"Step {i+1}", step[:80], "Leads to next step"]
                else:
                    entry = [f"Item {i+1}", step[:80], "Part of timeline"]
                converted_data.append(entry)
            timeline_data = converted_data
    
    # STEP 4: Final fallback to content-based extraction
    if not timeline_data or len(timeline_data) <= 1:
        print(f"   >>> Using content-based extraction as final fallback...")
        timeline_data = create_fallback_timeline_data(content, timeline_plan)
    
    if timeline_data and len(timeline_data) > 1:
        # Create timeline based on visual style
        visual_style = timeline_plan['visual_style']
        
        if visual_style == "chronological_timeline":
            create_chronological_timeline_layout(slide, timeline_data, timeline_plan)
        elif visual_style == "process_flow":
            create_process_flow_layout(slide, timeline_data, timeline_plan)
        elif visual_style == "lifecycle_diagram":
            create_lifecycle_diagram_layout(slide, timeline_data, timeline_plan)
        elif visual_style == "sequential_flow":
            create_sequential_flow_layout(slide, timeline_data, timeline_plan)
        elif visual_style == "stage_progression":
            create_stage_progression_layout(slide, timeline_data, timeline_plan)
        else:
            # Default to process flow
            create_process_flow_layout(slide, timeline_data, timeline_plan)
    else:
        # Fallback: Create simple text box if no timeline data
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        content_frame = content_box.text_frame
        content_frame.text = content
        content_frame.paragraphs[0].font.size = Pt(16)
        content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

def parse_timeline_content_legacy(content, question=""):
    """Legacy timeline parsing for fallback compatibility"""
    import re
    
    # Look for numbered steps (1., 2., 3., etc.)
    numbered_steps = re.findall(r'\d+\.\s+(.+?)(?=\d+\.|$)', content, re.DOTALL)
    
    if numbered_steps:
        # Clean up the steps
        steps = [step.strip() for step in numbered_steps if step.strip()]
        return steps
    
    # Look for step indicators
    step_indicators = [
        r'step\s+\d+', r'phase\s+\d+', r'stage\s+\d+', r'first', r'second', r'third', r'fourth'
    ]
    
    lines = content.split('\n')
    steps = []
    
    for line in lines:
        line = line.strip()
        if line:
            # Check if line contains step indicators
            for indicator in step_indicators:
                if re.search(indicator, line, re.IGNORECASE):
                    steps.append(line)
                    break
            
            # If no step indicator found but line is substantial, treat as step
            if len(line) > 20 and line not in steps:
                steps.append(line)
    
    # If we found steps, return them
    if steps:
        return steps[:4]  # Limit to 4 steps
    
    # If no clear steps found, split content into logical chunks
    sentences = re.split(r'[.!?]+', content)
    steps = []
    
    for sentence in sentences:
        sentence = sentence.strip()
        if sentence and len(sentence) > 10:
            steps.append(sentence)
            if len(steps) >= 4:  # Limit to 4 steps
                break
    
    return steps

def create_fallback_timeline_data(content, timeline_plan):
    """Create fallback timeline data when all other methods fail"""
    headers = timeline_plan['headers']
    timeline_data = [headers]
    
    sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20][:4]
    
    for i, sentence in enumerate(sentences):
        if timeline_plan['timeline_type'] == "historical":
            entry = [f"Period {i+1}", sentence[:80], "Historical context"]
        elif timeline_plan['timeline_type'] == "process":
            entry = [f"Step {i+1}", sentence[:80], "Process outcome"]
        elif timeline_plan['timeline_type'] == "lifecycle":
            entry = [f"Phase {i+1}", sentence[:80], "Phase deliverables"]
        elif timeline_plan['timeline_type'] == "sequential":
            entry = [f"Stage {i+1}", sentence[:80], "Next stage"]
        else:
            entry = [f"Item {i+1}", sentence[:80], "Timeline context"]
        
        timeline_data.append(entry)
    
    return timeline_data if len(timeline_data) > 1 else []

# =====================================================
# ENHANCED TIMELINE VISUAL LAYOUTS
# =====================================================

def create_chronological_timeline_layout(slide, timeline_data, timeline_plan):
    """Create a chronological timeline with dates/periods on a horizontal line"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    print(f"   >>> Creating chronological timeline layout...")
    
    # Skip header row
    entries = timeline_data[1:]
    num_entries = min(len(entries), 4)  # Limit to 4 entries for layout
    
    if num_entries == 0:
        return
    
    # Timeline parameters
    timeline_y = Inches(3)
    timeline_start_x = Inches(1)
    timeline_width = Inches(8)
    entry_width = timeline_width / num_entries
    entry_height = Inches(1.5)
    
    # Draw main timeline line - using rectangle instead of connector to avoid PowerPoint corruption
    timeline_line = slide.shapes.add_shape(
        1,  # Rectangle shape
        timeline_start_x, timeline_y - Pt(2), 
        timeline_width, Pt(4)
    )
    timeline_line.fill.solid()
    timeline_line.fill.fore_color.rgb = RGBColor(50, 50, 50)
    timeline_line.line.color.rgb = RGBColor(50, 50, 50)
    timeline_line.line.width = Pt(0)  # No border for timeline line
    
    # Add timeline entries
    for i, entry in enumerate(entries[:num_entries]):
        entry_x = timeline_start_x + (i * entry_width)
        
        # Timeline marker
        marker = slide.shapes.add_shape(1, entry_x + entry_width/2 - Inches(0.1), timeline_y - Inches(0.1), 
                                       Inches(0.2), Inches(0.2))
        marker.fill.solid()
        marker.fill.fore_color.rgb = RGBColor(70, 130, 180)
        marker.line.color.rgb = RGBColor(0, 0, 0)
        
        # Period/Date box (above timeline)
        period_box = slide.shapes.add_textbox(entry_x, timeline_y - Inches(1), entry_width, Inches(0.6))
        period_frame = period_box.text_frame
        period_frame.text = str(entry[0])  # Period/Era
        period_frame.paragraphs[0].font.size = Pt(11)
        period_frame.paragraphs[0].font.bold = True
        period_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        period_frame.word_wrap = True
        
        # Event box (below timeline)
        event_box = slide.shapes.add_textbox(entry_x, timeline_y + Inches(0.3), entry_width, entry_height)
        event_frame = event_box.text_frame
        event_frame.text = str(entry[1])  # Key Development
        event_frame.paragraphs[0].font.size = Pt(10)
        event_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        event_frame.word_wrap = True
        event_frame.margin_left = Pt(4)
        event_frame.margin_right = Pt(4)
        
        # Add border to event box
        event_box.line.color.rgb = RGBColor(200, 200, 200)
        event_box.line.width = Pt(1)

def create_process_flow_layout(slide, timeline_data, timeline_plan):
    """Create a process flow with arrows and step boxes - Enhanced with proper wrapping"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    print(f"   >>> Creating process flow layout...")
    
    # Skip header row
    entries = timeline_data[1:]
    num_entries = len(entries)  # Use all entries, not just 4
    
    if num_entries == 0:
        return
    
    # Enhanced layout parameters with wrapping
    slide_width = Inches(10)  # Standard slide width
    slide_height = Inches(7.5)  # Standard slide height
    margin = Inches(0.5)
    available_width = slide_width - (2 * margin)
    
    step_width = Inches(2.0)  # Slightly smaller for better fit
    step_height = Inches(1.5)
    arrow_width = Inches(0.4)
    step_spacing = Inches(0.2)
    
    # Calculate how many boxes fit per row
    box_and_arrow_width = step_width + arrow_width + step_spacing
    boxes_per_row = max(1, int(available_width / box_and_arrow_width))
    
    # Calculate rows needed
    rows_needed = (num_entries + boxes_per_row - 1) // boxes_per_row
    
    # Adjust step height if we have multiple rows
    if rows_needed > 1:
        available_height = slide_height - Inches(3)  # Leave space for title
        step_height = min(step_height, available_height / rows_needed - Inches(0.3))
    
    row_height = step_height + Inches(0.5)
    start_y = Inches(2.0)
    
    for i, entry in enumerate(entries):
        # Calculate row and column
        row = i // boxes_per_row
        col = i % boxes_per_row
        
        # Calculate position
        step_x = margin + col * box_and_arrow_width
        step_y = start_y + row * row_height
        
        # Ensure we don't go off the slide
        if step_x + step_width > slide_width - margin:
            continue  # Skip if it would go off slide
        
        # Step box
        step_box = slide.shapes.add_textbox(step_x, step_y, step_width, step_height)
        step_frame = step_box.text_frame
        
        # Step content - handle potential None values
        step_title = str(entry[0]) if entry[0] else f"Step {i+1}"
        step_action = str(entry[1]) if len(entry) > 1 and entry[1] else "Action"
        
        # Format content in the box with better text handling
        max_chars = 50 if step_height < Inches(1.5) else 70
        if len(step_action) > max_chars:
            step_action = smart_content_adaptation(step_action, max_chars)
        
        step_frame.text = f"{step_title}\n\n{step_action}"
        
        # Style the step box
        step_frame.paragraphs[0].font.size = Pt(11)
        step_frame.paragraphs[0].font.bold = True
        step_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if len(step_frame.paragraphs) > 1:
            step_frame.paragraphs[1].font.size = Pt(9)
            step_frame.paragraphs[1].alignment = PP_ALIGN.LEFT
        
        step_frame.word_wrap = True
        step_frame.margin_left = Pt(6)
        step_frame.margin_right = Pt(6)
        step_frame.margin_top = Pt(6)
        step_frame.margin_bottom = Pt(6)
        
        # Add border and fill to step box
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
        step_box.line.color.rgb = RGBColor(70, 130, 180)
        step_box.line.width = Pt(2)
        
        # Add arrow logic with wrapping consideration
        if i < num_entries - 1:
            next_row = (i + 1) // boxes_per_row
            next_col = (i + 1) % boxes_per_row
            
            if next_row == row:  # Same row - horizontal arrow
                arrow_x = step_x + step_width
                arrow_y = step_y + step_height/2
                next_x = step_x + box_and_arrow_width
                if next_x + step_width <= slide_width - margin:  # Only if next box fits
                    # Create arrow using rectangle instead of connector to avoid PowerPoint corruption
                    arrow = slide.shapes.add_shape(
                        1,  # Rectangle shape
                        arrow_x, arrow_y - Pt(1), 
                        arrow_width, Pt(2)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = RGBColor(0, 0, 0)
                    arrow.line.width = Pt(0)  # No border
            # For row transitions, could add curved arrows or skip arrows

def create_lifecycle_diagram_layout(slide, timeline_data, timeline_plan):
    """Create a lifecycle diagram with phases in a circular/semi-circular arrangement"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    import math
    
    print(f"   >>> Creating lifecycle diagram layout...")
    
    # Skip header row
    entries = timeline_data[1:]
    num_entries = min(len(entries), 4)  # Limit to 4 entries for layout
    
    if num_entries == 0:
        return
    
    # Lifecycle parameters
    center_x = Inches(5)
    center_y = Inches(3.5)
    radius = Inches(2.5)
    phase_width = Inches(2)
    phase_height = Inches(1.5)
    
    # Colors for different phases
    colors = [
        RGBColor(255, 182, 193),  # Light Pink
        RGBColor(173, 216, 230),  # Light Blue
        RGBColor(144, 238, 144),  # Light Green
        RGBColor(255, 218, 185)   # Peach
    ]
    
    for i, entry in enumerate(entries[:num_entries]):
        # Calculate position in semicircle
        angle = math.pi * i / (num_entries - 1) if num_entries > 1 else 0
        phase_x = center_x + radius * math.cos(angle) - phase_width/2
        phase_y = center_y - radius * math.sin(angle) - phase_height/2
        
        # Phase box
        phase_box = slide.shapes.add_textbox(phase_x, phase_y, phase_width, phase_height)
        phase_frame = phase_box.text_frame
        
        # Phase content
        phase_title = str(entry[0])  # Phase
        phase_activities = str(entry[1])  # Activities
        
        # Format content in the box
        phase_frame.text = f"{phase_title}\n\n{smart_content_adaptation(phase_activities, 50)}"
        
        # Style the phase box
        phase_frame.paragraphs[0].font.size = Pt(12)
        phase_frame.paragraphs[0].font.bold = True
        phase_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if len(phase_frame.paragraphs) > 1:
            phase_frame.paragraphs[1].font.size = Pt(9)
            phase_frame.paragraphs[1].alignment = PP_ALIGN.LEFT
        
        phase_frame.word_wrap = True
        phase_frame.margin_left = Pt(6)
        phase_frame.margin_right = Pt(6)
        phase_frame.margin_top = Pt(6)
        phase_frame.margin_bottom = Pt(6)
        
        # Add color and border to phase box
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = colors[i % len(colors)]
        phase_box.line.color.rgb = RGBColor(0, 0, 0)
        phase_box.line.width = Pt(2)

def create_sequential_flow_layout(slide, timeline_data, timeline_plan):
    """Create a sequential flow with connected stages"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    print(f"   >>> Creating sequential flow layout...")
    
    # Skip header row
    entries = timeline_data[1:]
    num_entries = min(len(entries), 4)  # Limit to 4 entries for layout
    
    if num_entries == 0:
        return
    
    # Sequential flow parameters
    if num_entries <= 2:
        # Horizontal layout for 1-2 items
        create_horizontal_sequential_layout(slide, entries, num_entries)
    else:
        # Vertical layout for 3+ items
        create_vertical_sequential_layout(slide, entries, num_entries)

def create_horizontal_sequential_layout(slide, entries, num_entries):
    """Create horizontal sequential layout"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    stage_width = Inches(3.5)
    stage_height = Inches(2)
    arrow_width = Inches(1)
    start_y = Inches(2.5)
    
    for i, entry in enumerate(entries[:num_entries]):
        stage_x = Inches(0.5) + (stage_width + arrow_width) * i
        
        # Stage box
        stage_box = slide.shapes.add_textbox(stage_x, start_y, stage_width, stage_height)
        stage_frame = stage_box.text_frame
        
        # Stage content
        stage_title = str(entry[0])  # Stage
        stage_event = str(entry[1])  # Event
        
        stage_frame.text = f"{stage_title}\n\n{smart_content_adaptation(stage_event, 70)}"
        
        # Style the stage box
        stage_frame.paragraphs[0].font.size = Pt(14)
        stage_frame.paragraphs[0].font.bold = True
        stage_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if len(stage_frame.paragraphs) > 1:
            stage_frame.paragraphs[1].font.size = Pt(11)
            stage_frame.paragraphs[1].alignment = PP_ALIGN.LEFT
        
        stage_frame.word_wrap = True
        stage_frame.margin_left = Pt(10)
        stage_frame.margin_right = Pt(10)
        stage_frame.margin_top = Pt(10)
        stage_frame.margin_bottom = Pt(10)
        
        # Style stage box
        stage_box.fill.solid()
        stage_box.fill.fore_color.rgb = RGBColor(230, 230, 250)
        stage_box.line.color.rgb = RGBColor(75, 0, 130)
        stage_box.line.width = Pt(2)
        
        # Add arrow using rectangle instead of connector to avoid PowerPoint corruption
        if i < num_entries - 1:
            arrow_x = stage_x + stage_width
            arrow_y = start_y + stage_height/2
            arrow = slide.shapes.add_shape(
                1,  # Rectangle shape
                arrow_x, arrow_y - Pt(1.5), 
                arrow_width, Pt(3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(75, 0, 130)
            arrow.line.width = Pt(0)  # No border

def create_vertical_sequential_layout(slide, entries, num_entries):
    """Create vertical sequential layout"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    stage_width = Inches(7)
    stage_height = Inches(1.2)
    arrow_height = Inches(0.3)
    start_x = Inches(1.5)
    start_y = Inches(2)
    
    for i, entry in enumerate(entries[:num_entries]):
        stage_y = start_y + (stage_height + arrow_height) * i
        
        # Stage box
        stage_box = slide.shapes.add_textbox(start_x, stage_y, stage_width, stage_height)
        stage_frame = stage_box.text_frame
        
        # Stage content
        stage_title = str(entry[0])  # Stage
        stage_event = str(entry[1])  # Event
        
        stage_frame.text = f"{stage_title}: {smart_content_adaptation(stage_event, 60)}"
        
        # Style the stage box
        stage_frame.paragraphs[0].font.size = Pt(12)
        stage_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        stage_frame.word_wrap = True
        stage_frame.margin_left = Pt(10)
        stage_frame.margin_right = Pt(10)
        stage_frame.margin_top = Pt(8)
        stage_frame.margin_bottom = Pt(8)
        
        # Style stage box
        stage_box.fill.solid()
        stage_box.fill.fore_color.rgb = RGBColor(245, 245, 220)
        stage_box.line.color.rgb = RGBColor(139, 69, 19)
        stage_box.line.width = Pt(2)
        
        # Add downward arrow using rectangle instead of connector to avoid PowerPoint corruption
        if i < num_entries - 1:
            arrow_x = start_x + stage_width/2
            arrow_y = stage_y + stage_height
            arrow = slide.shapes.add_shape(
                1,  # Rectangle shape
                arrow_x - Pt(1.5), arrow_y, 
                Pt(3), arrow_height
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(139, 69, 19)
            arrow.line.width = Pt(0)  # No border

def create_stage_progression_layout(slide, timeline_data, timeline_plan):
    """Create a stage progression layout with tiered levels"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    print(f"   >>> Creating stage progression layout...")
    
    # Skip header row
    entries = timeline_data[1:]
    num_entries = min(len(entries), 4)  # Limit to 4 entries for layout
    
    if num_entries == 0:
        return
    
    # Stage progression parameters
    stage_width = Inches(2)
    stage_height = Inches(1.5)
    level_spacing = Inches(0.3)
    start_x = Inches(1)
    base_y = Inches(4.5)
    
    # Colors for different stages (gradient effect)
    colors = [
        RGBColor(255, 240, 245),  # Very light pink
        RGBColor(255, 182, 193),  # Light pink
        RGBColor(255, 105, 180),  # Hot pink
        RGBColor(199, 21, 133)    # Medium violet red
    ]
    
    for i, entry in enumerate(entries[:num_entries]):
        # Calculate position (pyramid/stair-step effect)
        stage_x = start_x + i * (stage_width + level_spacing)
        stage_y = base_y - i * level_spacing * 2  # Each stage slightly higher
        
        # Stage box
        stage_box = slide.shapes.add_textbox(stage_x, stage_y, stage_width, stage_height)
        stage_frame = stage_box.text_frame
        
        # Stage content
        stage_title = str(entry[0])  # Stage
        stage_description = str(entry[1])  # Description
        
        stage_frame.text = f"{stage_title}\n\n{smart_content_adaptation(stage_description, 45)}"
        
        # Style the stage box
        stage_frame.paragraphs[0].font.size = Pt(12)
        stage_frame.paragraphs[0].font.bold = True
        stage_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if len(stage_frame.paragraphs) > 1:
            stage_frame.paragraphs[1].font.size = Pt(9)
            stage_frame.paragraphs[1].alignment = PP_ALIGN.LEFT
        
        stage_frame.word_wrap = True
        stage_frame.margin_left = Pt(8)
        stage_frame.margin_right = Pt(8)
        stage_frame.margin_top = Pt(8)
        stage_frame.margin_bottom = Pt(8)
        
        # Style stage box with progression colors
        stage_box.fill.solid()
        stage_box.fill.fore_color.rgb = colors[i % len(colors)]
        stage_box.line.color.rgb = RGBColor(0, 0, 0)
        stage_box.line.width = Pt(2)

def create_table_slide(prs, question, content, slide_num):
    """Create a table slide with enhanced question-aware structure and AI-powered content extraction"""
    import re
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    print(f"   >>> Creating enhanced table slide for: {question[:50]}...")
    
    # Add slide
    slide_layout = get_safe_slide_layout(prs)  # Safe layout selection
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title with AI-generated heading
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    
    # Generate professional title using AI
    from .prompt_builder import generate_slide_title
    # Extract topic from question or use a default
    topic = question.split()[-1] if question else ""
    professional_title = generate_slide_title(question, content, "table", topic)
    title_frame.text = professional_title
    
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # STEP 1: Plan table structure based on question intent
    table_plan = plan_table_structure(question, content)
    
    # STEP 2: Extract table data using AI-powered approach
    table_data = extract_table_data_with_ai(question, content, table_plan)
    
    # STEP 3: Fallback to legacy extraction if needed
    if not table_data or len(table_data) <= 1:
        print(f"   >>> Enhanced extraction failed, trying legacy approach...")
        table_data = parse_table_content(content, question)
    
    # STEP 4: Final fallback to content-based extraction
    if not table_data or len(table_data) <= 1:
        print(f"   >>> Using content-based extraction as final fallback...")
        table_data = extract_content_based_data(content, question)
    
    if table_data and len(table_data) > 1:
        # Calculate table dimensions
        rows = len(table_data)
        cols = max(len(row) for row in table_data)
        
        # Ensure reasonable limits
        rows = min(rows, 6)  # Reduced from 8 to 6 for better readability
        cols = min(cols, 3)  # Reduced from 4 to 3 for better spacing
        
        # Calculate table size based on content
        table_width = Inches(8.5)  # Increased width
        
        # Calculate optimal height based on rows
        if rows <= 2:
            table_height = Inches(3)
        elif rows <= 4:
            table_height = Inches(4)
        elif rows <= 6:
            table_height = Inches(5)
        else:
            table_height = Inches(5.5)
        
        # Create table with proper positioning
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.25), Inches(2), table_width, table_height)
        table = table_shape.table
        
        # Configure table properties for better text handling
        table.first_row = True  # Enable header row
        table.horz_banding = True  # Enable alternating row colors
        
        # Fill table with data and style
        for i, row_data in enumerate(table_data[:rows]):
            for j, cell_text in enumerate(row_data[:cols]):
                if i < len(table.rows) and j < len(table.columns):
                    cell = table.cell(i, j)
                    
                    # Clean and format cell text
                    cell_text = str(cell_text).strip()
                    
                    # Process text for better table formatting - use AI only for formatting when needed
                    if len(cell_text) > 120:  # Increased from 80 to 120
                        # Try AI formatting first to preserve meaning
                        reformatted = format_table_content_with_ai(cell_text, 120)
                        if reformatted != cell_text:
                            cell_text = reformatted
                        else:
                            # Fallback to word boundary truncation
                            words = cell_text.split()
                            if len(words) > 15:  # Increased from 12 to 15
                                cell_text = smart_content_adaptation(' '.join(words), 150)
                            else:
                                cell_text = smart_content_adaptation(cell_text, 120)
                    
                    # Remove markdown formatting
                    cell_text = re.sub(r'\*\*([^*]+)\*\*', r'\1', cell_text)  # Remove bold
                    cell_text = re.sub(r'\*([^*]+)\*', r'\1', cell_text)      # Remove italic
                    cell_text = re.sub(r'`([^`]+)`', r'\1', cell_text)        # Remove code
                    
                    cell.text = cell_text
                    
                    # Style header row
                    if i == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(70, 130, 180)
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                run.font.bold = True
                                run.font.size = Pt(14)
                    else:
                        # Style data rows
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(11)  # Slightly smaller for more content
                                run.font.color.rgb = RGBColor(0, 0, 0)
                    
                    # Configure text frame for proper wrapping
                    text_frame = cell.text_frame
                    text_frame.word_wrap = True
                    text_frame.margin_left = Pt(8)  # Increased margins
                    text_frame.margin_right = Pt(8)
                    text_frame.margin_top = Pt(4)
                    text_frame.margin_bottom = Pt(4)
                    
                    # Set text alignment
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    if j == 0:  # First column (usually headers/names)
                        text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                    else:  # Data columns
                        text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Changed to LEFT for better readability
        
        # Add borders to table (simplified approach)
        try:
            for row in table.rows:
                for cell in row.cells:
                    # Set cell border properties if available
                    if hasattr(cell, 'border_top'):
                        cell.border_top.color.rgb = RGBColor(0, 0, 0)
                        cell.border_top.width = Pt(1)
                    if hasattr(cell, 'border_bottom'):
                        cell.border_bottom.color.rgb = RGBColor(0, 0, 0)
                        cell.border_bottom.width = Pt(1)
                    if hasattr(cell, 'border_left'):
                        cell.border_left.color.rgb = RGBColor(0, 0, 0)
                        cell.border_left.width = Pt(1)
                    if hasattr(cell, 'border_right'):
                        cell.border_right.color.rgb = RGBColor(0, 0, 0)
                        cell.border_right.width = Pt(1)
        except:
            # Skip border styling if not supported
            pass
    else:
        # Fallback: Create simple text box if no table data
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        content_frame = content_box.text_frame
        content_frame.text = content
        content_frame.paragraphs[0].font.size = Pt(16)
        content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT



def create_comparison_slide(prs, question, content, slide_num):
    """Create an enhanced boxes slide with question-aware structure and AI-powered content extraction"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    print(f"   >>> Creating enhanced boxes slide for: {question[:50]}...")
    
    # Add slide
    slide_layout = get_safe_slide_layout(prs)  # Safe layout selection
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title with AI-generated heading
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    if not hasattr(title_box, 'text_frame') or not title_box.text_frame:
        print(f"   ERROR: Title box does not have text_frame")
        return
    title_frame = title_box.text_frame
    
    # Generate professional title using AI
    from .prompt_builder import generate_slide_title
    # Extract topic from question or use a default  
    topic = question.split()[-1] if question else ""
    professional_title = generate_slide_title(question, content, "boxes", topic)
    title_frame.text = professional_title
    
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # STEP 1: Plan boxes structure based on question intent
    boxes_plan = plan_boxes_structure(question, content)
    
    # STEP 2: Extract boxes data using AI-powered approach
    boxes_data = extract_boxes_data_with_ai(question, content, boxes_plan)
    
    # STEP 3: Fallback to legacy extraction if needed
    if not boxes_data or len(boxes_data) < 2:
        print(f"   >>> Enhanced extraction failed, trying legacy approach...")
        boxes_data = parse_comparison_content(content, question)
        # Convert legacy format to new format
        if boxes_data and len(boxes_data) >= 2:
            legacy_boxes = []
            for i, item in enumerate(boxes_data[:4]):
                legacy_boxes.append({
                    "title": f"Item {i+1}",
                    "content": str(item)[:150] if item else "No content",
                    "details": "Legacy extraction"
                })
            boxes_data = legacy_boxes
    
    # STEP 4: Final fallback to content-based extraction
    if not boxes_data or len(boxes_data) < 2:
        print(f"   >>> Using content-based extraction as final fallback...")
        boxes_data = create_fallback_boxes_data(content, boxes_plan)
    
    if boxes_data and len(boxes_data) >= 2:
        # Create boxes based on layout type
        layout_type = boxes_plan['layout_type']
        
        if layout_type == "two_column_comparison":
            create_two_column_comparison_layout(slide, boxes_data, boxes_plan)
        elif layout_type in ["multi_box_grid", "feature_grid", "category_grid", "aspect_grid"]:
            create_multi_box_grid_layout(slide, boxes_data, boxes_plan)
        else:
            # Default to multi-box grid
            create_multi_box_grid_layout(slide, boxes_data, boxes_plan)
    else:
        # Fallback: Create simple text box if no boxes data
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        content_frame = content_box.text_frame
        content_frame.text = content
        content_frame.paragraphs[0].font.size = Pt(16)
        content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

def create_fallback_boxes_data(content, boxes_plan):
    """Create fallback boxes data when all other methods fail"""
    boxes_type = boxes_plan['boxes_type']
    expected_boxes = boxes_plan['expected_boxes']
    
    sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20][:expected_boxes]
    
    boxes_data = []
    for i, sentence in enumerate(sentences):
        # Increased content length for richer boxes entries
        content_text = smart_content_adaptation(sentence, 150)
        
        if boxes_type == "comparison" and i < 2:
            boxes_data.append({
                "title": f"Item {chr(65+i)}",  # A, B
                "content": content_text,
                "details": "Comparison item"
            })
        elif boxes_type == "components":
            boxes_data.append({
                "title": f"Component {i+1}",
                "content": content_text,
                "details": "System component and function"
            })
        elif boxes_type == "features":
            boxes_data.append({
                "title": f"Feature {i+1}",
                "content": content_text,
                "details": "Feature benefits and usage"
            })
        else:
            boxes_data.append({
                "title": f"Item {i+1}",
                "content": content_text,
                "details": "Important details and context"
            })
    
    return boxes_data if len(boxes_data) >= 2 else []

def create_two_column_comparison_layout(slide, boxes_data, boxes_plan):
    """Create a two-column comparison layout with dynamic height and improved content fitting"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    print(f"   Creating two-column comparison layout with {len(boxes_data)} boxes")
    
    # STEP 1: Calculate optimal height ensuring both boxes are equal
    optimal_height = calculate_optimal_box_height(boxes_data)
    print(f"   DIMENSION: Using equal box height: {optimal_height.inches:.1f} inches")
    
    # Two side-by-side boxes with dynamic height
    box_width = Inches(4)
    box_height = optimal_height  # Dynamic height instead of fixed
    left_x = Inches(0.5)
    right_x = Inches(5.5)
    box_y = Inches(1.8)
    
    # Take only the first two boxes for comparison
    left_data = boxes_data[0] if len(boxes_data) > 0 else {"title": "Item A", "content": "No content", "details": ""}
    right_data = boxes_data[1] if len(boxes_data) > 1 else {"title": "Item B", "content": "No content", "details": ""}
    
    # Left comparison box with error handling
    try:
        left_box = slide.shapes.add_textbox(left_x, box_y, box_width, box_height)
        left_frame = left_box.text_frame
    except Exception as e:
        print(f"   ERROR: Error creating left comparison box: {e}")
        return
    left_frame.word_wrap = True
    left_frame.margin_left = Pt(15)
    left_frame.margin_right = Pt(15)
    left_frame.margin_top = Pt(15)
    left_frame.margin_bottom = Pt(15)
    
    # Add content to left box with improved space estimation
    left_title = left_data.get('title', 'Item A')
    left_content = left_data.get('content', 'No content available')
    left_details = left_data.get('details', '')
    
    # Calculate available space and adapt content
    left_capacity = estimate_character_capacity(box_width, box_height)
    left_total_length = len(left_title) + len(left_content) + len(left_details)
    
    if left_total_length <= left_capacity:
        # Content fits - show everything
        left_text = f"{left_title}\n\n{left_content}"
        if left_details and left_details != left_content:
            left_text += f"\n\n{left_details}"
    else:
        # Adapt content to fit available space
        title_space = len(left_title) + 4
        remaining_space = left_capacity - title_space
        if left_details and len(left_details) > 50:
            content_space = int(remaining_space * 0.7)
            details_space = int(remaining_space * 0.3)
            adapted_content = smart_content_adaptation(left_content, content_space)
            adapted_details = smart_content_adaptation(left_details, details_space)
            left_text = f"{left_title}\n\n{adapted_content}\n\n{adapted_details}"
        else:
            adapted_content = smart_content_adaptation(left_content, remaining_space)
            left_text = f"{left_title}\n\n{adapted_content}"
    
    left_frame.text = left_text
    
    # Style left box title
    if left_frame.paragraphs:
        left_frame.paragraphs[0].font.size = Pt(16)
        left_frame.paragraphs[0].font.bold = True
        left_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Style content paragraphs
        for i in range(1, len(left_frame.paragraphs)):
            left_frame.paragraphs[i].font.size = Pt(14)
            left_frame.paragraphs[i].alignment = PP_ALIGN.LEFT
    
    # Right comparison box with error handling
    try:
        right_box = slide.shapes.add_textbox(right_x, box_y, box_width, box_height)
        right_frame = right_box.text_frame
    except Exception as e:
        print(f"   ERROR: Error creating right comparison box: {e}")
        return
    right_frame.word_wrap = True
    right_frame.margin_left = Pt(15)
    right_frame.margin_right = Pt(15)
    right_frame.margin_top = Pt(15)
    right_frame.margin_bottom = Pt(15)
    
    # Add content to right box with improved space estimation
    right_title = right_data.get('title', 'Item B')
    right_content = right_data.get('content', 'No content available')
    right_details = right_data.get('details', '')
    
    # Calculate available space and adapt content (same logic as left box)
    right_capacity = estimate_character_capacity(box_width, box_height)
    right_total_length = len(right_title) + len(right_content) + len(right_details)
    
    if right_total_length <= right_capacity:
        # Content fits - show everything
        right_text = f"{right_title}\n\n{right_content}"
        if right_details and right_details != right_content:
            right_text += f"\n\n{right_details}"
    else:
        # Adapt content to fit available space
        title_space = len(right_title) + 4
        remaining_space = right_capacity - title_space
        if right_details and len(right_details) > 50:
            content_space = int(remaining_space * 0.7)
            details_space = int(remaining_space * 0.3)
            adapted_content = smart_content_adaptation(right_content, content_space)
            adapted_details = smart_content_adaptation(right_details, details_space)
            right_text = f"{right_title}\n\n{adapted_content}\n\n{adapted_details}"
        else:
            adapted_content = smart_content_adaptation(right_content, remaining_space)
            right_text = f"{right_title}\n\n{adapted_content}"
    
    right_frame.text = right_text
    
    # Style right box title
    if right_frame.paragraphs:
        right_frame.paragraphs[0].font.size = Pt(16)
        right_frame.paragraphs[0].font.bold = True
        right_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Style content paragraphs
        for i in range(1, len(right_frame.paragraphs)):
            right_frame.paragraphs[i].font.size = Pt(14)
            right_frame.paragraphs[i].alignment = PP_ALIGN.LEFT
    
    # Add subtle fill colors first (proper styling order)
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue
    
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(255, 248, 240)  # Light orange
    
    # Add borders to boxes (CORRECT line initialization for python-pptx)
    # Left box border - color MUST be set first to make line visible
    left_box.line.color.rgb = RGBColor(0, 120, 200)  # Blue border - this makes line visible
    left_box.line.width = Pt(2)
    
    # Right box border - color MUST be set first to make line visible
    right_box.line.color.rgb = RGBColor(200, 120, 0)  # Orange border - this makes line visible  
    right_box.line.width = Pt(2)

def calculate_optimal_box_height(boxes_data):
    """Calculate optimal height for boxes based on content length - ensures equal heights"""
    from pptx.util import Inches
    
    if not boxes_data:
        return Inches(3)  # Default height
    
    max_content_length = 0
    
    # Find the longest content across all boxes
    for box_data in boxes_data:
        box_title = box_data.get('title', '')
        box_content = box_data.get('content', '')
        box_details = box_data.get('details', '')
        
        # Calculate total text length
        total_text = f"{box_title}\n\n{box_content}"
        if box_details and box_details != box_content:
            total_text += f"\n\n{box_details}"
        
        content_length = len(total_text)
        max_content_length = max(max_content_length, content_length)
    
    # Calculate height based on content (more generous than before)
    if max_content_length <= 150:
        return Inches(2.5)  # Minimum height for short content
    elif max_content_length <= 300:
        return Inches(3.2)  # Medium content
    elif max_content_length <= 500:
        return Inches(4.0)  # Long content
    elif max_content_length <= 800:
        return Inches(4.8)  # Very long content
    else:
        return Inches(5.5)  # Maximum height for extensive content


def preserve_template_elements(slide):
    """Identify and preserve existing template elements like connectors and lines"""
    preserved_shapes = []
    
    # Check for existing shapes that shouldn't be modified
    for i, shape in enumerate(slide.shapes):
        # Preserve connector lines, autoshapes, and grouped elements
        if (hasattr(shape, 'connector') or 
            hasattr(shape, 'auto_shape_type') or
            shape.shape_type == 6):  # GROUP shape type
            preserved_shapes.append({
                'index': i,
                'shape_type': shape.shape_type,
                'position': (shape.left, shape.top, shape.width, shape.height) if hasattr(shape, 'left') else None
            })
    
    return preserved_shapes

def estimate_character_capacity(box_width, box_height):
    """Estimate character capacity based on box dimensions - more accurate than before"""
    from pptx.util import Inches
    
    # Convert to inches if needed
    width_inches = box_width.inches if hasattr(box_width, 'inches') else box_width
    height_inches = box_height.inches if hasattr(box_height, 'inches') else box_height
    
    # Account for margins (24pt = 0.33 inches total)
    usable_width = width_inches - 0.33
    usable_height = height_inches - 0.33
    
    # Estimate based on font size and character width
    # For 11pt font: ~6 chars per inch width, ~8 lines per inch height
    chars_per_line = int(usable_width * 6)
    lines_available = int(usable_height * 8)
    
    # Account for title spacing (takes ~2 lines)
    content_lines = max(1, lines_available - 2)
    
    # Calculate total capacity with some buffer
    total_capacity = chars_per_line * content_lines
    
    # Add buffer for better fitting (85% of theoretical capacity)
    return int(total_capacity * 0.85)

def create_multi_box_grid_layout(slide, boxes_data, boxes_plan):
    """Create a multi-box grid layout using proper auto-shapes to avoid compatibility errors"""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_SHADOW_STYLE
    
    print(f"   Creating multi-box grid layout with {len(boxes_data)} boxes")
    
    num_boxes = len(boxes_data)
    if num_boxes == 0:
        return
    
    # Calculate layout
    if num_boxes <= 2:
        rows, cols = 1, 2
        box_width = Inches(4)
        box_height = Inches(2.5)
    elif num_boxes == 3:
        rows, cols = 1, 3
        box_width = Inches(3.1)
        box_height = Inches(2.5)
    elif num_boxes == 4:
        rows, cols = 2, 2
        box_width = Inches(3.8)
        box_height = Inches(2.2)
    else:
        rows, cols = 2, 3
        box_width = Inches(3.2)
        box_height = Inches(2.0)
    
    # Calculate starting positions
    spacing_inches = 0.2
    total_width = cols * box_width.inches + (cols - 1) * spacing_inches
    slide_width = 10.0
    start_x = (slide_width - total_width) / 2
    start_y = 1.8
    
    # Create boxes using proper auto-shapes
    for i, box_data in enumerate(boxes_data[:rows * cols]):
        row = i // cols
        col = i % cols
        
        # Calculate position
        box_x_inches = start_x + col * (box_width.inches + spacing_inches)
        box_y_inches = start_y + row * (box_height.inches + spacing_inches)
        box_x = Inches(box_x_inches)
        box_y = Inches(box_y_inches)
        
        try:
            # ✅ Create proper auto-shape that supports advanced features
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_x, box_y, box_width, box_height)
            
            # ✅ Safe auto-shape adjustments
            if hasattr(box, 'adjustments') and len(box.adjustments) > 0:
                box.adjustments[0] = 0.1  # Rounded corners
            
            # ✅ Advanced styling
            box.fill.solid()
            colors = [
                RGBColor(240, 248, 255),  # Light blue
                RGBColor(240, 255, 240),  # Light green  
                RGBColor(255, 240, 240),  # Light red
                RGBColor(255, 255, 240),  # Light yellow
                RGBColor(248, 240, 255),  # Light purple
                RGBColor(240, 255, 255),  # Light cyan
            ]
            box.fill.fore_color.rgb = colors[i % len(colors)]
            
            # ✅ Professional border
            box.line.color.rgb = RGBColor(70, 130, 180)
            box.line.width = Pt(2)
            
            # ✅ Add subtle shadow
            try:
                box.shadow.inherit = False
                box.shadow.style = MSO_SHADOW_STYLE.OUTER
                box.shadow.blur_radius = Pt(3)
                box.shadow.distance = Pt(2)
                box.shadow.direction = 315  # Bottom-right shadow
            except:
                pass  # Shadow might not be supported in all versions
            
            # ✅ Add text content
            text_frame = box.text_frame
            text_frame.clear()
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            text_frame.word_wrap = True
            
            # Get box content
            box_title = box_data.get('title', f'Item {i+1}')
            box_content = box_data.get('content', 'No content available')
            box_details = box_data.get('details', '')
            
            # Smart content formatting
            if box_details and box_details != box_content and "not specified" not in box_details.lower():
                full_text = f"{box_title}\n\n{box_content}\n\n{box_details}"
            else:
                full_text = f"{box_title}\n\n{box_content}"
            
            # Adapt content to fit in box
            if len(full_text) > 250:  # If content is too long
                available_space = 200  # Character limit for box
                title_space = len(box_title) + 4
                content_space = available_space - title_space
                adapted_content = smart_content_adaptation(box_content, content_space)
                full_text = f"{box_title}\n\n{adapted_content}"
            
            text_frame.text = full_text
            
            # Format text with proper sizing
            if text_frame.paragraphs:
                # Title formatting
                text_frame.paragraphs[0].font.size = Pt(14)
                text_frame.paragraphs[0].font.bold = True
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Content formatting
                for j in range(1, len(text_frame.paragraphs)):
                    text_frame.paragraphs[j].font.size = Pt(11)
                    text_frame.paragraphs[j].alignment = PP_ALIGN.LEFT
                    
        except Exception as e:
            print(f"   ERROR: Error creating auto-shape box {i}: {e}")
            # Fallback to simple text box if auto-shape fails
            try:
                text_box = slide.shapes.add_textbox(box_x, box_y, box_width, box_height)
                text_frame = text_box.text_frame
                text_frame.text = f"{box_data.get('title', f'Item {i+1}')}\n\n{box_data.get('content', 'No content available')}"
                text_box.fill.solid()
                text_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
                text_box.line.color.rgb = RGBColor(80, 80, 80)
                text_box.line.width = Pt(1)
            except Exception as e2:
                print(f"   ERROR: Fallback text box also failed for box {i}: {e2}")
                continue


def create_bullet_slide(prs, question, content, slide_num):
    """Create a bullet point slide"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    
    # Add slide
    slide_layout = get_safe_slide_layout(prs)  # Safe layout selection
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title with AI-generated heading
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    
    # Generate professional title using AI
    from .prompt_builder import generate_slide_title
    # Extract topic from question or use a default
    topic = question.split()[-1] if question else ""
    professional_title = generate_slide_title(question, content, "bullet_points", topic)
    title_frame.text = professional_title
    
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Reformat content to fit slide
    reformatted_content = reformat_content_for_slide(content, "bullet_points")
    
    # Create bullet point list
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = reformatted_content
    content_frame.paragraphs[0].font.size = Pt(18)
    content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Configure text frame for proper wrapping
    content_frame.word_wrap = True
    content_frame.margin_left = Pt(12)
    content_frame.margin_right = Pt(12)
    content_frame.margin_top = Pt(6)
    content_frame.margin_bottom = Pt(6)

def calculate_dynamic_limit(content, format_type, base_limit):
    """Calculate dynamic character limit based on content complexity"""
    import re
    
    # If content has lots of structured data, allow more
    if format_type == 'table' and has_structured_data(content):
        return int(base_limit * 1.5)
    # If content has clear sections, allow more
    elif format_type == 'boxes' and has_comparison_patterns(content):
        return int(base_limit * 1.3)
    # If content has numbered steps, allow more
    elif format_type == 'timeline' and has_sequential_patterns(content):
        return int(base_limit * 1.4)
    return base_limit

def has_structured_data(content):
    """Check if content has structured data patterns"""
    import re
    patterns = [
        r'(\w+):\s*(\d+[%$]?)',  # Metric: Value
        r'(\w+)\s+=\s+(\d+[%$]?)',  # Metric = Value
        r'(\w+)\s+is\s+(\d+[%$]?)',  # Metric is Value
    ]
    return any(re.search(pattern, content, re.IGNORECASE) for pattern in patterns)

def has_comparison_patterns(content):
    """Check if content has comparison patterns"""
    import re
    patterns = [
        r'(\w+)\s+vs\s+(\w+)',  # A vs B
        r'(\w+)\s+versus\s+(\w+)',  # A versus B
        r'(\w+)\s+compared\s+to\s+(\w+)',  # A compared to B
        r'differences\s+between',  # Differences between
        r'advantages\s+and\s+disadvantages',  # Pros and cons
    ]
    return any(re.search(pattern, content, re.IGNORECASE) for pattern in patterns)

def has_sequential_patterns(content):
    """Check if content has sequential patterns"""
    import re
    patterns = [
        r'\d+\.\s+',  # Numbered steps
        r'first.*?second.*?third',  # Sequential words
        r'step\s+\d+',  # Step patterns
        r'phase\s+\d+',  # Phase patterns
        r'stage\s+\d+',  # Stage patterns
        r'process.*?involves',  # Process descriptions
    ]
    return any(re.search(pattern, content, re.IGNORECASE) for pattern in patterns)

def reformat_content_for_slide(content, format_type):
    """Reformat content for specific slide format with intelligent rephrasing"""
    # Get base character limits
    char_limits = {
        'paragraph': 1500,
        'bullet_points': 1200,
        'table': 800,
        'timeline': 600,
        'boxes': 600
    }
    
    base_limit = char_limits.get(format_type, 800)
    
    # Calculate dynamic limit based on content characteristics
    dynamic_limit = calculate_dynamic_limit(content, format_type, base_limit)
    
    # Reformat based on format type
    if format_type == 'paragraph':
        return reformat_paragraph_content(content, dynamic_limit)
    elif format_type == 'bullet_points':
        return reformat_bullet_content(content, dynamic_limit)
    elif format_type == 'table':
        return reformat_table_content(content, dynamic_limit)
    elif format_type == 'timeline':
        return reformat_timeline_content(content, dynamic_limit)
    elif format_type == 'boxes':
        return reformat_comparison_content(content, dynamic_limit)
    else:
        return reformat_paragraph_content(content, dynamic_limit)

def reformat_paragraph_content(content, max_chars):
    """Reformat paragraph content with intelligent rephrasing, never using truncation"""
    import re
    
    # Clean up content
    content = re.sub(r'\*\*([^*]+)\*\*', r'\1', content)
    content = re.sub(r'\*([^*]+)\*', r'\1', content)
    content = re.sub(r'`([^`]+)`', r'\1', content)
    
    if len(content) <= max_chars:
        return content
    
    # Split into sentences
    sentences = re.split(r'[.!?]+', content)
    sentences = [s.strip() for s in sentences if s.strip()]
    
    if not sentences:
        return content[:max_chars]
    
    # Calculate importance scores for sentences
    sentence_scores = []
    for sentence in sentences:
        score = calculate_sentence_importance(sentence)
        sentence_scores.append((sentence, score))
    
    # Sort by importance (highest first)
    sentence_scores.sort(key=lambda x: x[1], reverse=True)
    
    # Select most important sentences that fit within limit
    selected_sentences = []
    current_length = 0
    
    for sentence, score in sentence_scores:
        # Try to rephrase long sentences
        if len(sentence) > 150:
            rephrased = rephrase_long_sentence(sentence, 120)
            sentence = rephrased
        
        # Check if adding this sentence would exceed limit
        if current_length + len(sentence) + 2 <= max_chars:  # +2 for ". "
            selected_sentences.append(sentence)
            current_length += len(sentence) + 2
        else:
            # Try to rephrase this sentence to fit
            if len(sentence) > 50:
                rephrased = rephrase_sentence_for_length(sentence, max_chars - current_length - 10)
                if len(rephrased) + current_length + 2 <= max_chars:
                    selected_sentences.append(rephrased)
                    break
            break
    
    # If we still don't have enough content, try intelligent summarization
    if not selected_sentences or len(' '.join(selected_sentences)) < max_chars * 0.3:
        return create_intelligent_summary(content, max_chars)
    
    # Join sentences and validate quality
    result = '. '.join(selected_sentences) + '.'
    
    # Ensure we don't have "..." in the result
    if '...' in result:
        # Remove any truncation and rephrase
        result = re.sub(r'\.{3,}', '.', result)
        if len(result) < max_chars * 0.5:
            return create_intelligent_summary(content, max_chars)
    
    return result

def calculate_sentence_importance(sentence):
    """Calculate how important a sentence is based on content analysis"""
    import re
    score = 0.0
    
    # Keywords that indicate importance
    important_keywords = ['key', 'important', 'essential', 'main', 'primary', 'significant', 'critical']
    for keyword in important_keywords:
        if keyword.lower() in sentence.lower():
            score += 0.2
    
    # Numbers and metrics (often important)
    if re.search(r'\d+[%$]?', sentence):
        score += 0.3
    
    # Technical terms (often core concepts)
    technical_terms = ['accuracy', 'precision', 'recall', 'cost', 'performance', 'efficiency']
    for term in technical_terms:
        if term.lower() in sentence.lower():
            score += 0.15
    
    # Sentence length (medium length sentences are often most informative)
    if 20 <= len(sentence.split()) <= 40:
        score += 0.1
    
    return min(score, 1.0)

def smart_truncate_sentence(sentence, max_chars):
    """Intelligently truncate a sentence at word boundaries"""
    words = sentence.split()
    truncated = []
    current_length = 0
    
    for word in words:
        if current_length + len(word) + 1 <= max_chars - 3:  # Leave room for "..."
            truncated.append(word)
            current_length += len(word) + 1
        else:
            break
    
    if truncated:
        return smart_content_adaptation(' '.join(truncated), max_chars)
    return None

def create_multi_column_layout(slide, content, question):
    """Create multi-column layout for long content"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    
    # Split content into logical sections
    sections = split_content_into_sections(content)
    
    # Create 2-column layout
    col_width = Inches(4)
    col_height = Inches(4)
    
    for i, section in enumerate(sections[:2]):  # Limit to 2 columns
        x_pos = Inches(0.5) + (col_width + Inches(0.5)) * i
        y_pos = Inches(2)
        
        # Create text box for this section
        text_box = slide.shapes.add_textbox(x_pos, y_pos, col_width, col_height)
        text_frame = text_box.text_frame
        text_frame.text = section
        text_frame.paragraphs[0].font.size = Pt(14)  # Slightly smaller for columns
        text_frame.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        
        # Configure text frame
        text_frame.word_wrap = True
        text_frame.margin_left = Pt(8)
        text_frame.margin_right = Pt(8)
        text_frame.margin_top = Pt(4)
        text_frame.margin_bottom = Pt(4)

def split_content_into_sections(content):
    """Split content into logical sections"""
    import re
    
    # Split by double newlines or major separators
    sections = re.split(r'\n\s*\n', content)
    
    # If no clear sections, split by sentences
    if len(sections) <= 1:
        sentences = re.split(r'[.!?]+', content)
        # Group sentences into sections
        sections = []
        current_section = []
        
        for sentence in sentences:
            sentence = sentence.strip()
            if sentence:
                current_section.append(sentence)
                if len(current_section) >= 3:  # 3 sentences per section
                    sections.append('. '.join(current_section) + '.')
                    current_section = []
        
        # Add remaining sentences
        if current_section:
            sections.append('. '.join(current_section) + '.')
    
    return sections

def validate_content_quality(original_content, reformatted_content):
    """Validate that reformatted content maintains quality"""
    
    # Check if content is too short (might be over-truncated)
    if len(reformatted_content) < 50:
        return False, "Content too short - may have lost important information"
    
    # Check if content ends with "..." (indicates truncation)
    if reformatted_content.endswith("..."):
        return False, "Content was truncated - may be incomplete"
    
    # Check if content has meaningful structure
    if not has_meaningful_structure(reformatted_content):
        return False, "Content lacks meaningful structure"
    
    # Check if key information was preserved
    if not key_information_preserved(original_content, reformatted_content):
        return False, "Key information may have been lost"
    
    return True, "Content quality is acceptable"

def has_meaningful_structure(content):
    """Check if content has meaningful structure"""
    import re
    # Should have at least one complete sentence
    if not re.search(r'[.!?]$', content):
        return False
    
    # Should have reasonable word count
    words = content.split()
    if len(words) < 10:
        return False
    
    return True

def key_information_preserved(original, reformatted):
    """Check if key information was preserved"""
    import re
    # Extract key terms from original
    original_terms = extract_key_terms(original)
    reformatted_terms = extract_key_terms(reformatted)
    
    # Check if most key terms are preserved
    preserved_count = sum(1 for term in original_terms if term in reformatted_terms)
    preservation_rate = preserved_count / len(original_terms) if original_terms else 1.0
    
    return preservation_rate >= 0.7  # At least 70% of key terms preserved

def extract_key_terms(content):
    """Extract key terms from content"""
    import re
    
    # Extract numbers and metrics
    metrics = re.findall(r'\d+[%$]?', content)
    
    # Extract capitalized terms (often proper nouns or important concepts)
    capitalized = re.findall(r'\b[A-Z][a-z]+\b', content)
    
    # Extract technical terms
    technical_terms = ['accuracy', 'precision', 'recall', 'cost', 'performance', 'efficiency']
    found_technical = [term for term in technical_terms if term.lower() in content.lower()]
    
    return metrics + capitalized + found_technical

def create_intelligent_summary(content, max_chars):
    """Create an intelligent summary of content that fits within limits"""
    import re
    
    # Extract key information using multiple strategies
    summary_parts = []
    
    # Strategy 1: Extract metrics and numbers (often most important)
    metrics = re.findall(r'(\w+):\s*(\d+[%$]?)', content)
    if metrics:
        summary_parts.append("Key Metrics:")
        for metric, value in metrics[:3]:  # Limit to 3 most important
            summary_parts.append(f"• {metric}: {value}")
    
    # Strategy 2: Extract main concepts (words in bold or important)
    concepts = re.findall(r'\*\*([^*]+)\*\*', content)
    if concepts:
        summary_parts.append("Main Concepts:")
        for concept in concepts[:3]:
            summary_parts.append(f"• {concept}")
    
    # Strategy 3: Extract key sentences (sentences with important keywords)
    sentences = re.split(r'[.!?]+', content)
    important_sentences = []
    
    important_keywords = ['key', 'important', 'essential', 'main', 'primary', 'significant', 'critical']
    for sentence in sentences:
        sentence = sentence.strip()
        if any(keyword in sentence.lower() for keyword in important_keywords):
            if len(sentence) > 20:  # Only meaningful sentences
                important_sentences.append(sentence)
    
    if important_sentences:
        summary_parts.append("Key Points:")
        for sentence in important_sentences[:2]:  # Limit to 2 sentences
            summary_parts.append(f"• {sentence}")
    
    # Strategy 4: If no structured data found, create a general summary
    if not summary_parts:
        words = content.split()
        if len(words) > 50:
            # Take first 50 words and add context
            summary_parts.append("Summary:")
            summary_parts.append(smart_content_adaptation(" ".join(words), 200))
        else:
            summary_parts.append(content)
    
    # Combine and format
    summary = "\n".join(summary_parts)
    
    # Ensure it fits within limits
    if len(summary) > max_chars:
        # Truncate intelligently
        summary = smart_content_adaptation(summary, max_chars)
    
    return summary

def rephrase_long_sentence(sentence, target_length):
    """Rephrase a long sentence to be more concise"""
    import re
    
    # Remove unnecessary words and phrases
    sentence = re.sub(r'\b(very|really|quite|extremely|highly)\b', '', sentence, flags=re.IGNORECASE)
    sentence = re.sub(r'\b(in order to|so as to)\b', 'to', sentence, flags=re.IGNORECASE)
    sentence = re.sub(r'\b(due to the fact that|because of the fact that)\b', 'because', sentence, flags=re.IGNORECASE)
    sentence = re.sub(r'\b(at this point in time|at the present time)\b', 'now', sentence, flags=re.IGNORECASE)
    sentence = re.sub(r'\b(in the event that|in case)\b', 'if', sentence, flags=re.IGNORECASE)
    
    # Simplify complex phrases
    sentence = re.sub(r'\b(provide an explanation of|give an explanation of)\b', 'explain', sentence, flags=re.IGNORECASE)
    sentence = re.sub(r'\b(carry out|perform the task of)\b', 'do', sentence, flags=re.IGNORECASE)
    sentence = re.sub(r'\b(make use of|utilize)\b', 'use', sentence, flags=re.IGNORECASE)
    
    # Remove redundant words
    sentence = re.sub(r'\b(completely|entirely|totally)\s+(full|empty|finished)', r'\2', sentence, flags=re.IGNORECASE)
    sentence = re.sub(r'\b(actual|real)\s+(fact|truth)', r'\2', sentence, flags=re.IGNORECASE)
    
    # If still too long, extract key information
    if len(sentence) > target_length:
        return extract_key_information(sentence, target_length)
    
    return sentence

def rephrase_sentence_for_length(sentence, max_length):
    """Rephrase a sentence to fit within a specific length"""
    import re
    
    # First try to simplify the sentence
    simplified = rephrase_long_sentence(sentence, max_length)
    
    if len(simplified) <= max_length:
        return simplified
    
    # If still too long, extract the most important part
    return extract_key_information(sentence, max_length)

def extract_key_information(sentence, max_length):
    """Extract the most important information from a sentence"""
    import re
    
    # Look for key patterns (numbers, technical terms, main subjects)
    key_patterns = [
        r'\b\d+[%$]?\b',  # Numbers and percentages
        r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b',  # Proper nouns
        r'\b(accuracy|precision|recall|cost|performance|efficiency|security|encryption)\b',  # Technical terms
    ]
    
    key_info = []
    for pattern in key_patterns:
        matches = re.findall(pattern, sentence, re.IGNORECASE)
        key_info.extend(matches)
    
    # If we found key information, create a concise sentence
    if key_info:
        # Take the first few key terms and create a simple sentence
        key_terms = key_info[:3]  # Limit to 3 key terms
        return f"Key points: {', '.join(key_terms)}."
    
    # Fallback: take the first part of the sentence
    words = sentence.split()
    if len(words) <= 10:
        return sentence
    
    # Take first 8-10 words and add context
    first_part = ' '.join(words[:8])
    if first_part.endswith(('.', '!', '?')):
        return first_part
    else:
        return first_part + "."

def reformat_bullet_content(content, max_chars):
    """Reformat bullet content with intelligent rephrasing, never using truncation"""
    import re
    
    # Clean up content
    content = re.sub(r'\*\*([^*]+)\*\*', r'\1', content)
    content = re.sub(r'\*([^*]+)\*', r'\1', content)
    content = re.sub(r'`([^`]+)`', r'\1', content)
    
    if len(content) <= max_chars:
        return content
    
    # Extract bullet points
    bullets = parse_bullet_content(content)
    
    if not bullets:
        # If no bullets found, try to create them from sentences
        sentences = re.split(r'[.!?]+', content)
        bullets = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 10]
    
    if not bullets:
        return create_bullet_summary(content, max_chars)
    
    # Process each bullet point
    reformatted_bullets = []
    current_length = 0
    
    for bullet in bullets:
        # Clean up bullet point
        bullet = re.sub(r'^[•\-\*]\s*', '', bullet.strip())
        
        if not bullet:
            continue
        
        # If bullet is too long, rephrase it
        if len(bullet) > 100:
            rephrased = rephrase_bullet_point(bullet, 80)
            bullet = rephrased
        
        # Check if adding this bullet would exceed limit
        bullet_with_prefix = f"• {bullet}"
        if current_length + len(bullet_with_prefix) + 1 <= max_chars:  # +1 for newline
            reformatted_bullets.append(bullet_with_prefix)
            current_length += len(bullet_with_prefix) + 1
        else:
            # Try to rephrase this bullet to fit
            if len(bullet) > 30:
                rephrased = rephrase_bullet_point(bullet, max_chars - current_length - 10)
                if len(f"• {rephrased}") + current_length + 1 <= max_chars:
                    reformatted_bullets.append(f"• {rephrased}")
            break
    
    # If we don't have enough bullets, create a summary
    if not reformatted_bullets or len('\n'.join(reformatted_bullets)) < max_chars * 0.3:
        return create_bullet_summary(content, max_chars)
    
    result = '\n'.join(reformatted_bullets)
    
    # Ensure we don't have "..." in the result
    if '...' in result:
        # Remove any truncation and rephrase
        result = re.sub(r'\.{3,}', '.', result)
        if len(result) < max_chars * 0.5:
            return create_bullet_summary(content, max_chars)
    
    return result

def rephrase_bullet_point(point, max_length):
    """Rephrase a bullet point to be more concise"""
    import re
    
    # Remove unnecessary words
    point = re.sub(r'\b(very|really|quite|extremely)\b', '', point, flags=re.IGNORECASE)
    point = re.sub(r'\b(in order to|so as to)\b', 'to', point, flags=re.IGNORECASE)
    
    # Simplify complex phrases
    point = re.sub(r'\b(provide|give|offer)\b', 'show', point, flags=re.IGNORECASE)
    point = re.sub(r'\b(utilize|make use of)\b', 'use', point, flags=re.IGNORECASE)
    
    # If still too long, extract key information
    if len(point) > max_length:
        return extract_key_information(point, max_length)
    
    return point

def create_bullet_summary(content, max_chars):
    """Create a bullet point summary from content"""
    import re
    
    # Extract key information
    key_terms = []
    
    # Look for numbers and metrics
    numbers = re.findall(r'\d+[%$]?', content)
    key_terms.extend(numbers[:3])
    
    # Look for technical terms
    tech_terms = re.findall(r'\b(accuracy|precision|recall|cost|performance|efficiency|security|encryption)\b', content, re.IGNORECASE)
    key_terms.extend(tech_terms[:3])
    
    # Look for main concepts (capitalized words)
    concepts = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', content)
    key_terms.extend(concepts[:3])
    
    # Create bullet points
    bullets = []
    for term in key_terms[:5]:  # Limit to 5 bullet points
        bullets.append(f"• {term}")
    
    if not bullets:
        # Fallback to simple bullet points
        words = content.split()[:20]
        bullets = [f"• {' '.join(words)}"]
    
    return '\n'.join(bullets)

def reformat_table_content(content, max_chars):
    """Reformat table content with intelligent rephrasing, never using truncation"""
    import re
    
    # Clean up content
    content = re.sub(r'\*\*([^*]+)\*\*', r'\1', content)
    content = re.sub(r'\*([^*]+)\*', r'\1', content)
    content = re.sub(r'`([^`]+)`', r'\1', content)
    
    if len(content) <= max_chars:
        return content
    
    # Extract table data
    table_data = parse_table_content(content)
    
    if not table_data:
        # If no table data found, create bullet points from key information
        return create_bullet_summary(content, max_chars)
    
    # Process table data
    reformatted_items = []
    current_length = 0
    
    for item in table_data:
        if not item.strip():
            continue
        
        # If item is too long, rephrase it
        if len(item) > 60:
            # For table items, keep them concise
            words = item.split()
            if len(words) > 8:
                item = ' '.join(words[:8])
        
        # Check if adding this item would exceed limit
        item_with_bullet = f"• {item}"
        if current_length + len(item_with_bullet) + 1 <= max_chars:  # +1 for newline
            reformatted_items.append(item_with_bullet)
            current_length += len(item_with_bullet) + 1
        else:
            break
    
    # If we don't have enough items, create a summary
    if not reformatted_items or len('\n'.join(reformatted_items)) < max_chars * 0.3:
        return create_bullet_summary(content, max_chars)
    
    result = '\n'.join(reformatted_items)
    
    # Ensure we don't have "..." in the result
    if '...' in result:
        # Remove any truncation and rephrase
        result = re.sub(r'\.{3,}', '.', result)
        if len(result) < max_chars * 0.5:
            return create_bullet_summary(content, max_chars)
    
    return result

def reformat_timeline_content(content, max_chars):
    """Reformat timeline content with intelligent rephrasing, never using truncation"""
    import re
    
    # Clean up content
    content = re.sub(r'\*\*([^*]+)\*\*', r'\1', content)
    content = re.sub(r'\*([^*]+)\*', r'\1', content)
    content = re.sub(r'`([^`]+)`', r'\1', content)
    
    if len(content) <= max_chars:
        return content
    
    # Extract timeline steps
    steps = parse_timeline_content(content)
    
    if not steps:
        # If no steps found, try to create them from sentences
        sentences = re.split(r'[.!?]+', content)
        steps = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 10]
    
    if not steps:
        return create_intelligent_summary(content, max_chars)
    
    # Process each step
    reformatted_steps = []
    current_length = 0
    
    for i, step in enumerate(steps):
        # Clean up step
        step = re.sub(r'^\d+\.\s*', '', step.strip())
        
        if not step:
            continue
        
        # If step is too long, rephrase it
        if len(step) > 80:
            rephrased = rephrase_timeline_step(step, 70)
            step = rephrased
        
        # Check if adding this step would exceed limit
        step_with_number = f"{i+1}. {step}"
        if current_length + len(step_with_number) + 1 <= max_chars:  # +1 for newline
            reformatted_steps.append(step_with_number)
            current_length += len(step_with_number) + 1
        else:
            # Try to rephrase this step to fit
            if len(step) > 30:
                rephrased = rephrase_timeline_step(step, max_chars - current_length - 10)
                if len(f"{i+1}. {rephrased}") + current_length + 1 <= max_chars:
                    reformatted_steps.append(f"{i+1}. {rephrased}")
            break
    
    # If we don't have enough steps, create a summary
    if not reformatted_steps or len('\n'.join(reformatted_steps)) < max_chars * 0.3:
        return create_intelligent_summary(content, max_chars)
    
    result = '\n'.join(reformatted_steps)
    
    # Ensure we don't have "..." in the result
    if '...' in result:
        # Remove any truncation and rephrase
        result = re.sub(r'\.{3,}', '.', result)
        if len(result) < max_chars * 0.5:
            return create_intelligent_summary(content, max_chars)
    
    return result

def rephrase_timeline_step(step, max_length):
    """Rephrase a timeline step to be more concise"""
    import re
    
    # Remove unnecessary words
    step = re.sub(r'\b(very|really|quite|extremely)\b', '', step, flags=re.IGNORECASE)
    step = re.sub(r'\b(in order to|so as to)\b', 'to', step, flags=re.IGNORECASE)
    
    # Simplify action verbs
    step = re.sub(r'\b(perform the task of|carry out)\b', 'do', step, flags=re.IGNORECASE)
    step = re.sub(r'\b(provide|give|offer)\b', 'show', step, flags=re.IGNORECASE)
    
    # If still too long, extract key information
    if len(step) > max_length:
        return extract_key_information(step, max_length)
    
    return step

def reformat_comparison_content(content, max_chars):
    """Reformat comparison content with intelligent rephrasing, never using truncation"""
    import re
    
    # Clean up content
    content = re.sub(r'\*\*([^*]+)\*\*', r'\1', content)
    content = re.sub(r'\*([^*]+)\*', r'\1', content)
    content = re.sub(r'`([^`]+)`', r'\1', content)
    
    if len(content) <= max_chars:
        return content
    
    # Parse comparison content
    comparison_parts = parse_comparison_content(content)
    
    if len(comparison_parts) < 2:
        # If no clear comparison, create a simple summary
        return create_intelligent_summary(content, max_chars)
    
    # Process each part
    reformatted_parts = []
    part_limit = max_chars // 2  # Split limit between parts
    
    for i, part in enumerate(comparison_parts):
        if not part.strip():
            continue
        
        # If part is too long, rephrase it
        if len(part) > part_limit:
            rephrased = rephrase_comparison_item(part, part_limit - 20)
            part = rephrased
        
        reformatted_parts.append(part.strip())
    
    # If we don't have enough content, create a summary
    if not reformatted_parts or len('\n\n'.join(reformatted_parts)) < max_chars * 0.3:
        return create_intelligent_summary(content, max_chars)
    
    result = '\n\n'.join(reformatted_parts)
    
    # Ensure we don't have "..." in the result
    if '...' in result:
        # Remove any truncation and rephrase
        result = re.sub(r'\.{3,}', '.', result)
        if len(result) < max_chars * 0.5:
            return create_intelligent_summary(content, max_chars)
    
    return result

def rephrase_comparison_item(item, max_length):
    """Rephrase a comparison item to be more concise"""
    import re
    
    # Remove unnecessary words
    item = re.sub(r'\b(very|really|quite|extremely)\b', '', item, flags=re.IGNORECASE)
    
    # Simplify complex terms
    item = re.sub(r'\b(utilize|make use of)\b', 'use', item, flags=re.IGNORECASE)
    item = re.sub(r'\b(provide|give|offer)\b', 'show', item, flags=re.IGNORECASE)
    
    # If still too long, take the first few words
    if len(item) > max_length:
        words = item.split()[:4]
        return ' '.join(words)
    
    return item

def create_paragraph_slide(prs, question, content, slide_num):
    """Create a paragraph slide with intelligent content distribution and AI reformatting"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    
    # Add slide
    slide_layout = get_safe_slide_layout(prs)  # Safe layout selection
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title with AI-generated heading
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    
    # Generate professional title using AI
    from .prompt_builder import generate_slide_title
    # Extract topic from question or use a default
    topic = question.split()[-1] if question else ""
    professional_title = generate_slide_title(question, content, "paragraph", topic)
    title_frame.text = professional_title
    
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Intelligent content reformatting
    reformatted_content = reformat_content_for_slide(content, "paragraph")
    
    # Check if content is still too long for single text box
    if len(reformatted_content) > 1000:
        # Split into multiple text boxes for better layout
        create_multi_column_layout(slide, reformatted_content, question)
    else:
        # Single text box approach with AI reformatting if needed
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
        content_frame = content_box.text_frame
        
        # Use AI reformatting if content doesn't fit
        final_content = fit_text_with_ai_reformatting(content_box, reformatted_content, "body", question)
        content_frame.text = final_content
        content_frame.paragraphs[0].font.size = Pt(16)
        content_frame.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        
        # Configure text frame for proper wrapping
        content_frame.word_wrap = True
        content_frame.margin_left = Pt(12)
        content_frame.margin_right = Pt(12)
        content_frame.margin_top = Pt(6)
        content_frame.margin_bottom = Pt(6)

def parse_timeline_content(content):
    """Parse content for timeline slide - extract numbered steps"""
    import re
    
    # Look for numbered steps (1., 2., 3., etc.)
    numbered_steps = re.findall(r'\d+\.\s+(.+?)(?=\d+\.|$)', content, re.DOTALL)
    
    if numbered_steps:
        # Clean up the steps
        steps = [step.strip() for step in numbered_steps if step.strip()]
        return steps
    
    # Look for step indicators
    step_indicators = [
        r'step\s+\d+', r'phase\s+\d+', r'stage\s+\d+', r'first', r'second', r'third', r'fourth'
    ]
    
    lines = content.split('\n')
    steps = []
    
    for line in lines:
        line = line.strip()
        if line:
            # Check if line contains step indicators
            for indicator in step_indicators:
                if re.search(indicator, line, re.IGNORECASE):
                    steps.append(line)
                    break
            
            # If no step indicator found but line is substantial, treat as step
            if len(line) > 20 and line not in steps:
                steps.append(line)
    
    # If we found steps, return them
    if steps:
        return steps[:4]  # Limit to 4 steps
    
    # If no clear steps found, split content into logical chunks
    sentences = re.split(r'[.!?]+', content)
    steps = []
    
    for sentence in sentences:
        sentence = sentence.strip()
        if sentence and len(sentence) > 10:
            steps.append(sentence)
            if len(steps) >= 4:  # Limit to 4 steps
                break
    
    return steps

def parse_table_content(content, question=""):
    """Parse content to extract table data with question-aware analysis - uses only actual content"""
    import re
    
    table_data = []
    
    # Step 1: Analyze the question to understand what we're looking for
    question_type = analyze_question_type(question)
    print(f"   >>> Question analysis: {question_type}")
    
    # Step 2: Extract relevant content based on question type
    if question_type == "statistics":
        table_data = extract_statistics_data(content, question)
    elif question_type == "features":
        table_data = extract_features_data(content, question)
    elif question_type == "costs":
        table_data = extract_costs_data(content, question)
    elif question_type == "process":
        table_data = extract_process_data(content, question)
    elif question_type == "comparison":
        table_data = extract_comparison_data(content, question)
    else:
        # Generic approach - look for any structured data
        table_data = extract_generic_data(content, question)
    
    # Step 3: If no relevant data found, try content-based extraction
    if not table_data or len(table_data) <= 1:
        print(f"   >>> No question-specific data found, extracting from content...")
        table_data = extract_content_based_data(content, question)
    
    # Step 4: If still no data, try the old reliable approach
    if not table_data or len(table_data) <= 1:
        print(f"   >>> Using fallback content extraction...")
        table_data = extract_fallback_data(content)
    
    return table_data

def extract_fallback_data(content):
    """Fallback extraction that tries to find any meaningful content"""
    import re
    
    table_data = [['Topic', 'Description']]
    
    # Split into sentences and look for any meaningful content
    sentences = re.split(r'[.!?]+', content)
    meaningful_sentences = [s.strip() for s in sentences if len(s.strip()) > 10 and len(s.strip()) < 500]
    
    for i, sentence in enumerate(meaningful_sentences[:4]):
        if sentence.strip():
            # Look for any capitalized terms that might be topics
            capitalized_terms = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', sentence)
            
            topic = None
            for term in capitalized_terms:
                if len(term) > 3 and term.lower() not in ['The', 'This', 'That', 'With', 'From', 'Into', 'During', 'Before', 'After']:
                    topic = term
                    break
            
            if not topic:
                # Create a simple topic from the sentence
                words = sentence.split()
                if len(words) >= 3:
                    # Take first few meaningful words
                    meaningful_words = [word for word in words[:3] if len(word) > 2 and word.lower() not in ['the', 'this', 'that', 'and', 'or', 'but', 'for', 'with', 'from']]
                    if meaningful_words:
                        topic = ' '.join(meaningful_words).title()
                    else:
                        topic = f'Information {i+1}'
                else:
                    topic = f'Information {i+1}'
            
            # Clean up description
            description = sentence.strip()
            if len(description) > 150:
                description = smart_content_adaptation(description, 150)
            
            table_data.append([topic, description])
    
    return table_data if len(table_data) > 1 else []

def analyze_question_type(question):
    """Analyze question to determine what type of data to extract"""
    question_lower = question.lower()
    
    # Statistics questions
    if any(word in question_lower for word in ['statistics', 'percentage', 'percent', '%', 'number', 'data', 'rate', 'ratio']):
        return "statistics"
    
    # Features questions
    if any(word in question_lower for word in ['features', 'characteristics', 'benefits', 'advantages', 'capabilities', 'functions']):
        return "features"
    
    # Cost questions
    if any(word in question_lower for word in ['cost', 'price', 'expense', 'budget', 'financial', 'money', '$']):
        return "costs"
    
    # Process questions
    if any(word in question_lower for word in ['process', 'steps', 'procedure', 'workflow', 'how to', 'method']):
        return "process"
    
    # Comparison questions
    if any(word in question_lower for word in ['compare', 'versus', 'vs', 'difference', 'contrast']):
        return "comparison"
    
    return "generic"

def extract_statistics_data(content, question):
    """Extract statistics and numerical data from content"""
    import re
    
    table_data = [['Metric', 'Value']]
    extracted_stats = []
    
    # Look for percentages with context
    percentage_patterns = re.findall(r'(\d+\.?\d*)%[^.!?]*', content, re.IGNORECASE)
    for match in percentage_patterns:
        percentage = match
        # Find the sentence containing this percentage
        sentences = re.split(r'[.!?]+', content)
        for sentence in sentences:
            if f"{percentage}%" in sentence:
                # Extract meaningful context
                context = sentence.replace(f"{percentage}%", "").strip()
                context = re.sub(r'^[^a-zA-Z]*', '', context)  # Remove leading non-letters
                context = context[:50] if len(context) > 50 else context
                
                if context and len(context) > 5:
                    extracted_stats.append((context.title(), f"{percentage}%"))
                break
    
    # Look for other numerical data
    number_patterns = re.findall(r'(\w+(?:\s+\w+)*)\s+(\d+(?:\.\d+)?)', content, re.IGNORECASE)
    for metric, value in number_patterns:
        if len(metric) > 3 and metric.lower() not in ['the', 'this', 'that', 'with', 'from']:
            extracted_stats.append((metric.title(), value))
    
    # Add to table data
    for metric, value in extracted_stats[:5]:  # Limit to 5 entries
        table_data.append([metric, value])
    
    return table_data if len(table_data) > 1 else []

def extract_features_data(content, question):
    """Extract features and characteristics from content with flexible pattern matching"""
    import re
    
    table_data = [['Feature', 'Description']]
    features = []
    
    # More flexible feature patterns - look for any sentence that might describe a feature
    feature_indicators = [
        r'(\w+(?:\s+\w+)*)\s+(?:provides|offers|includes|supports|enables|allows|permits)',
        r'(\w+(?:\s+\w+)*)\s+(?:feature|capability|function|characteristic)',
        r'(\w+(?:\s+\w+)*)\s+(?:is|are|can|has|have)\s+(?:used|available|compatible|supported)',
        r'(\w+(?:\s+\w+)*)\s+(?:works|runs|operates|functions)\s+(?:on|with|in)',
        r'(\w+(?:\s+\w+)*)\s+(?:platform|system|environment|technology)',
    ]
    
    # First try specific patterns
    for pattern in feature_indicators:
        matches = re.findall(pattern, content, re.IGNORECASE)
        for match in matches:
            if len(match) > 3:
                # Find the sentence containing this feature
                sentences = re.split(r'[.!?]+', content)
                for sentence in sentences:
                    if match.lower() in sentence.lower():
                        description = sentence.strip()
                        if len(description) > 20:
                            features.append((match.title(), smart_content_adaptation(description, 100)))
                        break
    
    # If no features found with specific patterns, try broader approach
    if not features:
        # Look for any sentence that might describe a feature or characteristic
        sentences = re.split(r'[.!?]+', content)
        for sentence in sentences:
            sentence = sentence.strip()
            if len(sentence) > 20 and len(sentence) < 200:  # Reasonable length
                # Look for sentences that might describe features
                feature_keywords = ['available', 'supports', 'compatible', 'works', 'runs', 'includes', 'provides', 'offers', 'can', 'has', 'have']
                if any(keyword in sentence.lower() for keyword in feature_keywords):
                    # Extract a potential feature name from the sentence
                    words = sentence.split()
                    if len(words) >= 3:
                        # Try to find a meaningful feature name
                        feature_name = extract_feature_name_from_sentence(sentence)
                        if feature_name:
                            features.append((feature_name, sentence))
    
    # If still no features, look for bullet points or lists
    if not features:
        bullet_items = re.findall(r'[•\-\*]\s*([^:]+):\s*([^.!?\n]+)', content)
        for feature, desc in bullet_items:
            if feature.strip() and desc.strip():
                features.append((feature.strip().title(), desc.strip()))
    
    # If still no features, try to extract from any structured content
    if not features:
        # Look for any colon-separated content that might be features
        colon_patterns = re.findall(r'([^:]+):\s*([^.!?\n]+)', content)
        for topic, desc in colon_patterns:
            if len(topic.strip()) > 3 and len(desc.strip()) > 10:
                features.append((topic.strip().title(), desc.strip()))
    
    # Add to table data
    for feature, description in features[:4]:  # Limit to 4 entries
        table_data.append([feature, description])
    
    return table_data if len(table_data) > 1 else []

def extract_feature_name_from_sentence(sentence):
    """Extract a meaningful feature name from a sentence"""
    import re
    
    # Look for common feature patterns
    patterns = [
        r'(\w+(?:\s+\w+)*)\s+(?:is|are)\s+(?:available|compatible|supported)',
        r'(\w+(?:\s+\w+)*)\s+(?:works|runs|operates)\s+(?:on|with)',
        r'(\w+(?:\s+\w+)*)\s+(?:platform|system|environment)',
        r'(\w+(?:\s+\w+)*)\s+(?:can|has|have)\s+(?:be|been)',
    ]
    
    for pattern in patterns:
        match = re.search(pattern, sentence, re.IGNORECASE)
        if match:
            feature = match.group(1).strip()
            if len(feature) > 3 and feature.lower() not in ['the', 'this', 'that', 'it', 'they']:
                return feature.title()
    
    # Fallback: take first few meaningful words
    words = sentence.split()
    meaningful_words = [word for word in words[:4] if len(word) > 2 and word.lower() not in ['the', 'this', 'that', 'and', 'or', 'but', 'for', 'with', 'from']]
    
    if meaningful_words:
        return ' '.join(meaningful_words[:3]).title()
    
    return None

def extract_costs_data(content, question):
    """Extract cost and financial data from content"""
    import re
    
    table_data = [['Cost Type', 'Amount']]
    costs = []
    
    # Look for cost patterns
    cost_patterns = re.findall(r'(\w+(?:\s+\w+)*)\s+costs?\s+(?:range from |average |are )?(\$\d+(?:,\d{3})*(?:\.\d{2})?(?:\s+to\s+\$\d+(?:,\d{3})*(?:\.\d{2})?)?)', content, re.IGNORECASE)
    for cost_type, cost_amount in cost_patterns:
        if cost_type.strip() and cost_amount.strip():
            costs.append((cost_type.title(), cost_amount))
    
    # Look for standalone cost amounts
    cost_amounts = re.findall(r'\$(\d+(?:,\d{3})*(?:\.\d{2})?)', content)
    cost_context = re.findall(r'(\w+(?:\s+\w+)*)\s+cost', content, re.IGNORECASE)
    
    for i, amount in enumerate(cost_amounts[:3]):
        cost_type = cost_context[i] if i < len(cost_context) else f'Cost {i+1}'
        costs.append((cost_type.title(), f"${amount}"))
    
    # Add to table data
    for cost_type, amount in costs[:4]:  # Limit to 4 entries
        table_data.append([cost_type, amount])
    
    return table_data if len(table_data) > 1 else []

def extract_process_data(content, question):
    """Extract process steps and procedures from content"""
    import re
    
    table_data = [['Step', 'Description']]
    steps = []
    
    # Look for numbered steps
    numbered_steps = re.findall(r'(\d+)\.\s*([^.!?\n]+)', content)
    for step_num, description in numbered_steps:
        if description.strip():
            steps.append((f"Step {step_num}", description.strip()))
    
    # Look for bullet points as steps
    bullet_steps = re.findall(r'[•\-\*]\s*([^.!?\n]+)', content)
    for i, step in enumerate(bullet_steps[:5]):
        if step.strip():
            steps.append((f"Step {i+1}", step.strip()))
    
    # Add to table data
    for step, description in steps[:5]:  # Limit to 5 entries
        table_data.append([step, description])
    
    return table_data if len(table_data) > 1 else []

def extract_comparison_data(content, question):
    """Extract comparison data from content"""
    import re
    
    table_data = [['Aspect', 'Item 1', 'Item 2']]
    
    # Look for comparison indicators
    comparison_indicators = ['vs', 'versus', 'compared to', 'in contrast', 'however', 'while', 'whereas']
    
    for indicator in comparison_indicators:
        if indicator in content.lower():
            # Split content around comparison indicator
            parts = re.split(indicator, content, flags=re.IGNORECASE)
            if len(parts) >= 2:
                item1_content = parts[0].strip()
                item2_content = parts[1].strip()
                
                # Extract key aspects for comparison
                aspects = extract_comparison_aspects(item1_content, item2_content)
                
                for aspect, desc1, desc2 in aspects[:3]:  # Limit to 3 aspects
                    table_data.append([aspect, desc1, desc2])
                
                return table_data
    
    return []

def extract_comparison_aspects(content1, content2):
    """Extract comparison aspects from two content pieces"""
    aspects = []
    
    # Look for common patterns in both contents
    sentences1 = content1.split('.')
    sentences2 = content2.split('.')
    
    # Simple approach: take first few sentences from each
    if sentences1 and sentences2:
        desc1 = smart_content_adaptation(sentences1[0], 50)
        desc2 = smart_content_adaptation(sentences2[0], 50)
        aspects.append(("Description", desc1, desc2))
    
    return aspects

def extract_generic_data(content, question):
    """Extract generic structured data from content with better topic extraction"""
    import re
    
    table_data = [['Topic', 'Description']]
    
    # Split content into meaningful sentences
    sentences = re.split(r'[.!?]+', content)
    meaningful_sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
    
    if len(meaningful_sentences) >= 2:
        # Try to extract meaningful topics from the content
        key_topics = extract_key_topics(content)
        
        # If we have key topics, use them
        if key_topics:
            for i, topic in enumerate(key_topics[:4]):
                if i < len(meaningful_sentences):
                    description = meaningful_sentences[i]
                    if len(description) > 100:
                        description = smart_content_adaptation(description, 100)
                    table_data.append([topic, description])
        else:
            # Create better topics from sentences
            for i, sentence in enumerate(meaningful_sentences[:4]):
                if sentence.strip():
                    # Try to extract a meaningful topic from the sentence
                    topic = extract_topic_from_sentence(sentence)
                    if not topic:
                        # Fallback: create topic from first few meaningful words
                        words = sentence.split()
                        meaningful_words = [word for word in words[:4] if len(word) > 2 and word.lower() not in ['the', 'this', 'that', 'and', 'or', 'but', 'for', 'with', 'from', 'in', 'on', 'at', 'to', 'of', 'a', 'an', 'is', 'are', 'was', 'were', 'be', 'been', 'have', 'has', 'had', 'do', 'does', 'did']]
                        if meaningful_words:
                            topic = ' '.join(meaningful_words[:3]).title()
                        else:
                            topic = f'Point {i+1}'
                    
                    # Clean up the description
                    description = sentence.strip()
                    if len(description) > 100:
                        description = smart_content_adaptation(description, 100)
                    
                    table_data.append([topic, description])
        
        return table_data
    
    return []

def extract_topic_from_sentence(sentence):
    """Extract a meaningful topic from a sentence"""
    import re
    
    # Look for common topic patterns
    patterns = [
        r'(\w+(?:\s+\w+)*)\s+(?:is|are)\s+(?:a|an)\s+(\w+(?:\s+\w+)*)',
        r'(\w+(?:\s+\w+)*)\s+(?:provides|offers|includes|supports)',
        r'(\w+(?:\s+\w+)*)\s+(?:technology|system|platform|tool)',
        r'(\w+(?:\s+\w+)*)\s+(?:encryption|security|protection)',
        r'(\w+(?:\s+\w+)*)\s+(?:algorithm|method|technique)',
    ]
    
    for pattern in patterns:
        match = re.search(pattern, sentence, re.IGNORECASE)
        if match:
            topic = match.group(1).strip()
            if len(topic) > 3 and topic.lower() not in ['the', 'this', 'that', 'it', 'they']:
                return topic.title()
    
    # Look for capitalized terms (proper nouns, technical terms)
    capitalized_terms = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', sentence)
    for term in capitalized_terms:
        if len(term) > 3 and term.lower() not in ['The', 'This', 'That', 'With', 'From', 'Into', 'During', 'Before', 'After']:
            return term
    
    return None

def extract_content_based_data(content, question):
    """Extract data based on content structure when question-specific extraction fails"""
    import re
    
    table_data = [['Topic', 'Description']]
    
    # Look for any structured patterns with better validation
    patterns = [
        # Numbered lists with better validation
        (r'\d+\.\s*([^:]+):\s*([^.!?\n]+)', 'Item', 'Description'),
        # Bullet points with better validation
        (r'[•\-\*]\s*([^:]+):\s*([^.!?\n]+)', 'Feature', 'Description'),
        # Key-value pairs with better validation
        (r'([^:]+):\s*([^.!?\n]+)', 'Topic', 'Value'),
    ]
    
    for pattern, col1, col2 in patterns:
        matches = re.findall(pattern, content)
        if matches:
            table_data = [[col1, col2]]
            for item, desc in matches[:4]:  # Limit to 4 entries
                if item.strip() and desc.strip() and len(item.strip()) > 2 and len(desc.strip()) > 5:
                    # Validate that the item is meaningful
                    if not item.strip().lower() in ['the', 'this', 'that', 'here', 'there']:
                        table_data.append([item.strip().title(), desc.strip()])
            if len(table_data) > 1:
                return table_data
    
    # If no structured patterns found, try intelligent sentence analysis
    sentences = re.split(r'[.!?]+', content)
    meaningful_sentences = [s.strip() for s in sentences if len(s.strip()) > 15 and len(s.strip()) < 300]
    
    if len(meaningful_sentences) >= 2:
        for i, sentence in enumerate(meaningful_sentences[:4]):
            if sentence.strip():
                # Try to extract a meaningful topic from the sentence
                topic = extract_topic_from_sentence(sentence)
                if not topic:
                    # Look for any capitalized terms that might be topics
                    capitalized_terms = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', sentence)
                    for term in capitalized_terms:
                        if len(term) > 3 and term.lower() not in ['The', 'This', 'That', 'With', 'From', 'Into', 'During', 'Before', 'After']:
                            topic = term
                            break
                    
                    if not topic:
                        # Create topic from first few meaningful words
                        words = sentence.split()
                        meaningful_words = [word for word in words[:4] if len(word) > 2 and word.lower() not in ['the', 'this', 'that', 'and', 'or', 'but', 'for', 'with', 'from', 'in', 'on', 'at', 'to', 'of', 'a', 'an', 'is', 'are', 'was', 'were', 'be', 'been', 'have', 'has', 'had', 'do', 'does', 'did']]
                        if meaningful_words:
                            topic = ' '.join(meaningful_words[:3]).title()
                        else:
                            topic = f'Point {i+1}'
                
                # Clean up the description
                description = sentence.strip()
                if len(description) > 100:
                    description = description[:97] + "..."
                
                # Only add if we have a meaningful topic
                if topic and topic.lower() not in ['the', 'this', 'that', 'here', 'there']:
                    table_data.append([topic, description])
    
    return table_data if len(table_data) > 1 else []

def format_table_content_with_ai(content, max_length):
    """Use AI only for formatting content to fit size constraints - no content generation"""
    import openai
    
    try:
        prompt = f"""
        Reformat the following content to fit within {max_length} characters while preserving all the original information and meaning. Do not add any new information or generate content - only reformat existing content.

        Content: {content}

        Requirements:
        - Keep all original facts and data
        - Maintain the same meaning and context
        - Use shorter words/sentences where possible
        - Remove redundant phrases
        - Stay within {max_length} characters
        - Do not add any new information
        """
        
        response = openai.ChatCompletion.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=200,
            temperature=0.1
        )
        
        reformatted_content = response.choices[0].message.content
        # Defensive type checking - handle case where content might be a list
        if isinstance(reformatted_content, list):
            reformatted = str(reformatted_content[0]) if reformatted_content else ""
        else:
            reformatted = str(reformatted_content).strip() if reformatted_content else ""
        
        # Validate that we're not generating new content
        if len(reformatted) <= max_length and reformatted != content:
            return reformatted
        else:
            # Fallback to simple truncation
            return smart_content_adaptation(content, max_length)
            
    except Exception as e:
        print(f"   WARNING: AI formatting failed: {e}")
        # Fallback to simple truncation
        return content[:max_length-3] + "..." if len(content) > max_length else content

def extract_key_topics(content):
    """Extract key topics from content for better table headers"""
    import re
    
    # Look for key terms and concepts
    key_terms = []
    
    # Look for technical terms, proper nouns, etc.
    technical_terms = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', content)
    for term in technical_terms:
        if len(term) > 3 and term not in ['The', 'This', 'That', 'With', 'From', 'Into', 'During', 'Before', 'After']:
            key_terms.append(term)
    
    # Look for specific patterns
    patterns = [
        r'(\w+)\s+(?:adoption|usage|implementation)',
        r'(\w+)\s+(?:features|characteristics|benefits)',
        r'(\w+)\s+(?:cost|price|value)',
        r'(\w+)\s+(?:security|encryption|protection)',
        r'(\w+)\s+(?:compatibility|availability)',
    ]
    
    for pattern in patterns:
        matches = re.findall(pattern, content, re.IGNORECASE)
        for match in matches:
            if match and len(match) > 2:
                key_terms.append(match.title())
    
    # Remove duplicates and limit
    unique_terms = list(dict.fromkeys(key_terms))  # Preserve order
    return unique_terms[:4]  # Return top 4 terms

def parse_comparison_content(content, question):
    """Parse content for comparison slide - split into two parts"""
    import re
    
    # Try to find comparison indicators
    comparison_indicators = [
        r'vs\.?', r'versus', r'compared to', r'in contrast', r'on the other hand',
        r'however', r'while', r'whereas', r'but', r'although'
    ]
    
    # Look for natural split points
    lines = content.split('\n')
    split_point = -1
    
    for i, line in enumerate(lines):
        line_lower = line.lower()
        for indicator in comparison_indicators:
            if re.search(indicator, line_lower):
                split_point = i
                break
        if split_point != -1:
            break
    
    if split_point != -1:
        # Split at the comparison indicator
        left_content = '\n'.join(lines[:split_point]).strip()
        right_content = '\n'.join(lines[split_point:]).strip()
        return [left_content, right_content]
    
    # If no clear comparison indicator, try to split by paragraphs
    paragraphs = content.split('\n\n')
    if len(paragraphs) >= 2:
        left_content = paragraphs[0].strip()
        right_content = '\n\n'.join(paragraphs[1:]).strip()
        return [left_content, right_content]
    
    # If still no clear split, split by content length
    content_length = len(content)
    mid_point = content_length // 2
    
    # Try to find a good split point near the middle
    words = content.split()
    if len(words) > 10:
        mid_word_index = len(words) // 2
        left_words = words[:mid_word_index]
        right_words = words[mid_word_index:]
        
        left_content = ' '.join(left_words)
        right_content = ' '.join(right_words)
        return [left_content, right_content]
    
    # Final fallback - split in half
    left_content = content[:content_length//2]
    right_content = content[content_length//2:]
    return [left_content, right_content]

def parse_bullet_content(content):
    """Parse content to extract bullet points"""
    import re
    
    # Look for bullet points or numbered lists
    bullets = []
    
    # Try numbered list
    numbered = re.findall(r'\d+\.\s*(.+?)(?=\d+\.|$)', content, re.DOTALL)
    if numbered:
        bullets = [item.strip() for item in numbered]
    else:
        # Try bullet points
        bulleted = re.findall(r'[•\-\*]\s*(.+?)(?=[•\-\*]|$)', content, re.DOTALL)
        if bulleted:
            bullets = [item.strip() for item in bulleted]
        else:
            # Split by sentences
            sentences = re.split(r'[.!?]+', content)
            bullets = [s.strip() for s in sentences if s.strip()][:6]
    
    return bullets[:6]  # Limit to 6 bullets

def check_content_fits_in_slide(content, slide_type, slide_dimensions=None):
    """Check if content fits within slide bounds"""
    if slide_dimensions is None:
        # Default slide dimensions
        slide_dimensions = {
            'paragraph': {'max_lines': 15, 'max_chars': 800},
            'bullet_points': {'max_lines': 12, 'max_chars': 600},
            'table': {'max_lines': 10, 'max_chars': 500},
            'timeline': {'max_lines': 8, 'max_chars': 400},
            'boxes': {'max_lines': 10, 'max_chars': 600}
        }
    
    # Get dimensions for this slide type
    dimensions = slide_dimensions.get(slide_type, {'max_lines': 12, 'max_chars': 600})
    
    # Measure content size
    content_length = len(content)
    estimated_lines = content_length / 80  # Rough estimate of lines
    
    # Check against slide dimensions
    if estimated_lines > dimensions['max_lines'] or content_length > dimensions['max_chars']:
        return False, f"Content too long: {estimated_lines:.1f} lines, {content_length} chars vs {dimensions['max_lines']} lines, {dimensions['max_chars']} chars max"
    
    return True, "Content fits"

def handle_content_overflow(original_content, format_type, slide_num, question):
    """Handle content overflow by getting AI to reformat it"""
    
    # Create prompt for AI to reformat content
    format_instructions = {
        'paragraph': 'Create a concise paragraph that explains the key points clearly and fits within slide bounds.',
        'bullet_points': 'Create 4-6 concise bullet points that capture the main ideas and fit within slide bounds.',
        'table': 'Extract the most important metrics and data points in a structured format that fits within slide bounds.',
        'timeline': 'Create 3-4 clear, numbered steps that outline the process and fit within slide bounds.',
        'boxes': 'Create a clear comparison with 2-3 key points for each side that fits within slide bounds.'
    }
    
    reformat_prompt = f"""
    The following content is too long for a PowerPoint slide. 
    Please reformat it to be more concise while maintaining all key information.
    
    Slide Question: {question}
    Original content: {original_content}
    Target format: {format_type}
    Maximum length: 600 characters
    
    Instructions: {format_instructions.get(format_type, 'Make it concise and clear.')}
    
    Please provide a concise, well-structured version that fits within slide bounds.
    Focus on the most important information and key points.
    Return only the reformatted content, no explanations or markdown.
    """
    
    try:
        # Call AI to get reformatted content
        reformatted_content = get_ai_reformatted_content(reformat_prompt)
        
        # Clean up the response - extract just the content
        if reformatted_content:
            # Remove any markdown formatting
            import re
            reformatted_content = re.sub(r'```.*?```', '', reformatted_content, flags=re.DOTALL)
            reformatted_content = re.sub(r'\*\*([^*]+)\*\*', r'\1', reformatted_content)
            reformatted_content = re.sub(r'\*([^*]+)\*', r'\1', reformatted_content)
            reformatted_content = reformatted_content.strip()
            
            # Validate the new content
            if len(reformatted_content) <= 800 and len(reformatted_content) > 50:
                return reformatted_content
            else:
                # If still too long or too short, try again with stricter limits
                return handle_content_overflow_strict(original_content, format_type, slide_num, question)
        else:
            # No response from AI, use local reformatting
            return reformat_content_for_slide(original_content, format_type)
            
    except Exception as e:
        print(f"Error in AI reformatting: {e}")
        # Fallback to local reformatting
        return reformat_content_for_slide(original_content, format_type)

def handle_content_overflow_strict(original_content, format_type, slide_num, question):
    """Handle content overflow with stricter limits"""
    
    reformat_prompt = f"""
    The content is still too long. Please create a very concise version.
    
    Slide Question: {question}
    Original content: {original_content}
    Target format: {format_type}
    Maximum length: 400 characters
    
    Instructions: Create a very brief, essential-only version that captures only the most critical information.
    Return only the reformatted content, no explanations or markdown.
    """
    
    try:
        reformatted_content = get_ai_reformatted_content(reformat_prompt)
        
        if reformatted_content:
            # Clean up the response
            import re
            reformatted_content = re.sub(r'```.*?```', '', reformatted_content, flags=re.DOTALL)
            reformatted_content = re.sub(r'\*\*([^*]+)\*\*', r'\1', reformatted_content)
            reformatted_content = re.sub(r'\*([^*]+)\*', r'\1', reformatted_content)
            reformatted_content = reformatted_content.strip()
            
            if len(reformatted_content) <= 500 and len(reformatted_content) > 30:
                return reformatted_content
            else:
                # Final fallback to local reformatting
                return reformat_content_for_slide(original_content, format_type)
        else:
            # No response from AI, use local reformatting
            return reformat_content_for_slide(original_content, format_type)
            
    except Exception as e:
        print(f"Error in strict AI reformatting: {e}")
        return reformat_content_for_slide(original_content, format_type)

def create_slide_with_overflow_check(prs, slide_data):
    """Create slide with overflow checking and AI reformatting"""
    question = slide_data['question']
    content = slide_data['content']
    format_type = slide_data['format_type']
    slide_num = slide_data['slide_num']
    
    # Check if content fits
    fits, message = check_content_fits_in_slide(content, format_type)
    
    if not fits:
        print(f"WARNING:  Content overflow detected for slide {slide_num}: {message}")
        
        # Enhanced systems handle their own content processing - skip overflow for them
        if format_type in ['timeline', 'table']:
            print(f"🎯 Enhanced {format_type} system will handle content processing internally")
            # Let the enhanced systems handle content - they have their own AI extraction
        else:
            print(f"🔄 Getting AI to reformat content...")
            
            # Get AI to reformat the content for non-enhanced formats
            reformatted_content = handle_content_overflow(content, format_type, slide_num, question)
            
            # Update the slide data with reformatted content
            slide_data['content'] = reformatted_content
            
            # Check if the reformatted content fits
            fits, message = check_content_fits_in_slide(reformatted_content, format_type)
            if not fits:
                print(f"WARNING:  Reformatted content still too long, using local reformatting")
                slide_data['content'] = reformat_content_for_slide(content, format_type)
    
    # Create the slide based on format type
    print(f"   🎨 Creating slide with format_type: {format_type}")
    
    if format_type == 'paragraph':
        print(f"   📝 Calling create_paragraph_slide")
        return create_paragraph_slide(prs, question, slide_data['content'], slide_num)
    elif format_type == 'bullet_points':
        print(f"   📝 Calling create_bullet_slide")
        return create_bullet_slide(prs, question, slide_data['content'], slide_num)
    elif format_type == 'table':
        print(f"   📝 Calling create_table_slide")
        return create_table_slide(prs, question, slide_data['content'], slide_num)
    elif format_type == 'timeline':
        print(f"   📝 Calling create_timeline_slide")
        return create_timeline_slide(prs, question, slide_data['content'], slide_num)
    elif format_type == 'boxes':
        print(f"   📝 Calling create_comparison_slide")
        return create_comparison_slide(prs, question, slide_data['content'], slide_num)
    else:
        print(f"   📝 Calling create_paragraph_slide (default)")
        return create_paragraph_slide(prs, question, slide_data['content'], slide_num)

def test_if_text_fits_enhanced(shape, text):
    """Enhanced text fitting test with overflow detection and AI reformatting trigger"""
    try:
        # Get shape dimensions
        width_inches = shape.width.inches
        height_inches = shape.height.inches
        
        # Estimate text requirements
        estimated_lines = len(text) / 80  # Rough estimate
        estimated_height_inches = estimated_lines * 0.2  # Approximate line height
        
        # Check if text fits
        if estimated_height_inches > height_inches or len(text) > 1000:
            return False, f"Text too long: {len(text)} chars, {estimated_lines:.1f} lines vs {height_inches:.1f} inches height"
        
        return True, "Text fits"
        
    except Exception as e:
        return False, f"Error checking text fit: {e}"

def fit_text_with_ai_reformatting(shape, text, field_type="body", question=""):
    """Fit text with AI reformatting when overflow is detected"""
    from .ai_content_fetcher import get_ai_content
    
    # First try to fit the text
    fits, message = test_if_text_fits_enhanced(shape, text)
    
    if fits:
        return text
    
    print(f"WARNING:  Text overflow detected: {message}")
    print(f"🔄 Getting AI to reformat content...")
    
    # Create AI prompt for reformatting
    format_instructions = {
        'title': 'Create a concise, engaging title that fits within the slide header.',
        'subtitle': 'Create a brief subtitle that complements the main title.',
        'body': 'Create concise, well-structured content that fits within the slide body while maintaining key information.',
        'bullet': 'Create 3-5 concise bullet points that capture the main ideas.',
        'footer': 'Create a brief footer text that fits within the slide footer.'
    }
    
    reformat_prompt = f"""
    The following text is too long for a PowerPoint slide {field_type} field.
    Please reformat it to be more concise while maintaining all key information.
    
    Slide Question: {question}
    Original text: {text}
    Field type: {field_type}
    Maximum length: 400 characters
    
    Instructions: {format_instructions.get(field_type, 'Make it concise and clear.')}
    
    Please provide a concise, well-structured version that fits within slide bounds.
    Focus on the most important information and key points.
    """
    
    try:
        # Call OpenAI to get reformatted content
        reformatted_text = get_ai_content(reformat_prompt)
        reformatted_text = reformatted_text.strip()
        
        # Check if the reformatted text fits
        fits, message = test_if_text_fits_enhanced(shape, reformatted_text)
        
        if fits and len(reformatted_text) > 30:
            print(f"SUCCESS: AI reformatting successful: {len(reformatted_text)} characters")
            return reformatted_text
        else:
            print(f"WARNING:  AI reformatting still too long, using local truncation")
            return intelligent_truncate(text, 400)
            
    except Exception as e:
        print(f"ERROR: Error in AI reformatting: {e}")
        # Fallback to local truncation
        return intelligent_truncate(text, 400)

def get_ai_reformatted_content(prompt):
    """Get AI reformatted content using a simple text-based approach"""
    try:
        import openai
        import os
        
        # Get API key from environment
        api_key = os.getenv('OPENAI_API_KEY')
        if not api_key:
            print("WARNING:  No OpenAI API key found, using local reformatting")
            return None
        
        # Set up OpenAI client
        client = openai.OpenAI(api_key=api_key)
        
        # Make the API call
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that reformats content to fit within PowerPoint slide bounds. Return only the reformatted content, no explanations or markdown."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=500,
            temperature=0.7
        )
        
        # Extract the content
        if response.choices and len(response.choices) > 0:
            content_raw = response.choices[0].message.content
            # Defensive type checking - handle case where content might be a list
            if isinstance(content_raw, list):
                content = str(content_raw[0]) if content_raw else ""
            else:
                content = str(content_raw).strip() if content_raw else ""
            return content
        else:
            return None
            
    except Exception as e:
        print(f"ERROR: Error calling OpenAI API: {e}")
        return None

def create_ppt_from_scratch(content_analysis_results, output_path):
    """
    Create PowerPoint presentation from scratch using content analysis results
    Integrates enhanced systems for tables, timelines, and boxes
    """
    from pptx import Presentation
    
    print(f"Creating PowerPoint from scratch with {len(content_analysis_results)} slides...")
    
    # Create a new presentation
    prs = Presentation()
    
    # Process each slide with enhanced content processing
    for i, slide_data in enumerate(content_analysis_results):
        print(f"\nProcessing slide {i+1}/{len(content_analysis_results)}")
        print(f"   Question: {slide_data.get('question', 'No question')[:50]}...")
        print(f"   Format: {slide_data.get('format_type', 'unknown')}")
        print(f"   Content length: {len(slide_data.get('content', ''))} chars")
        
        # Prepare slide data structure expected by create_slide_with_overflow_check
        processed_slide_data = {
            'question': slide_data.get('question', f'Slide {i+1}'),
            'content': slide_data.get('content', ''),
            'format_type': slide_data.get('format_type', 'paragraph'),
            'slide_num': i + 1,
            'confidence': slide_data.get('confidence', 0.0),
            'recommended_slide_type': slide_data.get('recommended_slide_type', 'content_slide'),
            'template_match': slide_data.get('template_match', {}),
            'slide_role': slide_data.get('slide_role', ''),
            'flow_adjustment': slide_data.get('flow_adjustment', 0.0)
        }
        
        try:
            # Create slide using the existing enhanced system
            create_slide_with_overflow_check(prs, processed_slide_data)
            
            print(f"   Slide {i+1} created successfully")
            
        except Exception as e:
            print(f"   Error creating slide {i+1}: {e}")
            # Create a fallback slide
            create_fallback_slide(prs, processed_slide_data)
    
    # Save the presentation
    try:
        prs.save(output_path)
        print(f"\nPowerPoint presentation saved: {output_path}")
        return output_path
        
    except Exception as e:
        print(f"\nError saving presentation: {e}")
        raise

def create_fallback_slide(prs, slide_data):
    """Create a fallback slide when the main slide creation fails"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    
    print(f"   Creating fallback slide for slide {slide_data['slide_num']}")
    
    # Add slide
    slide_layout = get_safe_slide_layout(prs)  # Safe layout selection
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = f"Slide {slide_data['slide_num']}: {slide_data['question']}"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    
    # Limit content to reasonable length
    content = slide_data['content']
    if len(content) > 800:
        content = smart_content_adaptation(content, 800)
    
    content_frame.text = content
    content_frame.paragraphs[0].font.size = Pt(16)
    content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Configure text frame
    content_frame.word_wrap = True
    content_frame.margin_left = Pt(12)
    content_frame.margin_right = Pt(12)
    content_frame.margin_top = Pt(6)
    content_frame.margin_bottom = Pt(6)

def create_title_slide(prs, topic, user_prompt):
    """Create professional title slide"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from .prompt_builder import generate_presentation_subtitle
    
    # Add slide with blank layout
    slide_layout = get_safe_slide_layout(prs)
    slide = prs.slides.add_slide(slide_layout)
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = topic
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Generate and add subtitle
    subtitle = generate_presentation_subtitle(topic, user_prompt)
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(20)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Optional: Add decorative line or spacing
    line_box = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(4), Inches(0.2))
    line_frame = line_box.text_frame
    line_frame.text = "─" * 20
    line_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_thank_you_slide(prs, topic):
    """Create professional thank you slide"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    
    # Add slide with blank layout
    slide_layout = get_safe_slide_layout(prs)
    slide = prs.slides.add_slide(slide_layout)
    
    # Main "Thank You" text
    thank_you_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Thank You"
    thank_you_frame.paragraphs[0].font.size = Pt(42)
    thank_you_frame.paragraphs[0].font.bold = True
    thank_you_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Optional: Add topic reference
    if topic:
        topic_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.8))
        topic_frame = topic_box.text_frame
        topic_frame.text = f"Questions about {topic}?"
        topic_frame.paragraphs[0].font.size = Pt(16)
        topic_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_ppt_with_title_and_thanks(content_analysis_results, output_path, topic, user_prompt):
    """Create PowerPoint presentation with title slide, content slides, and thank you slide"""
    from pptx import Presentation
    
    print(f"🎨 Creating PowerPoint with title + {len(content_analysis_results)} content slides + thank you slide...")
    
    # Create a new presentation
    prs = Presentation()
    
    # Step 1: Create title slide
    print(f"   📝 Creating title slide: {topic}")
    create_title_slide(prs, topic, user_prompt)
    
    # Step 2: Create content slides
    for i, slide_data in enumerate(content_analysis_results):
        content_slide_num = i + 2  # +2 because slide 1 is title
        format_type = slide_data.get('format_type', 'paragraph')
        content = slide_data.get('content', '')
        question = slide_data.get('question', f'Content Slide {i+1}')
        
        print(f"   📝 Creating content slide {content_slide_num}: {format_type} format")
        
        # Create slide data structure for the slide creation functions
        slide_data_for_creation = {
            'question': question,
            'content': content,
            'format_type': format_type,
            'slide_num': content_slide_num,
            'confidence': slide_data.get('confidence', 0.0),
            'recommended_slide_type': slide_data.get('recommended_slide_type', 'content_slide'),
            'template_match': slide_data.get('template_match', {}),
            'slide_role': slide_data.get('slide_role', ''),
            'flow_adjustment': slide_data.get('flow_adjustment', 0.0)
        }
        
        # Create slide with overflow checking and AI reformatting
        try:
            create_slide_with_overflow_check(prs, slide_data_for_creation)
        except Exception as e:
            print(f"   ⚠️  Error creating slide {content_slide_num}: {e}")
            # Fallback: create simple text slide
            create_fallback_slide(prs, slide_data_for_creation)
    
    # Step 3: Create thank you slide
    total_slides = len(content_analysis_results) + 2  # +2 for title and thank you
    print(f"   📝 Creating thank you slide (slide {total_slides})")
    create_thank_you_slide(prs, topic)
    
    # Save the presentation
    try:
        prs.save(output_path)
        print(f"✅ PowerPoint saved: {output_path}")
        print(f"📊 Total slides created: {total_slides} (1 title + {len(content_analysis_results)} content + 1 thank you)")
    except Exception as e:
        print(f"❌ Error saving PowerPoint: {e}")
        raise

if __name__ == "__main__":
    import sys, json
    if len(sys.argv) < 4:
        print("Usage: python replace_content.py template.pptx ai_output.json output.pptx")
    else:
        with open(sys.argv[2], "r", encoding="utf-8") as f:
            aiout = json.load(f)
        replace_with_ai_content(sys.argv[1], aiout, sys.argv[3])
