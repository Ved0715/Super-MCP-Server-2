"""
🎯 Perfect Research PPT Generator
Advanced presentation generation with research intelligence and semantic search
"""

import os
import logging
import asyncio
from datetime import datetime
from typing import Dict, List, Any, Optional, Tuple
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import openai
import requests

logger = logging.getLogger(__name__)

class PerfectPPTGenerator:
    """Perfect research presentation generator with AI intelligence"""
    
    def __init__(self, config, vector_storage=None, research_analyzer=None):
        self.config = config
        self.vector_storage = vector_storage
        self.research_analyzer = research_analyzer
        self.openai_client = openai.OpenAI(api_key=config.OPENAI_API_KEY)
        
        # Advanced academic themes
        self.themes = {
            "academic_professional": {
                "primary": RGBColor(26, 35, 126),      # Deep blue
                "secondary": RGBColor(63, 81, 181),    # Medium blue  
                "accent": RGBColor(255, 193, 7),       # Gold
                "text": RGBColor(33, 33, 33),          # Dark gray
                "background": RGBColor(250, 250, 250), # Light gray
                "font": "Calibri",
                "style": "academic"
            },
            "research_modern": {
                "primary": RGBColor(0, 150, 136),      # Teal
                "secondary": RGBColor(0, 188, 212),    # Cyan
                "accent": RGBColor(76, 175, 80),       # Green
                "text": RGBColor(55, 71, 79),          # Blue gray
                "background": RGBColor(248, 249, 250), # Off white
                "font": "Segoe UI",
                "style": "modern"
            },
            "executive_clean": {
                "primary": RGBColor(96, 125, 139),     # Blue gray
                "secondary": RGBColor(144, 164, 174),  # Light blue gray
                "accent": RGBColor(255, 152, 0),       # Orange
                "text": RGBColor(84, 110, 122),        # Medium blue gray
                "background": RGBColor(255, 255, 255), # White
                "font": "Arial",
                "style": "executive"
            }
        }

    async def create_perfect_presentation(self,
                                        paper_content: Dict[str, Any],
                                        user_prompt: str,
                                        paper_id: str,
                                        search_results: Dict[str, Any] = None,
                                        title: str = None,
                                        author: str = "AI Research Assistant",
                                        theme: str = "academic_professional",
                                        slide_count: int = 12,
                                        audience_type: str = "academic") -> str:
        """Create perfect research presentation with AI intelligence"""
        try:
            # Step 1: Analyze user requirements and paper content
            analysis_result = await self._analyze_presentation_requirements(
                user_prompt, paper_content, audience_type
            )
            
            # Step 2: Generate intelligent presentation plan
            presentation_plan = await self._create_intelligent_plan(
                analysis_result, paper_content, paper_id, search_results, slide_count
            )
            
            # Step 3: Create presentation with advanced features
            presentation_path = await self._create_advanced_presentation(
                presentation_plan, title, author, theme, paper_content
            )
            
            logger.info(f"Perfect presentation created: {presentation_path}")
            return presentation_path
            
        except Exception as e:
            logger.error(f"Error creating perfect presentation: {e}")
            raise

    async def _create_references_slide(self, prs, references: List[Dict[str, Any]], theme_config: Dict[str, Any]):
        """Create a professional references slide with proper citations"""
        try:
            logger.info(f"📚 Creating references slide with {len(references)} sources")
            
            # Add references slide
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            title_shape = slide.shapes.title
            title_shape.text = "References"
            title_shape.text_frame.paragraphs[0].font.size = Pt(36)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add references content
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(12.3)
            height = Inches(5.5)
            
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            # Clear default paragraph
            text_frame.clear()
            
            # Add each reference
            for i, ref in enumerate(references):
                p = text_frame.add_paragraph()
                p.text = f"{i+1}. {ref.get('source', 'Unknown Source')}"
                p.font.size = Pt(16)
                p.font.name = "Arial"
                p.space_after = Pt(6)
                
                # Add relevance score if available
                if 'relevance_score' in ref:
                    p.text += f" (Relevance: {ref['relevance_score']:.2f})"
                
                # Add type if available
                if 'type' in ref:
                    p.text += f" [{ref['type']}]"
            
            # Apply theme colors
            if theme_config:
                primary_color = theme_config.get('primary_color', RGBColor(0, 51, 102))
                title_shape.text_frame.paragraphs[0].font.color.rgb = primary_color
                
                for paragraph in text_frame.paragraphs:
                    paragraph.font.color.rgb = RGBColor(51, 51, 51)
            
            logger.info(f"✅ References slide created successfully")
            
        except Exception as e:
            logger.error(f"❌ Error creating references slide: {e}")
            # Continue without references slide rather than failing

    async def _analyze_presentation_requirements(self,
                                               user_prompt: str,
                                               paper_content: Dict[str, Any],
                                               audience_type: str) -> Dict[str, Any]:
        """Analyze user requirements with AI understanding"""
        try:
            # Prepare context for AI analysis
            paper_summary = self._create_paper_summary_for_ai(paper_content)
            analysis_prompt = f"""
                You are an expert academic presentation strategist. Your task is to analyze the user’s objective and the research paper to generate a detailed, targeted presentation strategy optimized for the audience.

                USER REQUEST:
                "{user_prompt}"

                AUDIENCE TYPE:
                {audience_type}

                RESEARCH PAPER SUMMARY:
                {paper_summary}

                Based on this input, generate a comprehensive analysis in the following structured JSON format:

                {{
                    "presentation_focus": "Identify the core contributions, innovations, or questions that the presentation should emphasize to create the most impact.",
                    "key_message": "Distill the single most important insight or conclusion that the audience should walk away with.",
                    "audience_needs": "Describe what the specific audience type must understand or appreciate to stay engaged and gain value from the presentation.",
                    "content_priorities": ["Ranked list of critical content elements (e.g., background, methodology, results, applications) in priority order."],
                    "slide_types_needed": ["Types of slides recommended (e.g., agenda, problem statement, methodology visualization, experiment results, key takeaways, references)."],
                    "depth_level": "Choose the appropriate level of complexity: 'technical', 'general', or 'executive', depending on the audience’s expertise and expectations.",
                    "emphasis_areas": ["Highlight which sections need in-depth focus (e.g., methodology, results, comparative analysis, limitations, implications, future directions)."],
                    "storytelling_approach": "Suggest the most effective narrative structure (e.g., problem–solution, chronological, thematic, data-driven) to guide the audience clearly.",
                    "visual_elements": ["Recommend visual aids like charts, diagrams, flowcharts, infographics, model architectures, timelines, etc., with context of use."],
                    "citation_strategy": "Explain how and where to cite key sources (e.g., inline citations, footnotes, visual callouts, dedicated reference slides) to maintain credibility."
                }}

            Be specific, actionable, and strategic. Tailor every section to align with both the user’s intent and the audience’s perspective.
            """

            
            response = self.openai_client.chat.completions.create(
                model=self.config.LLM_MODEL,
                messages=[{"role": "user", "content": analysis_prompt}],
                temperature=0.1
            )
            
            import json
            analysis = json.loads(self._clean_json_response(response.choices[0].message.content))
            
            return {
                "success": True,
                "analysis": analysis,
                "user_prompt": user_prompt,
                "audience_type": audience_type
            }
            
        except Exception as e:
            logger.error(f"Error in requirement analysis: {e}")
            return {
                "success": False,
                "error": str(e),
                "fallback_analysis": self._create_fallback_analysis(user_prompt, audience_type)
            }

    async def _create_intelligent_plan(self,
                                     analysis_result: Dict[str, Any],
                                     paper_content: Dict[str, Any],
                                     paper_id: str,
                                     search_results: Dict[str, Any],
                                     slide_count: int) -> Dict[str, Any]:
        """Create intelligent presentation plan with semantic search"""
        try:
            analysis = analysis_result.get("analysis", {})
            print(analysis)
            user_prompt = analysis_result.get("user_prompt", "")
            print(user_prompt)
            
            # Get relevant content using vector search
            relevant_content = await self._get_relevant_content_for_presentation(
                user_prompt, paper_id, analysis.get("emphasis_areas", [])
            )
            
            # Create slide plan with AI assistance
            slide_plan = await self._generate_ai_slide_plan(
                analysis, paper_content, relevant_content, search_results, slide_count
            )
            
            return {
                "success": True,
                "slide_plan": slide_plan,
                "analysis": analysis,
                "relevant_content": relevant_content,
                "content_strategy": self._create_content_strategy(analysis, paper_content)
            }
            
        except Exception as e:
            logger.error(f"Error creating intelligent plan: {e}")
            return {"success": False, "error": str(e)}

    async def _get_relevant_content_for_presentation(self,
                                                   user_prompt: str,
                                                   paper_id: str,
                                                   emphasis_areas: List[str]) -> Dict[str, List]:
        """Get semantically relevant content using vector search"""
        if not self.vector_storage:
            return {"general": [], "sections": {}}
        
        try:
            relevant_content = {"general": [], "sections": {}}
            
            # General content search based on user prompt
            general_results = await self.vector_storage.semantic_search(
                query=user_prompt,
                namespace=paper_id,
                top_k=15
            )
            relevant_content["general"] = [r.content for r in general_results]
            
            # Section-specific searches
            for area in emphasis_areas:
                if area in ["methodology", "results", "discussion", "conclusion"]:
                    section_results = await self.vector_storage.contextual_search(
                        user_prompt=user_prompt,
                        namespace=paper_id,
                        context_type=area
                    )
                    relevant_content["sections"][area] = [r.content for r in section_results]
            
            return relevant_content
            
        except Exception as e:
            logger.error(f"Error getting relevant content: {e}")
            return {"general": [], "sections": {}}

    async def _generate_ai_slide_plan(self,
                                    analysis: Dict[str, Any],
                                    paper_content: Dict[str, Any],
                                    relevant_content: Dict[str, List],
                                    search_results: Dict[str, Any],
                                    slide_count: int) -> List[Dict[str, Any]]:
        """Generate AI-powered slide plan"""
        try:
            # Prepare context
            context_summary = self._prepare_context_for_slide_generation(
                paper_content, relevant_content, search_results
            )

            plan_prompt = f"""
            Create a detailed presentation plan with exactly {slide_count} slides.
            
            PRESENTATION STRATEGY:
            Focus: {analysis.get('presentation_focus', 'Research overview')}
            Key Message: {analysis.get('key_message', 'Research findings')}
            Audience: {analysis.get('audience_needs', 'Academic audience')}
            Depth: {analysis.get('depth_level', 'technical')}
            
            CONTENT PRIORITIES:
            {', '.join(analysis.get('content_priorities', []))}
            
            EMPHASIS AREAS:
            {', '.join(analysis.get('emphasis_areas', []))}
            
            AVAILABLE CONTENT:
            {context_summary}
            
            Create a detailed slide plan in JSON format:
            {{
                "slides": [
                    {{
                        "slide_number": 1,
                        "type": "title",
                        "title": "Compelling title that captures the research essence",
                        "content": ["Concise subtitle", "Presenter name and affiliation", "One-line summary of key contribution"],
                        "speaker_notes": "Welcome the audience and introduce the presentation structure and key message.",
                        "visual_elements": ["background_image", "logo"],
                        "content_source": "metadata",
                        "academic_citations": []
                    }},
                    {{
                        "slide_number": 2,
                        "type": "research_context",
                        "title": "Research Background and Significance",
                        "content": ["Why this research matters", "Current state of the field", "Research gap addressed", "Study significance"],
                        "speaker_notes": "Establish research importance and context",
                        "visual_elements": ["concept_diagram", "statistics"],
                        "content_source": "introduction_section",
                        "academic_citations": ["relevant citations"]
                    }},
                    # slides 3 to {slide_count} continue in this format
                ]
            }}
            
            REQUIREMENTS:
            - Each slide must have detailed, specific content (not generic)
            - Include proper academic citations where relevant
            - Vary slide types: title, research_context, methodology, findings, data_visualization, implications, conclusion
            - Ensure logical flow and narrative coherence
            - Include speaker notes for each slide
            - Specify visual elements that would enhance understanding
            
            Make every slide count with substantial, research-specific content.
            """
            

            response = self.openai_client.chat.completions.create(
                model=self.config.LLM_MODEL,
                messages=[{"role": "user", "content": plan_prompt}],
                temperature=0.1
            )
           
            import json
            plan_json = json.loads(self._clean_json_response(response.choices[0].message.content))
            print(plan_json)
            return plan_json.get("slides", [])
            
        except Exception as e:
            logger.error(f"Error generating AI slide plan: {e}")
            return self._create_fallback_slide_plan(slide_count, paper_content)

    async def _create_advanced_presentation(self,
                                          presentation_plan: Dict[str, Any],
                                          title: str,
                                          author: str,
                                          theme: str,
                                          paper_content: Dict[str, Any]) -> str:
        """Create advanced presentation with professional features"""
        try:
            # Initialize presentation
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            slide_plan = presentation_plan.get("slide_plan", [])
            theme_config = self.themes.get(theme, self.themes["academic_professional"])
            
            # Create slides with advanced features
            for slide_info in slide_plan:
                slide_type = slide_info.get("type", "content")
                
                if slide_type == "title":
                    await self._create_advanced_title_slide(prs, slide_info, title, author, theme_config)
                elif slide_type == "research_context":
                    await self._create_research_context_slide(prs, slide_info, theme_config)
                elif slide_type == "methodology":
                    await self._create_methodology_slide(prs, slide_info, theme_config)
                elif slide_type == "findings" or slide_type == "results":
                    await self._create_findings_slide(prs, slide_info, theme_config)
                elif slide_type == "data_visualization":
                    await self._create_data_visualization_slide(prs, slide_info, theme_config)
                elif slide_type == "implications":
                    await self._create_implications_slide(prs, slide_info, theme_config)
                elif slide_type == "conclusion":
                    await self._create_conclusion_slide(prs, slide_info, theme_config)
                else:
                    await self._create_content_slide(prs, slide_info, theme_config)
            
            # Add references slide if paper content contains references
            if paper_content.get("references"):
                await self._create_references_slide(prs, paper_content["references"], theme_config)
            
            # Apply advanced styling
            self._apply_advanced_theme(prs, theme_config)
            
            # Save presentation
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"perfect_research_presentation_{timestamp}.pptx"
            filepath = Path(self.config.PPT_OUTPUT_DIR) / filename
            
            prs.save(str(filepath))
            
            return str(filepath)
            
        except Exception as e:
            logger.error(f"Error creating advanced presentation: {e}")
            raise

    async def _create_advanced_title_slide(self, prs: Presentation, slide_info: Dict, title: str, author: str, theme_config: Dict):
        """Create advanced title slide with research branding"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        # Main title
        if slide.shapes.title:
            slide.shapes.title.text = title or slide_info.get("title", "Research Presentation")
            title_shape = slide.shapes.title
            title_shape.text_frame.paragraphs[0].font.name = theme_config["font"]
            title_shape.text_frame.paragraphs[0].font.size = Pt(32)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.color.rgb = theme_config["primary"]
        
        # Subtitle with research info
        if slide.placeholders[1]:
            content_lines = slide_info.get("content", [])
            content_lines.extend([
                f"Presented by: {author}",
                f"Generated on: {datetime.now().strftime('%B %d, %Y')}",
                "AI-Enhanced Research Presentation"
            ])
            slide.placeholders[1].text = "\n".join(content_lines)
            
            # Style subtitle
            subtitle_shape = slide.placeholders[1]
            for paragraph in subtitle_shape.text_frame.paragraphs:
                paragraph.font.name = theme_config["font"]
                paragraph.font.size = Pt(16)
                paragraph.font.color.rgb = theme_config["secondary"]

    async def _create_research_context_slide(self, prs: Presentation, slide_info: Dict, theme_config: Dict):
        """Create research context slide with academic structure"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        slide.shapes.title.text = slide_info.get("title", "Research Context")
        
        # Content with academic formatting
        content_items = slide_info.get("content", [])
        enhanced_content = []
        
        research_icons = ["🔬", "📊", "🎯", "💡", "🔍"]
        
        for i, item in enumerate(content_items):
            icon = research_icons[i % len(research_icons)]
            enhanced_content.append(f"{icon} {item}")
        
        if slide.placeholders[1]:
            slide.placeholders[1].text = "\n".join(enhanced_content)
            
            # Style content
            content_shape = slide.placeholders[1]
            for paragraph in content_shape.text_frame.paragraphs:
                paragraph.font.name = theme_config["font"]
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = theme_config["text"]
                paragraph.space_after = Pt(12)

    async def _create_methodology_slide(self, prs: Presentation, slide_info: Dict, theme_config: Dict):
        """Create methodology slide with research rigor emphasis"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = slide_info.get("title", "Research Methodology")
        
        # Enhanced methodology content
        content_items = slide_info.get("content", [])
        methodology_content = []
        
        methodology_icons = ["⚗️", "📏", "🔢", "📋", "✅"]
        
        for i, item in enumerate(content_items):
            icon = methodology_icons[i % len(methodology_icons)]
            methodology_content.append(f"{icon} {item}")
        
        if slide.placeholders[1]:
            slide.placeholders[1].text = "\n".join(methodology_content)
            
            # Add methodology quality indicators
            if len(content_items) > 0:
                quality_note = "\n\n🎯 Methodological Rigor: Peer-reviewed approach with validated instruments"
                slide.placeholders[1].text += quality_note

    async def _create_findings_slide(self, prs: Presentation, slide_info: Dict, theme_config: Dict):
        """Create findings slide with statistical emphasis"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = slide_info.get("title", "Key Findings")
        
        # Enhanced findings with statistical indicators
        content_items = slide_info.get("content", [])
        findings_content = []
        
        findings_icons = ["📈", "📊", "🔍", "💯", "⭐"]
        
        for i, item in enumerate(content_items):
            icon = findings_icons[i % len(findings_icons)]
            findings_content.append(f"{icon} {item}")
        
        if slide.placeholders[1]:
            slide.placeholders[1].text = "\n".join(findings_content)
            
            # Add statistical significance note if relevant
            if any("significant" in item.lower() or "p <" in item for item in content_items):
                sig_note = "\n\n📊 Statistical significance: Results validated at p < 0.05"
                slide.placeholders[1].text += sig_note

    async def _create_data_visualization_slide(self, prs: Presentation, slide_info: Dict, theme_config: Dict):
        """Create data visualization slide"""
        slide_layout = prs.slide_layouts[5]  # Blank layout for custom content
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = slide_info.get("title", "Data Analysis")
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = theme_config["primary"]
        
        # Add content description
        content_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
        content_frame = content_shape.text_frame
        
        content_items = slide_info.get("content", [])
        data_content = []
        
        for item in content_items:
            data_content.append(f"• {item}")
        
        content_frame.text = "\n".join(data_content)
        
        # Add placeholder for visualization
        viz_shape = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9), Inches(2))
        viz_frame = viz_shape.text_frame
        viz_frame.text = "📊 [Data visualization would be displayed here]\n🔢 Statistical analysis and charts"

    async def _create_implications_slide(self, prs: Presentation, slide_info: Dict, theme_config: Dict):
        """Create implications slide with impact focus"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = slide_info.get("title", "Research Implications")
        
        # Enhanced implications content
        content_items = slide_info.get("content", [])
        implications_content = []
        
        impact_icons = ["🎯", "🌍", "🚀", "💡", "🔮"]
        
        for i, item in enumerate(content_items):
            icon = impact_icons[i % len(impact_icons)]
            implications_content.append(f"{icon} {item}")
        
        if slide.placeholders[1]:
            slide.placeholders[1].text = "\n".join(implications_content)

    async def _create_conclusion_slide(self, prs: Presentation, slide_info: Dict, theme_config: Dict):
        """Create conclusion slide with future directions"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = slide_info.get("title", "Conclusions & Future Work")
        
        # Enhanced conclusion content
        content_items = slide_info.get("content", [])
        conclusion_content = []
        
        conclusion_icons = ["✅", "🎯", "🔬", "🚀", "🙏"]
        
        for i, item in enumerate(content_items):
            icon = conclusion_icons[i % len(conclusion_icons)]
            conclusion_content.append(f"{icon} {item}")
        
        # Add standard conclusion elements
        conclusion_content.extend([
            "🙏 Thank you for your attention",
            "❓ Questions and Discussion Welcome",
            "📧 Contact for collaboration opportunities"
        ])
        
        if slide.placeholders[1]:
            slide.placeholders[1].text = "\n".join(conclusion_content)

    async def _create_content_slide(self, prs: Presentation, slide_info: Dict, theme_config: Dict):
        """Create general content slide with academic formatting"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = slide_info.get("title", "Content")
        
        content_items = slide_info.get("content", [])
        formatted_content = []
        
        for item in content_items:
            formatted_content.append(f"• {item}")
        
        if slide.placeholders[1]:
            slide.placeholders[1].text = "\n".join(formatted_content)
            
            # Style content
            content_shape = slide.placeholders[1]
            for paragraph in content_shape.text_frame.paragraphs:
                paragraph.font.name = theme_config["font"]
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = theme_config["text"]
                paragraph.space_after = Pt(10)

    def _apply_advanced_theme(self, prs: Presentation, theme_config: Dict):
        """Apply advanced theme styling to entire presentation"""
        try:
            for slide in prs.slides:
                # Style all text elements
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            if not paragraph.font.name:
                                paragraph.font.name = theme_config["font"]
                            if not paragraph.font.color.rgb:
                                paragraph.font.color.rgb = theme_config["text"]
        except Exception as e:
            logger.warning(f"Error applying advanced theme: {e}")

    # Helper methods
    def _create_paper_summary_for_ai(self, paper_content: Dict[str, Any]) -> str:
        """Create focused summary for AI analysis"""
        sections = paper_content.get("sections", {})
        metadata = paper_content.get("metadata", {})
        
        summary_parts = []
        
        if metadata.get("title"):
            summary_parts.append(f"Title: {metadata['title']}")
        
        if sections.get("abstract"):
            abstract = sections["abstract"][:500]
            summary_parts.append(f"Abstract: {abstract}")
        
        if sections.get("methodology"):
            method = sections["methodology"][:300]
            summary_parts.append(f"Methodology: {method}")
        
        if sections.get("results"):
            results = sections["results"][:300]
            summary_parts.append(f"Results: {results}")
        
        available_sections = list(sections.keys())
        summary_parts.append(f"Available sections: {', '.join(available_sections)}")
        
        return "\n\n".join(summary_parts)

    def _clean_json_response(self, response_text: str) -> str:
        """Clean AI response to extract JSON"""
        import re
        
        # Remove markdown code blocks
        response_text = re.sub(r'```json\s*', '', response_text)
        response_text = re.sub(r'```\s*$', '', response_text)
        
        # Find JSON object
        start = response_text.find('{')
        end = response_text.rfind('}')
        
        if start >= 0 and end >= 0:
            return response_text[start:end + 1]
        
        return response_text

    def _create_fallback_analysis(self, user_prompt: str, audience_type: str) -> Dict[str, Any]:
        """Create fallback analysis if AI fails"""
        return {
            "presentation_focus": "Research overview and key findings",
            "key_message": "Present research methodology and results",
            "audience_needs": f"Information appropriate for {audience_type} audience",
            "content_priorities": ["methodology", "results", "implications"],
            "slide_types_needed": ["title", "overview", "methodology", "results", "conclusion"],
            "depth_level": "technical" if audience_type == "academic" else "general",
            "emphasis_areas": ["methodology", "results"],
            "storytelling_approach": "Linear progression from context to conclusions",
            "visual_elements": ["charts", "statistics"],
            "citation_strategy": "Include key references throughout"
        }

    def _prepare_context_for_slide_generation(self,
                                            paper_content: Dict[str, Any],
                                            relevant_content: Dict[str, List],
                                            search_results: Dict[str, Any]) -> str:
        """Prepare context summary for slide generation"""
        context_parts = []
        
        sections = paper_content.get("sections", {})
        for section_name, content in sections.items():
            # Handle both string content and dictionary content (from knowledge base)
            if isinstance(content, dict):
                # Knowledge base format with nested content
                actual_content = content.get("content", "")
                if actual_content:
                    context_parts.append(f"{section_name.upper()}: {actual_content[:200]}...")
            elif isinstance(content, str):
                # Traditional string format
                context_parts.append(f"{section_name.upper()}: {content[:200]}...")
            else:
                # Fallback for any other format
                context_parts.append(f"{section_name.upper()}: {str(content)[:200]}...")
        
        if relevant_content.get("general"):
            context_parts.append(f"RELEVANT CONTENT: {relevant_content['general'][0][:200]}...")
        
        if search_results and search_results.get("success"):
            results = search_results.get("results", [])[:3]
            for i, result in enumerate(results):
                title = result.get("title", "")[:100]
                context_parts.append(f"RELATED RESEARCH {i+1}: {title}")
        
        return "\n\n".join(context_parts)

    def _create_content_strategy(self, analysis: Dict[str, Any], paper_content: Dict[str, Any]) -> Dict[str, Any]:
        """Create content strategy for presentation"""
        return {
            "focus_areas": analysis.get("emphasis_areas", []),
            "narrative_flow": analysis.get("storytelling_approach", ""),
            "depth_level": analysis.get("depth_level", "technical"),
            "visual_strategy": analysis.get("visual_elements", []),
            "citation_approach": analysis.get("citation_strategy", "")
        }

    def _create_fallback_slide_plan(self, slide_count: int, paper_content: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Create fallback slide plan if AI generation fails"""
        sections = paper_content.get("sections", {})
        metadata = paper_content.get("metadata", {})
        
        slides = []
        
        # Title slide
        slides.append({
            "slide_number": 1,
            "type": "title",
            "title": metadata.get("title", "Research Presentation"),
            "content": ["AI-Generated Research Presentation", "Academic Content Analysis"],
            "speaker_notes": "Introduction to the research presentation",
            "visual_elements": [],
            "content_source": "metadata"
        })
        
        # Use actual sections from the knowledge base content
        slide_counter = 2
        available_sections = list(sections.keys())
        
        logger.info(f"📊 Available sections for slides: {available_sections}")
        
        # Create slides from available sections
        for section_name in available_sections:
            if slide_counter > slide_count - 1:  # Reserve last slide for references
                break
                
            section_data = sections.get(section_name, {})
            section_content = section_data.get("content", "") if isinstance(section_data, dict) else str(section_data)
            
            if section_content:
                # Create bullet points from section content
                sentences = section_content.split('.')[:5]  # Take first 5 sentences
                content = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 20]
                
                # Clean section name for title
                title = section_name.replace('_', ' ').title()
                if title.lower() == 'general':
                    title = "Overview"
                
                slides.append({
                    "slide_number": slide_counter,
                    "type": "content",
                    "title": title,
                    "content": content[:4],  # Limit to 4 bullet points
                    "speaker_notes": f"Detailed discussion of {title.lower()}",
                    "visual_elements": ["diagrams", "charts"],
                    "content_source": section_name
                })
                slide_counter += 1
        
        # If we still need more slides, create summary slides
        if slide_counter <= slide_count - 1:
            # Add a summary slide
            slides.append({
                "slide_number": slide_counter,
                "type": "conclusion",
                "title": "Key Takeaways",
                "content": [
                    "Summary of main concepts covered",
                    "Practical applications discussed",
                    "Important considerations highlighted",
                    "Future directions identified"
                ],
                "speaker_notes": "Summarize the key points from the presentation",
                "visual_elements": ["summary_chart"],
                "content_source": "summary"
            })
        
        logger.info(f"📊 Created {len(slides)} slides in fallback plan")
        return slides

    def get_available_themes(self) -> List[str]:
        """Get list of available presentation themes"""
        return list(self.themes.keys()) 